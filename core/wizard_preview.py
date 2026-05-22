"""Wizard step-1: deck profile cards + pipeline-specific previews."""

from __future__ import annotations

import sys
from pathlib import Path
from typing import Any

from .ppt_test_path import ensure_engine_on_path
from .slide_limits import slide_limit_payload

ensure_engine_on_path()

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

from scripts.convert_legacy_deck_to_academy import (  # noqa: E402
    analyze_specs,
    convert_presentation,
)
from scripts.deck_migrate_config import migrate_config_for_source  # noqa: E402
from scripts.deck_profile import analyze_deck_structure, detect_deck_profile  # noqa: E402

PROFILE_CARDS: list[dict[str, str]] = [
    {
        "id": "migrate_cmp",
        "title": "도형 이식 (§7)",
        "subtitle": "migrate_cmp",
        "description": "파트너·CONTRABASS·CMP형 — 원본 도형·표·차트·그림을 유지하며 아카데미 레이아웃에 맞춥니다.",
    },
    {
        "id": "spec",
        "title": "텍스트·placeholder (§5)",
        "subtitle": "spec",
        "description": "강의안·Google export형 — 텍스트를 읽어 표지·목차·본문 placeholder에 채웁니다.",
    },
]


def _shape_counts(slide) -> dict[str, int]:
    """Local copy — audit_slide_mapping lives under ppt-academizer/scripts, not ppt-test."""
    counts = {"picture": 0, "chart": 0, "table": 0, "group": 0, "other": 0}
    for sh in slide.shapes:
        st = sh.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            counts["picture"] += 1
        elif st == MSO_SHAPE_TYPE.CHART or (getattr(sh, "has_chart", False) and sh.has_chart):
            counts["chart"] += 1
        elif st == MSO_SHAPE_TYPE.TABLE:
            counts["table"] += 1
        elif st == MSO_SHAPE_TYPE.GROUP:
            counts["group"] += 1
        else:
            counts["other"] += 1
    return counts


def _text_sample(slide, max_len: int = 80) -> str:
    parts: list[str] = []
    for sh in slide.shapes:
        if sh.has_text_frame and (sh.text or "").strip():
            parts.append(sh.text.strip().replace("\n", " ")[:max_len])
    return " | ".join(parts[:3])[:max_len]


def _structure_bullets(st: Any) -> list[str]:
    return [
        f"원본 {st.slide_count}장 · 슬라이드당 도형 약 {st.shapes_per_slide:.1f}개",
        f"그림 {st.picture_count} · 표 {st.table_count} · 차트 {st.chart_count}",
        f"텍스트 블록 {st.total_text_blocks}개 · 이미지-only 비율 {int(st.google_image_ratio * 100)}%",
    ]


def _preview_spec(source: Path, *, title: str, deck_subtitle: str) -> dict[str, Any]:
    specs = convert_presentation(
        source,
        deck_title=title,
        deck_subtitle=deck_subtitle,
        include_default_front_matter=True,
        front_matter_mode="auto",
    )
    rows = analyze_specs(specs)
    preview_slides: list[dict[str, Any]] = []
    for row in rows[:6]:
        texts = row.get("texts_preview") or []
        preview_slides.append(
            {
                "out_slide": row.get("slide"),
                "layout": row.get("layout"),
                "layout_reason": row.get("layout_reason"),
                "ingest_kind": row.get("ingest_kind"),
                "lines": texts[:2],
            }
        )
    kinds = []
    for s in specs[:3]:
        ing = s.get("_ingest") or s.get("ingest") or {}
        kinds.append(ing.get("kind") if isinstance(ing, dict) else None)
    summary = [
        f"출력 예상 {len(specs)}장 (표지·목차 자동 추출 포함)",
        "앞 장 표지/목차는 원본 1~2장에서 뽑을 수 있습니다",
    ]
    if "google_image_slide" in kinds:
        summary.append("일부 슬라이드는 배경 이미지로 옮깁니다 (편집 가능 텍스트 없음)")
    return {
        "output_slide_count": len(specs),
        "summary": summary,
        "slides": preview_slides,
    }


def _preview_migrate(source: Path) -> dict[str, Any]:
    from scripts.build_cmp_academy import build_slide_plan, classify_slide

    cfg = migrate_config_for_source(source)
    prs = Presentation(str(source))
    plan = build_slide_plan(prs, cfg)
    plan_len = len(plan)

    preview_slides: list[dict[str, Any]] = []
    shown = 0
    for src_idx, slide in enumerate(prs.slides):
        if shown >= 6:
            break
        kind = classify_slide(slide, src_idx, cfg.part_cover_indices)
        in_plan = next((k for si, k in plan if si == src_idx), None)
        if not in_plan and kind == "empty":
            continue
        plan_kind = in_plan or kind
        layout_hint = {
            "cover": "2_표지",
            "section": "간지",
            "content": "내지_거버닝 O / 1_내지_거버닝 X",
            "empty": "(생략)",
        }.get(plan_kind, plan_kind)
        preview_slides.append(
            {
                "src_slide": src_idx + 1,
                "plan_kind": plan_kind,
                "layout_hint": layout_hint,
                "lines": [_text_sample(slide, 100)],
                "shapes": _shape_counts(slide),
            }
        )
        shown += 1

    summary = [
        f"출력 예상 약 {plan_len}장 (원본 {len(prs.slides)}장 기준 §7 플랜)",
        f"deck_kind: {cfg.deck_kind}"
        + (" · CMP 자동 목차" if cfg.auto_toc_after_cover else ""),
        "도형·표·차트는 원본 위치를 유지하며 이식합니다",
    ]
    if cfg.skip_src_indices:
        summary.append(f"설정상 제외 슬라이드: {sorted(cfg.skip_src_indices)}")
    return {
        "output_slide_count": plan_len,
        "summary": summary,
        "slides": preview_slides,
    }


def build_wizard_preview(
    source: Path,
    *,
    deck_title: str | None = None,
    deck_subtitle: str = "PPT 아카데미화",
) -> dict[str, Any]:
    """Step-1 payload: detected profile, selectable cards, per-card previews."""
    source = Path(source).resolve()
    title = deck_title or source.stem
    prs = Presentation(str(source))
    detected, profile_meta = detect_deck_profile(prs, filename_hint=source.name)
    st = analyze_deck_structure(prs)

    if detected == "google_image":
        recommended = "spec"
    else:
        recommended = detected

    structure_bullets = _structure_bullets(st)
    spec_preview = _preview_spec(source, title=title, deck_subtitle=deck_subtitle)
    migrate_preview = _preview_migrate(source)

    cards: list[dict[str, Any]] = []
    for meta in PROFILE_CARDS:
        pid = meta["id"]
        card = {**meta, "recommended": pid == recommended}
        if pid == "migrate_cmp":
            card["preview"] = migrate_preview
        else:
            card["preview"] = spec_preview
            if detected == "google_image":
                card["note"] = "자동 감지: Google 이미지 export → spec + 배경 이미지"
        cards.append(card)

    est_out = None
    if recommended == "migrate_cmp":
        est_out = migrate_preview.get("output_slide_count")
    else:
        est_out = spec_preview.get("output_slide_count")

    return {
        "source": str(source),
        "deck_title": title,
        "detected_profile": detected,
        "recommended_profile": recommended,
        "profile_meta": profile_meta,
        "structure_bullets": structure_bullets,
        "cards": cards,
        "slide_limits": slide_limit_payload(
            source_slides=st.slide_count,
            estimated_output=est_out,
        ),
    }
