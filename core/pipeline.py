"""Academize: source .pptx → academy template .pptx."""

from __future__ import annotations

import json
import shutil
import sys
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any

from .ppt_test_path import ensure_engine_on_path
from .version import SERVICE_VERSION
from .postprocess_pptx import maybe_mac_powerpoint_repair
from .slide_limits import QualityMode, validate_quality_mode
from .warnings_display import format_warnings

ensure_engine_on_path()

from pptx import Presentation  # noqa: E402

from scripts.academy_deck_build_lib import (  # noqa: E402
    apply_background_images,
    build_from_json_specs,
    save_academy_deck,
)
from scripts.academy_template import resolve_academy_template_path  # noqa: E402
from scripts.convert_legacy_deck_to_academy import convert_presentation  # noqa: E402
from scripts.deck_profile import detect_deck_profile  # noqa: E402

Profile = str  # auto | spec | migrate_cmp | google_image


def _stamp() -> str:
    return datetime.now().strftime("%Y%m%d-%H%M%S")


def collect_warnings(specs: list[dict], extra: list[dict] | None = None) -> list[dict]:
    out: list[dict] = list(extra or [])
    for i, spec in enumerate(specs, start=1):
        ing = spec.get("_ingest") or {}
        if ing.get("kind") == "google_image_slide":
            out.append(
                {
                    "slide": i,
                    "code": "GOOGLE_IMAGE_SLIDE",
                    "message": "Google Slides 보내기 형식 — 슬라이드 내용을 배경 이미지로 옮겼습니다.",
                }
            )
        if ing.get("warning"):
            out.append(
                {
                    "slide": i,
                    "code": "TEXT_NOT_EXTRACTED",
                    "message": str(ing["warning"]),
                }
            )
    return out


def _save_specs_to_pptx(
    specs: list[dict],
    *,
    tmp: Path,
    stem: str,
    extra_warnings: list[dict] | None = None,
) -> Path:
    template = resolve_academy_template_path()
    out = tmp / f"academy-{stem}-{_stamp()}.pptx"
    shutil.copy2(template, out)
    prs = Presentation(str(out))
    build_from_json_specs(prs, specs)
    apply_background_images(prs, specs)
    if len(prs.slides) != len(specs):
        raise RuntimeError(f"slide count mismatch: built {len(prs.slides)}, expected {len(specs)}")
    save_academy_deck(prs, out)
    return out


def _academize_spec_from_specs(
    specs: list[dict],
    *,
    tmp: Path,
    stem: str,
    route: str,
    quality_mode: QualityMode = "standard",
    source_format: str = "pptx",
    ingest_warnings: list[dict] | None = None,
) -> tuple[Path, list[dict], int, dict[str, Any]]:
    out = _save_specs_to_pptx(specs, tmp=tmp, stem=stem)

    from scripts.migrate_version import MIGRATE_ENGINE_VERSION

    raw = collect_warnings(specs)
    for w in ingest_warnings or []:
        raw.append(w)
    raw.insert(0, {"code": "PIPELINE_SPEC", "message": ""})
    if quality_mode == "unlimited":
        raw.append({"code": "QUALITY_MODE_UNLIMITED", "message": ""})
    meta = {
        "pipeline": "spec_json",
        "route_profile": route,
        "source_format": source_format,
        "service_version": SERVICE_VERSION,
        "migrate_engine_version": MIGRATE_ENGINE_VERSION,
        "quality_mode": quality_mode,
    }
    return out, format_warnings(raw), len(specs), meta


def _academize_spec_path(
    source: Path,
    *,
    title: str,
    deck_subtitle: str,
    tmp: Path,
    route: str,
    quality_mode: QualityMode = "standard",
) -> tuple[Path, list[dict], int, dict[str, Any]]:
    assets_dir = tmp / "assets"
    specs = convert_presentation(
        source,
        deck_title=title,
        deck_subtitle=deck_subtitle,
        assets_dir=assets_dir,
        include_default_front_matter=True,
        front_matter_mode="auto",
    )
    stem = source.stem.replace(" ", "-")[:40]
    return _academize_spec_from_specs(
        specs,
        tmp=tmp,
        stem=stem,
        route=route,
        quality_mode=quality_mode,
        source_format="pptx",
    )


def _academize_migrate_cmp(
    source: Path,
    *,
    tmp: Path,
    quality_mode: QualityMode = "standard",
) -> tuple[Path, list[dict], int, dict[str, Any]]:
    from scripts.build_cmp_academy import migrate_cmp_deck

    stem = source.stem.replace(" ", "-")[:40]
    out = tmp / f"academy-cmp-{stem}-{_stamp()}.pptx"
    output, migrate_warnings, slide_count = migrate_cmp_deck(
        source, out, skip_ooxml_repair=False
    )
    raw_warnings = [
        dict(w)
        for w in migrate_warnings
        if w.get("code") not in ("MIGRATE_META",)
    ]
    from scripts.academy_deck_build_lib import finalize_academy_package

    svg_n = finalize_academy_package(output)
    if svg_n:
        raw_warnings.append(
            {
                "code": "SVG_RASTERIZED",
                "message": f"{svg_n}개 SVG를 PNG로 변환했습니다(레이아웃 복구 팝업 방지).",
            }
        )
    raw_warnings.extend(maybe_mac_powerpoint_repair(output))
    finalize_academy_package(output)
    raw_warnings.insert(
        0,
        {"code": "PIPELINE_MIGRATE_CMP", "message": ""},
    )
    if quality_mode == "unlimited":
        raw_warnings.append({"code": "QUALITY_MODE_UNLIMITED", "message": ""})
    warnings = format_warnings(raw_warnings)
    from scripts.migrate_version import MIGRATE_ENGINE_VERSION

    deck_kind = "generic_partner"
    for w in migrate_warnings:
        if w.get("code") == "MIGRATE_META":
            deck_kind = w.get("deck_kind", deck_kind)
            break

    meta = {
        "pipeline": "migrate_cmp",
        "route_profile": "migrate_cmp",
        "deck_kind": deck_kind,
        "service_version": SERVICE_VERSION,
        "migrate_engine_version": MIGRATE_ENGINE_VERSION,
        "quality_mode": quality_mode,
    }
    return output, warnings, slide_count, meta


def academize_pptx(
    source: Path,
    *,
    deck_title: str | None = None,
    deck_subtitle: str = "PPT 아카데미화",
    work_dir: Path | None = None,
    profile: Profile = "auto",
    quality_mode: str = "standard",
    estimated_output_slides: int | None = None,
) -> tuple[Path, list[dict], int, dict[str, Any]]:
    """Convert *source* pptx to academy deck.

    Returns (output_path, warnings, slide_count, meta).
    """
    source = Path(source).resolve()
    if not source.is_file():
        raise FileNotFoundError(source)

    title = deck_title or source.stem
    tmp = work_dir or Path(tempfile.mkdtemp(prefix="ppt-academizer-"))
    tmp.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(source))
    source_slides = len(prs.slides)
    detected, _ = detect_deck_profile(prs, filename_hint=source.name)
    route = detected if profile == "auto" else profile
    mode = validate_quality_mode(
        quality_mode,
        source_slides=source_slides,
        estimated_output=estimated_output_slides,
    )

    if route == "migrate_cmp":
        return _academize_migrate_cmp(source, tmp=tmp, quality_mode=mode)
    if route == "google_image":
        route = "spec"
    return _academize_spec_path(
        source,
        title=title,
        deck_subtitle=deck_subtitle,
        tmp=tmp,
        route=route,
        quality_mode=mode,
    )


def academize_pptx_legacy_return(
    source: Path,
    **kwargs: Any,
) -> tuple[Path, list[dict], int]:
    """Backward-compatible 3-tuple for callers that omit meta."""
    out, warnings, count, _meta = academize_pptx(source, **kwargs)
    return out, warnings, count
