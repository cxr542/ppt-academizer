#!/usr/bin/env python3
"""Slide mapping audit: source plan vs optional academy output (CONTRABASS pilot)."""

from __future__ import annotations

import argparse
import csv
import sys
from datetime import datetime
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from core.ppt_test_path import ensure_engine_on_path  # noqa: E402

ensure_engine_on_path()

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

from scripts.build_cmp_academy import (  # noqa: E402
    build_slide_plan,
    classify_slide,
    extract_header,
    migrate_config_for_source,
)
from scripts.deck_migrate_config import DeckMigrateConfig  # noqa: E402


def _shape_counts(slide) -> dict[str, int]:
    counts = {"picture": 0, "chart": 0, "table": 0, "group": 0, "other": 0}
    for sh in slide.shapes:
        st = sh.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            counts["picture"] += 1
        elif st == MSO_SHAPE_TYPE.CHART or (getattr(sh, "has_chart", False) and sh.has_chart):
            counts["chart"] += 1
        elif st == MSO_SHAPE_TYPE.TABLE:
            counts["table"] += 1
        elif st == MSO_SHAPE_TYPE.GROUP:
            counts["group"] += 1
        else:
            counts["other"] += 1
    return counts


def _text_sample(slide, max_len: int = 80) -> str:
    parts: list[str] = []
    for sh in slide.shapes:
        if sh.has_text_frame and (sh.text or "").strip():
            parts.append(sh.text.strip().replace("\n", " ")[:max_len])
    return " | ".join(parts[:3])[:max_len]


def audit_source(source: Path, cfg: DeckMigrateConfig) -> list[dict]:
    prs = Presentation(str(source))
    plan = build_slide_plan(prs, cfg)
    rows: list[dict] = []
    for src_idx, slide in enumerate(prs.slides):
        kind_raw = classify_slide(slide, src_idx, cfg.part_cover_indices)
        in_plan = next(((i, k) for i, (si, k) in enumerate(plan) if si == src_idx), None)
        title, _ = extract_header(slide, kind_raw if kind_raw != "empty" else "content")
        rows.append(
            {
                "src_index": src_idx,
                "classify": kind_raw,
                "in_plan": in_plan is not None,
                "plan_kind": in_plan[1] if in_plan else "",
                "plan_order": in_plan[0] if in_plan else "",
                "title": (title or "")[:120],
                "text_sample": _text_sample(slide),
                **_shape_counts(slide),
            }
        )
    return rows


def audit_output(output: Path) -> list[dict]:
    prs = Presentation(str(output))
    rows: list[dict] = []
    for out_idx, slide in enumerate(prs.slides):
        layout = slide.slide_layout.name if slide.slide_layout else ""
        title = ""
        for sh in slide.placeholders:
            if sh.placeholder_format.idx == 10 and (sh.text or "").strip():
                title = sh.text.strip().replace("\n", " ")[:120]
                break
        rows.append(
            {
                "out_index": out_idx,
                "layout": layout,
                "title": title,
                "text_sample": _text_sample(slide),
                **_shape_counts(slide),
            }
        )
    return rows


def write_csv(path: Path, rows: list[dict]) -> None:
    if not rows:
        return
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=list(rows[0].keys()))
        w.writeheader()
        w.writerows(rows)


def main() -> int:
    parser = argparse.ArgumentParser(description="Audit slide mapping for migrate_cmp")
    parser.add_argument("--source", type=Path, required=True)
    parser.add_argument("--output", type=Path, default=None, help="Academy output .pptx")
    parser.add_argument("--out-dir", type=Path, default=ROOT / "output" / "audit")
    args = parser.parse_args()

    if not args.source.is_file():
        print(f"Not found: {args.source}", file=sys.stderr)
        return 2

    cfg = migrate_config_for_source(args.source)
    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    stem = args.source.stem[:40]

    src_rows = audit_source(args.source.resolve(), cfg)
    src_csv = args.out_dir / f"{stem}-source-{stamp}.csv"
    write_csv(src_csv, src_rows)

    plan_len = sum(1 for r in src_rows if r["in_plan"])
    skipped = [r["src_index"] for r in src_rows if not r["in_plan"]]
    empty = [r["src_index"] for r in src_rows if r["classify"] == "empty"]

    print(f"deck_kind={cfg.deck_kind}")
    print(f"source_slides={len(src_rows)} plan_entries={plan_len} skipped_indices={skipped}")
    print(f"classify_empty={empty}")
    print(f"Wrote {src_csv}")

    if args.output and args.output.is_file():
        out_rows = audit_output(args.output.resolve())
        out_csv = args.out_dir / f"{stem}-output-{stamp}.csv"
        write_csv(out_csv, out_rows)
        print(f"output_slides={len(out_rows)} -> {out_csv}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
