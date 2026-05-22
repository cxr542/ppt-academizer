#!/usr/bin/env python3
"""Create tests/fixtures/cmp-like-partner.pptx for smoke (partner-heavy, migrate_cmp routing)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "tests" / "fixtures" / "cmp-like-partner.pptx"

# deck_profile.is_partner_shape_heavy: slide_count >= 8, shapes_per_slide >= 4
MIN_SLIDES = 8
BOXES_PER_SLIDE = 5


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    def box(slide, left, top, w, h, text):
        sh = slide.shapes.add_textbox(left, top, w, h)
        sh.text_frame.text = text

    sections = [
        ("클라우드 구현기술 (CMP)\n핵심기술 개요", "cover"),
        ("1. 가상화\n2. 컨테이너\n3. 오케스트레이션", "toc"),
        ("CONTAINER", "section"),
    ]
    for i in range(MIN_SLIDES):
        slide = prs.slides.add_slide(blank)
        if i == 0:
            box(slide, Inches(1), Inches(2.5), Inches(10), Inches(1), sections[0][0])
            continue
        if i == 1:
            box(slide, Inches(1), Inches(1), Inches(8), Inches(3), sections[1][0])
            continue
        if i == 2:
            box(slide, Inches(1), Inches(3), Inches(8), Inches(1), sections[2][0])
            title = f"주제 {i - 2}"
        else:
            title = f"Docker · Kubernetes · CMP 주제 {i - 2}"

        box(slide, Inches(0.6), Inches(0.5), Inches(9), Inches(0.5), title)
        box(slide, Inches(0.6), Inches(0.9), Inches(4), Inches(0.3), "개요")
        box(
            slide,
            Inches(0.6),
            Inches(1.4),
            Inches(5),
            Inches(2.5),
            "• 이미지 기반 배포\n• 격리된 프로세스\n• 레지스트리",
        )
        box(
            slide,
            Inches(7),
            Inches(1.4),
            Inches(5),
            Inches(2.5),
            "• 오케스트레이션\n• 서비스 메시\n• 관측성",
        )
        box(slide, Inches(0.6), Inches(4.2), Inches(11), Inches(0.4), f"슬라이드 {i + 1} 보조 설명")

    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides)} slides, ~{BOXES_PER_SLIDE} shapes/slide on content)")


if __name__ == "__main__":
    main()
