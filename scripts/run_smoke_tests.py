#!/usr/bin/env python3
"""Smoke tests for academy-design logic before full ppt-academizer implementation.

Runs without FastAPI. Uses bundled ``engine/`` (or ``PPT_ENGINE_ROOT``).

Usage:
  cd apps/ppt-academizer
  export TEMPLATE_PPTX="…"   # optional if Spotlight finds template
  .venv/bin/python scripts/run_smoke_tests.py
  .venv/bin/python scripts/run_smoke_tests.py --fixture tests/fixtures/apple-history.pptx
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import sys
import zipfile
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from core.ppt_test_path import ensure_engine_on_path  # noqa: E402

ENGINE = ensure_engine_on_path()

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from scripts.academy_deck_build_lib import (  # noqa: E402
    build_from_json_specs,
    save_academy_deck,
)
from scripts.academy_template import resolve_academy_template_path  # noqa: E402
from scripts.convert_legacy_deck_to_academy import (  # noqa: E402
    analyze_specs,
    convert_presentation,
    validate_spec,
)
from scripts.deck_profile import detect_deck_profile  # noqa: E402

OUT = ROOT / "output"
FIXTURES = ROOT / "tests" / "fixtures"
APPLE_HISTORY_PPTX = FIXTURES / "apple-history.pptx"
CMP_LIKE_PPTX = FIXTURES / "cmp-like-partner.pptx"
CONTRABASS_PPTX = FIXTURES / "contrabass-partner.pptx"
MINI_JSON = ENGINE / "docs" / "examples" / "academy-deck-mini.json"
APPLE_JSON = ENGINE / "docs" / "examples" / "apple-history-academy.json"


@dataclass
class CaseResult:
    name: str
    ok: bool
    detail: str = ""
    artifacts: list[str] = field(default_factory=list)


def _stamp() -> str:
    return datetime.now().strftime("%Y%m%d-%H%M%S")


def _read_json_slides(path: Path) -> list[dict]:
    data = json.loads(path.read_text(encoding="utf-8"))
    if isinstance(data, list):
        return data
    return data.get("slides", [])


def test_template_resolves() -> CaseResult:
    try:
        p = resolve_academy_template_path()
        return CaseResult("template_resolve", True, str(p))
    except FileNotFoundError as e:
        return CaseResult("template_resolve", False, str(e))


def test_spec_build_mini() -> CaseResult:
    """academy-design §5 spec path: JSON → build_from_json_specs → save_academy_deck."""
    try:
        template = resolve_academy_template_path()
        specs = _read_json_slides(MINI_JSON)
        OUT.mkdir(parents=True, exist_ok=True)
        out = OUT / f"smoke-spec-mini-{_stamp()}.pptx"
        shutil.copy2(template, out)
        prs = Presentation(str(out))
        build_from_json_specs(prs, specs)
        if len(prs.slides) != len(specs):
            return CaseResult(
                "spec_build_mini",
                False,
                f"slide count {len(prs.slides)} != {len(specs)}",
            )
        save_academy_deck(prs, out)
        return CaseResult("spec_build_mini", True, f"{len(specs)} slides", [str(out)])
    except Exception as e:
        return CaseResult("spec_build_mini", False, repr(e))


def test_spec_build_apple() -> CaseResult:
    try:
        template = resolve_academy_template_path()
        specs = _read_json_slides(APPLE_JSON)
        OUT.mkdir(parents=True, exist_ok=True)
        out = OUT / f"smoke-spec-apple-{_stamp()}.pptx"
        shutil.copy2(template, out)
        prs = Presentation(str(out))
        build_from_json_specs(prs, specs)
        if len(prs.slides) != len(specs):
            return CaseResult(
                "spec_build_apple",
                False,
                f"slide count {len(prs.slides)} != {len(specs)}",
            )
        save_academy_deck(prs, out)
        return CaseResult("spec_build_apple", True, f"{len(specs)} slides", [str(out)])
    except Exception as e:
        return CaseResult("spec_build_apple", False, repr(e))


def test_save_opens_slide_view(pptx_path: Path) -> CaseResult:
    """academy-design §6.3: lastView should be sldView after save_academy_deck."""
    try:
        with zipfile.ZipFile(pptx_path) as z:
            xml = z.read("ppt/viewProps.xml").decode("utf-8")
        if 'lastView="sldView"' in xml:
            return CaseResult("save_slide_view", True, "lastView=sldView")
        m = re.search(r'lastView="([^"]+)"', xml)
        view = m.group(1) if m else "unknown"
        return CaseResult("save_slide_view", False, f"lastView={view!r}")
    except Exception as e:
        return CaseResult("save_slide_view", False, repr(e))


def _make_ai_like_fixture(path: Path) -> None:
    """Minimal AI-style deck: blank layout + free text boxes (no academy placeholders)."""
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    def add_box(slide, left, top, w, h, text):
        box = slide.shapes.add_textbox(left, top, w, h)
        box.text_frame.text = text

    # slide 0 — cover-like
    s0 = prs.slides.add_slide(blank)
    add_box(s0, Inches(1), Inches(2.5), Inches(10), Inches(1.2), "Smoke Test Course\nAI Draft Cover")

    # slide 1 — toc-like
    s1 = prs.slides.add_slide(blank)
    add_box(s1, Inches(1), Inches(1), Inches(5), Inches(2), "1. Chapter One\n2. Chapter Two")

    # slide 2 — body
    s2 = prs.slides.add_slide(blank)
    add_box(s2, Inches(0.8), Inches(0.6), Inches(8), Inches(0.6), "Introduction")
    add_box(s2, Inches(0.8), Inches(1.0), Inches(4), Inches(0.4), "Overview")
    add_box(
        s2,
        Inches(0.8),
        Inches(1.6),
        Inches(11),
        Inches(4),
        "• First bullet point\n• Second bullet\n\nBody paragraph from AI export.",
    )

    prs.save(path)


def _legacy_pipeline(
    src: Path,
    *,
    case_name: str,
    deck_title: str,
    deck_subtitle: str,
    min_specs: int = 3,
) -> CaseResult:
    """A → convert_legacy (heuristic) → spec build."""
    try:
        if not src.is_file():
            return CaseResult(case_name, False, f"not found: {src}")

        stem = src.stem.replace(" ", "-")
        src_slides = len(Presentation(str(src)).slides)
        assets_dir = OUT / f"assets-{stem}"
        specs = convert_presentation(
            src,
            deck_title=deck_title,
            deck_subtitle=deck_subtitle,
            assets_dir=assets_dir,
            include_default_front_matter=True,
            front_matter_mode="auto",
        )
        paas_toc = any(
            "Docker" in str(s.get("texts", [""])[0]) and s.get("layout") == "목차" for s in specs
        )
        if paas_toc and "PaaS" not in deck_title:
            return CaseResult(case_name, False, "unexpected hardcoded PaaS TOC in specs")
        val_issues = [v for row in analyze_specs(specs) for v in row.get("validation", [])]
        if val_issues:
            return CaseResult(case_name, False, f"spec validation: {val_issues[:3]}")
        image_specs = sum(1 for s in specs if s.get("background_image"))
        if len(specs) < min_specs:
            return CaseResult(case_name, False, f"only {len(specs)} specs")

        template = resolve_academy_template_path()
        OUT.mkdir(parents=True, exist_ok=True)
        out = OUT / f"smoke-{stem}-{_stamp()}.pptx"
        shutil.copy2(template, out)
        prs = Presentation(str(out))
        build_from_json_specs(prs, specs)
        if len(prs.slides) != len(specs):
            return CaseResult(
                case_name,
                False,
                f"built {len(prs.slides)} slides, expected {len(specs)}",
            )
        save_academy_deck(prs, out)

        spec_path = OUT / f"smoke-{stem}-{_stamp()}.json"
        spec_path.write_text(json.dumps(specs, ensure_ascii=False, indent=2), encoding="utf-8")

        detail = (
            f"source {src_slides} → spec {len(specs)} slides ({src.name}); "
            f"background_image={image_specs}"
        )
        artifacts = [str(out), str(spec_path)]
        if assets_dir.is_dir():
            artifacts.append(str(assets_dir))
        return CaseResult(case_name, True, detail, artifacts)
    except Exception as e:
        return CaseResult(case_name, False, repr(e))


def test_legacy_ingest_to_spec_build() -> CaseResult:
    """AI-like fixture pptx."""
    FIXTURES.mkdir(parents=True, exist_ok=True)
    src = FIXTURES / "ai-like-draft.pptx"
    if not src.is_file():
        _make_ai_like_fixture(src)
    return _legacy_pipeline(
        src,
        case_name="legacy_to_spec",
        deck_title="Smoke AI Draft",
        deck_subtitle="academizer smoke test",
    )


def test_cmp_like_fixture() -> CaseResult:
    if not CMP_LIKE_PPTX.is_file():
        import subprocess

        gen = ROOT / "scripts" / "make_cmp_like_fixture.py"
        try:
            subprocess.run([sys.executable, str(gen)], check=True, cwd=ROOT)
        except Exception as e:
            return CaseResult("cmp_like_fixture", False, f"fixture gen failed: {e}")
    if not CMP_LIKE_PPTX.is_file():
        return CaseResult("cmp_like_fixture", False, f"missing {CMP_LIKE_PPTX}")
    profile, _ = detect_deck_profile(CMP_LIKE_PPTX)
    if profile != "migrate_cmp":
        return CaseResult("cmp_like_fixture", False, f"expected migrate_cmp, got {profile}")
    return _legacy_pipeline(
        CMP_LIKE_PPTX,
        case_name="cmp_like_spec_preview",
        deck_title="CMP-like",
        deck_subtitle="smoke",
        min_specs=3,
    )


def test_migrate_contrabass_fidelity() -> CaseResult:
    """§7 migrate: CONTRABASS deck_kind, slide count, no auto TOC."""
    if not CONTRABASS_PPTX.is_file():
        return CaseResult(
            "contrabass_migrate",
            True,
            f"SKIP: missing {CONTRABASS_PPTX} (run scripts/restore_fixtures.py)",
        )
    try:
        from scripts.build_cmp_academy import (
            build_slide_plan,
            expected_output_slide_count,
            migrate_cmp_deck,
        )
        from scripts.deck_migrate_config import migrate_config_for_source

        cfg = migrate_config_for_source(CONTRABASS_PPTX)
        if cfg.deck_kind != "contrabass":
            return CaseResult(
                "contrabass_migrate",
                False,
                f"expected deck_kind contrabass, got {cfg.deck_kind}",
            )
        if cfg.auto_toc_after_cover:
            return CaseResult("contrabass_migrate", False, "CONTRABASS must not auto-insert TOC")

        src_prs = Presentation(str(CONTRABASS_PPTX))
        plan = build_slide_plan(src_prs, cfg)
        expected = expected_output_slide_count(plan, cfg)
        out = OUT / f"smoke-contrabass-{_stamp()}.pptx"
        OUT.mkdir(parents=True, exist_ok=True)
        output, warnings, count = migrate_cmp_deck(CONTRABASS_PPTX, out, skip_ooxml_repair=True)
        mismatch = [w for w in warnings if w.get("code") == "SLIDE_COUNT_MISMATCH"]
        if mismatch:
            return CaseResult("contrabass_migrate", False, mismatch[0].get("message", ""))
        if count != expected:
            return CaseResult(
                "contrabass_migrate",
                False,
                f"slides {count} != expected {expected} (plan {len(plan)})",
            )
        return CaseResult(
            "contrabass_migrate",
            True,
            f"deck_kind=contrabass plan={len(plan)} out={count} -> {output.name}",
            [str(output)],
        )
    except Exception as e:
        return CaseResult("contrabass_migrate", False, repr(e))


def test_migrate_cmp_like_slide_count() -> CaseResult:
    if not CMP_LIKE_PPTX.is_file():
        return CaseResult("cmp_like_migrate", True, "SKIP: cmp-like missing")
    try:
        from scripts.build_cmp_academy import (
            build_slide_plan,
            expected_output_slide_count,
            migrate_cmp_deck,
        )
        from scripts.deck_migrate_config import migrate_config_for_source

        cfg = migrate_config_for_source(CMP_LIKE_PPTX)
        plan = build_slide_plan(Presentation(str(CMP_LIKE_PPTX)), cfg)
        expected = expected_output_slide_count(plan, cfg)
        out = OUT / f"smoke-cmp-like-migrate-{_stamp()}.pptx"
        _, warnings, count = migrate_cmp_deck(CMP_LIKE_PPTX, out, skip_ooxml_repair=True)
        if any(w.get("code") == "SLIDE_COUNT_MISMATCH" for w in warnings):
            return CaseResult("cmp_like_migrate", False, "SLIDE_COUNT_MISMATCH")
        if count != expected:
            return CaseResult("cmp_like_migrate", False, f"{count} != {expected}")
        return CaseResult("cmp_like_migrate", True, f"slides={count} deck_kind={cfg.deck_kind}")
    except Exception as e:
        return CaseResult("cmp_like_migrate", False, repr(e))


def test_apple_history_pptx_fixture() -> CaseResult:
    """Real-ish deck: tests/fixtures/apple-history.pptx."""
    return _legacy_pipeline(
        APPLE_HISTORY_PPTX,
        case_name="apple_history_pptx",
        deck_title="Apple History",
        deck_subtitle="Fixture academize smoke",
        min_specs=5,
    )


def test_fixture_pptx(path: Path, title: str, subtitle: str) -> CaseResult:
    stem = path.stem.replace(" ", "-")
    return _legacy_pipeline(
        path,
        case_name=f"fixture_{stem}",
        deck_title=title,
        deck_subtitle=subtitle,
        min_specs=3,
    )


def main() -> int:
    parser = argparse.ArgumentParser(description="ppt-academizer smoke tests")
    parser.add_argument(
        "--fixture",
        type=Path,
        help="Only run legacy pipeline on this .pptx (e.g. tests/fixtures/apple-history.pptx)",
    )
    parser.add_argument("--title", default="Academize Fixture", help="Cover title for --fixture")
    parser.add_argument("--subtitle", default="ppt-academizer smoke", help="Cover subtitle")
    args = parser.parse_args()

    print("ppt-academizer smoke tests (academy-design logic)\n")
    results: list[CaseResult] = []

    if args.fixture:
        results.append(test_template_resolves())
        if not results[-1].ok:
            _print_report(results)
            return 2
        p = args.fixture if args.fixture.is_absolute() else ROOT / args.fixture
        results.append(test_fixture_pptx(p.resolve(), args.title, args.subtitle))
        if results[-1].ok and results[-1].artifacts:
            results.append(test_save_opens_slide_view(Path(results[-1].artifacts[0])))
        _print_report(results)
        return 1 if any(not r.ok for r in results) else 0

    results.append(test_template_resolves())
    if not results[-1].ok:
        print("STOP: template required for remaining tests.\n")
        _print_report(results)
        return 2

    results.append(test_spec_build_mini())
    results.append(test_spec_build_apple())

    for r in results:
        if r.name == "spec_build_mini" and r.ok and r.artifacts:
            results.append(test_save_opens_slide_view(Path(r.artifacts[0])))
            break

    results.append(test_legacy_ingest_to_spec_build())
    results.append(test_cmp_like_fixture())
    results.append(test_migrate_cmp_like_slide_count())
    results.append(test_migrate_contrabass_fidelity())
    if APPLE_HISTORY_PPTX.is_file():
        results.append(test_apple_history_pptx_fixture())
    else:
        results.append(
            CaseResult("apple_history_pptx", False, f"missing {APPLE_HISTORY_PPTX}")
        )

    for r in results:
        if r.ok and r.artifacts and r.name in (
            "legacy_to_spec",
            "apple_history_pptx",
        ):
            results.append(test_save_opens_slide_view(Path(r.artifacts[0])))

    _print_report(results)
    failed = sum(1 for r in results if not r.ok)
    return 1 if failed else 0


def _print_report(results: list[CaseResult]) -> None:
    print("-" * 60)
    for r in results:
        mark = "PASS" if r.ok else "FAIL"
        print(f"[{mark}] {r.name}: {r.detail}")
        for a in r.artifacts:
            print(f"       → {a}")
    print("-" * 60)
    passed = sum(1 for r in results if r.ok)
    print(f"{passed}/{len(results)} passed")
    report = OUT / f"smoke-report-{_stamp()}.json"
    OUT.mkdir(parents=True, exist_ok=True)
    report.write_text(
        json.dumps(
            [{"name": r.name, "ok": r.ok, "detail": r.detail, "artifacts": r.artifacts} for r in results],
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    print(f"Report: {report}")


if __name__ == "__main__":
    raise SystemExit(main())
