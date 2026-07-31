"""Unit tests for academy 2/3/4-column card layout pattern detection."""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches

from scripts.academy_layout_patterns import (
    _is_light_fill,
    apply_card_column_pattern,
    detect_card_column_count,
)


def test_light_fill_accepts_soft_cards_only():
    assert _is_light_fill("F7F8FA")
    assert _is_light_fill("FFFFFF")
    assert _is_light_fill("E6F1FF")
    assert not _is_light_fill("051C2C")
    assert not _is_light_fill("00A3E0")
    assert not _is_light_fill(None)


def _blank_slide() -> tuple[Presentation, object]:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    return prs, prs.slides.add_slide(layout)


def _add_card(slide, left, top, width, height, hex_rgb: str):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hex_rgb)
    shape.line.fill.background()
    return shape


def test_detect_three_soft_cards():
    prs, slide = _blank_slide()
    sw = int(prs.slide_width)
    tops = 3_420_000
    for i, fill in enumerate(("F7F8FA", "F7F8FA", "051C2C")):
        _add_card(slide, 540_000 + i * 3_852_000, tops, 3_528_000, 1_980_000, fill)
    assert detect_card_column_count(slide, sw) == 3


def test_reject_two_dark_callout_strips():
    prs, slide = _blank_slide()
    sw = int(prs.slide_width)
    _add_card(slide, 504_000, 4_680_000, 5_580_000, 1_224_000, "0A3A4A")
    _add_card(slide, 6_300_000, 4_680_000, 5_400_000, 1_224_000, "00A3E0")
    assert detect_card_column_count(slide, sw) == 0


def test_reject_stacked_card_rows():
    prs, slide = _blank_slide()
    sw = int(prs.slide_width)
    for row, fill in ((1_224_000, "F7F8FA"), (3_672_000, "F7F8FA")):
        for i in range(3):
            _add_card(slide, 540_000 + i * 3_852_000, row, 3_672_000, 2_232_000, fill)
    assert detect_card_column_count(slide, sw) == 0


def test_apply_snaps_to_pattern_slots():
    prs, slide = _blank_slide()
    sw = int(prs.slide_width)
    for i in range(2):
        _add_card(slide, 540_000 + i * 5_652_000, 1_224_000, 5_400_000, 4_752_000, "F7F8FA")
    patterns = {
        "2": {
            "name": "레이아웃_2단",
            "cards": [
                {
                    "left": 513_025,
                    "top": 1_700_213,
                    "width": 5_403_587,
                    "height": 4_573_587,
                    "fill": "#E6F1FF",
                },
                {
                    "left": 6_283_464,
                    "top": 1_700_213,
                    "width": 5_403_587,
                    "height": 4_573_587,
                    "fill": "#DDE7FB",
                },
            ],
        }
    }
    name = apply_card_column_pattern(slide, patterns, sw, column_count=2)
    assert name == "레이아웃_2단"
    tops = sorted(int(sh.top) for sh in slide.shapes if not sh.is_placeholder)
    assert tops[0] == 1_700_213
