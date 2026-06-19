"""Tests for deck profile detection."""

from __future__ import annotations

import shutil
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from scripts.convert_legacy_deck_to_academy import convert_presentation, validate_spec
from scripts.deck_profile import (
    ai_filename_hint,
    analyze_deck_structure,
    detect_deck_profile,
    is_ai_freeform_text_deck,
    is_lab_lecture_deck,
    is_partner_shape_heavy,
)
from scripts.front_matter import _is_toc_like, extract_front_matter


def _blank_deck(path: Path, slides: int = 3, *, boxes_per_slide: int = 1) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    for i in range(slides):
        s = prs.slides.add_slide(layout)
        for j in range(boxes_per_slide):
            box = s.shapes.add_textbox(
                Inches(0.5 + j * 0.1),
                Inches(1 + j * 0.1),
                Inches(4),
                Inches(0.4),
            )
            box.text_frame.text = f"Slide {i + 1} block {j} with enough text"
    prs.save(path)


def test_partner_heavy_ignores_apple_history_filename(tmp_path: Path) -> None:
    """Many shapes/slides → migrate_cmp even if file is named apple-history.pptx."""
    p = tmp_path / "apple-history.pptx"
    _blank_deck(p, slides=12, boxes_per_slide=5)
    profile, meta = detect_deck_profile(p)
    assert profile == "migrate_cmp"
    assert meta.get("reason") == "partner_shape_heavy"
    assert meta.get("shapes_per_slide", 0) >= 4.0


def test_sparse_deck_stays_spec_even_with_contrabass_filename(tmp_path: Path) -> None:
    """Few shapes → spec even if filename says CONTRABASS."""
    p = tmp_path / "CONTRABASS-light.pptx"
    _blank_deck(p, slides=5, boxes_per_slide=1)
    profile, meta = detect_deck_profile(p)
    assert profile == "spec"
    assert meta.get("reason") in (
        "text_lecture_structure",
        "extractable_text_blocks",
        "default_spec_heuristic",
        "ai_freeform_export",
    )


def test_ai_like_deck_routes_spec(tmp_path: Path) -> None:
    """Blank layout + text boxes (ChatGPT-style) → spec, not migrate_cmp."""
    p = tmp_path / "chatgpt-course-draft.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    for title in ("AI Course", "1. One\n2. Two", "Body"):
        s = prs.slides.add_slide(blank)
        s.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2)).text_frame.text = title
    prs.save(p)
    profile, meta = detect_deck_profile(p)
    st = analyze_deck_structure(Presentation(str(p)))
    assert profile == "spec"
    assert meta.get("reason") in ("ai_freeform_export", "text_lecture_structure")
    assert ai_filename_hint(p.name)


def test_detect_cmp_filename(tmp_path: Path) -> None:
    p = tmp_path / "cloud-cmp-deck.pptx"
    _blank_deck(p, slides=10, boxes_per_slide=5)
    profile, meta = detect_deck_profile(p)
    assert profile == "migrate_cmp"
    assert meta.get("reason") == "partner_shape_heavy"


def test_lab_lecture_routes_spec_before_partner_heavy(tmp_path: Path) -> None:
    p = tmp_path / "k8s-dashboard-lab.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide_terms = [
        ("학습 목표", "Dashboard 실습 목표", "Deployment Service Ingress ConfigMap Volume Mount"),
        ("리소스 생성용 YAML", "apiVersion: v1\nkind: Service\nmetadata:\n  name: nginx-service", "kubectl"),
        ("강사용 진행 멘트", "Deployment", "“Deployment가 원하는 Pod 상태를 유지하게 합니다.”"),
        ("검증 체크리스트", "성공 체크리스트", "Pod Running Ready 1/1"),
        ("ConfigMap 적용 절차", "YAML Import", "volumeMounts containers"),
        ("마무리", "Ingress URL 확인", "실습 완료"),
        ("전체 실습 흐름", "Deployment Service Ingress", "ConfigMap Volume Mount"),
        ("사전 준비", "Namespace StorageClass", "Dashboard 권한"),
    ]
    for title, body, extra in slide_terms:
        slide = prs.slides.add_slide(blank)
        for idx in range(5):
            left = Inches(0.5 + idx * 1.8)
            top = Inches(0.4 + idx * 0.6)
            box = slide.shapes.add_textbox(left, top, Inches(1.5), Inches(0.4))
            box.text_frame.text = f"{title}\n{body}\n{extra}"
    prs.save(p)

    st = analyze_deck_structure(Presentation(str(p)))
    profile, meta = detect_deck_profile(p)

    assert is_partner_shape_heavy(st)
    assert is_lab_lecture_deck(st)
    assert profile == "spec"
    assert meta.get("reason") == "lab_lecture_structure"


@pytest.mark.skipif(
    not Path("/Users/yhkim/Downloads/CONTRABASS 기반기술@260504_수정 요청.pptx").is_file(),
    reason="CONTRABASS sample not on disk",
)
def test_contrabass_content_wrong_filename(tmp_path: Path) -> None:
    src = Path("/Users/yhkim/Downloads/CONTRABASS 기반기술@260504_수정 요청.pptx")
    p = tmp_path / "apple-history.pptx"
    shutil.copy(src, p)
    profile, meta = detect_deck_profile(p)
    st = analyze_deck_structure(Presentation(str(p)))
    assert is_partner_shape_heavy(st)
    assert profile == "migrate_cmp"
    assert meta.get("reason") == "partner_shape_heavy"


def test_chapter_bullets_are_not_toc() -> None:
    """Era body slide (title + • 연도별 사실) must not become academy 목차."""
    merged = (
        "창업과 초기 (1976–1984)\n"
        "• 1976년: 애플 컴퓨터 컴퍼니 창업\n"
        "• 1977년: 대중용 PC 시장을 연 Apple II 출시\n"
        "• 1983년: GUI를 탑재한 리사(Lisa) 선보임"
    )
    blocks = [{"text": merged, "len": len(merged), "top": 1.0, "left": 0.5}]
    assert _is_toc_like(blocks) is False


def test_extract_front_matter_skips_cover_toc(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    s0 = prs.slides.add_slide(layout)
    s0.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1)).text_frame.text = "My Course"
    s1 = prs.slides.add_slide(layout)
    s1.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2)).text_frame.text = (
        "1. Chapter A\n2. Chapter B"
    )
    s2 = prs.slides.add_slide(layout)
    s2.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1)).text_frame.text = "Body slide"
    prs.save(p)

    fm, skip = extract_front_matter(prs, deck_title="T", deck_subtitle="S")
    assert len(fm) == 2
    assert fm[0]["layout"] == "2_표지"
    assert fm[1]["layout"] == "목차"
    assert skip == {0, 1}

    specs = convert_presentation(
        p, deck_title="T", deck_subtitle="S", include_default_front_matter=True, front_matter_mode="auto"
    )
    assert specs[0]["layout"] == "2_표지"
    assert "Docker" not in specs[1]["texts"][0]
    assert len(specs) == 3
    for i, spec in enumerate(specs, start=1):
        assert not validate_spec(spec, slide_index=i)
