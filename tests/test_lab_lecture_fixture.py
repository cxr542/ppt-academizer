from __future__ import annotations

import importlib.util
import sys
from pathlib import Path
from xml.etree import ElementTree as ET
from zipfile import ZipFile

import pytest
from pptx import Presentation

from core.pipeline import academize_pptx
from scripts.deck_profile import detect_deck_profile

ROOT = Path(__file__).resolve().parents[1]
LAB_FIXTURE = ROOT / "tests" / "fixtures" / "real_world" / "01_k8s_dashboard_lab_lecture.pptx"
_VALIDATOR_SPEC = importlib.util.spec_from_file_location(
    "validate_pptx_integrity",
    ROOT / "scripts" / "validate_pptx_integrity.py",
)
assert _VALIDATOR_SPEC is not None
assert _VALIDATOR_SPEC.loader is not None
_VALIDATOR = importlib.util.module_from_spec(_VALIDATOR_SPEC)
sys.modules[_VALIDATOR_SPEC.name] = _VALIDATOR
_VALIDATOR_SPEC.loader.exec_module(_VALIDATOR)


def _slide_text(slide) -> str:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = (shape.text or "").strip()
            if text:
                texts.append(text)
    return "\n".join(texts)


_YAML_TOKENS = (
    "apiVersion:",
    "kind:",
    "metadata:",
    "spec:",
    "containers:",
    "volumeMounts:",
    "volumes:",
    "index.html: |",
)


def _yaml_shape_texts(slide) -> list[str]:
    texts: list[tuple[int, int, str]] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text or "").replace("\x0b", "\n").strip()
        if text and any(token in text for token in _YAML_TOKENS):
            texts.append((int(shape.top), int(shape.left), text))
    return [text for _top, _left, text in sorted(texts)]


def _assert_ordered(text: str, tokens: tuple[str, ...]) -> None:
    cursor = -1
    for token in tokens:
        next_cursor = text.find(token, cursor + 1)
        assert next_cursor > cursor
        cursor = next_cursor


def _notes_text(slide) -> str:
    try:
        return (slide.notes_slide.notes_text_frame.text or "").strip()
    except AttributeError:
        return ""


def _text_shape_boxes(slide) -> list[tuple[float, float, float, float, str]]:
    boxes: list[tuple[float, float, float, float, str]] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text or "").strip()
        if not text:
            continue
        boxes.append((float(shape.left), float(shape.top), float(shape.width), float(shape.height), text))
    return boxes


def _severe_overlap_count(slide) -> int:
    boxes = _text_shape_boxes(slide)
    severe = 0
    for index, box_a in enumerate(boxes):
        ax, ay, aw, ah, _ = box_a
        for box_b in boxes[index + 1 :]:
            bx, by, bw, bh, _ = box_b
            ix = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
            iy = max(0.0, min(ay + ah, by + bh) - max(ay, by))
            overlap = ix * iy
            if overlap <= 0:
                continue
            smaller = min(aw * ah, bw * bh)
            if smaller > 0 and overlap / smaller > 0.7:
                severe += 1
    return severe


def _placeholder_box(slide, idx: int) -> tuple[int, int, int, int]:
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return (int(shape.left), int(shape.top), int(shape.width), int(shape.height))
    raise AssertionError(f"placeholder {idx} not found")


def _non_placeholder_shape_count(slide) -> int:
    return sum(1 for shape in slide.shapes if not shape.is_placeholder)


def _text_outline_count(path: Path) -> int:
    count = 0
    targets = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}endParaRPr",
    }
    outlines = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ln",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
    }
    with ZipFile(path) as archive:
        for name in archive.namelist():
            if not name.startswith("ppt/") or not name.endswith(".xml"):
                continue
            root = ET.fromstring(archive.read(name))
            for prop in root.iter():
                if prop.tag in targets:
                    count += sum(1 for child in prop.iter() if child is not prop and child.tag in outlines)
    return count


def _visible_no_fill_text_shape_line_count(path: Path, slide_numbers: tuple[int, ...]) -> int:
    ns = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    count = 0
    with ZipFile(path) as archive:
        for slide_number in slide_numbers:
            root = ET.fromstring(archive.read(f"ppt/slides/slide{slide_number}.xml"))
            for shape in root.findall(".//p:sp", ns):
                if shape.find("p:nvSpPr/p:nvPr/p:ph", ns) is not None:
                    continue
                if not "".join(text.text or "" for text in shape.findall(".//a:t", ns)).strip():
                    continue
                sp_pr = shape.find("p:spPr", ns)
                if sp_pr is None or sp_pr.find("a:noFill", ns) is None:
                    continue
                line = sp_pr.find("a:ln", ns)
                if line is not None and line.find("a:noFill", ns) is None:
                    count += 1
    return count


def _visible_non_text_shape_line_count(path: Path, slide_number: int) -> int:
    ns = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    count = 0
    with ZipFile(path) as archive:
        root = ET.fromstring(archive.read(f"ppt/slides/slide{slide_number}.xml"))
    for shape in root.findall(".//p:sp", ns):
        if "".join(text.text or "" for text in shape.findall(".//a:t", ns)).strip():
            continue
        line = shape.find("p:spPr/a:ln", ns)
        if line is not None and line.find("a:noFill", ns) is None:
            count += 1
    return count


@pytest.mark.skipif(not LAB_FIXTURE.is_file(), reason="real-world lab fixture is local-only")
def test_lab_fixture_routes_to_spec() -> None:
    profile, meta = detect_deck_profile(LAB_FIXTURE)

    assert profile == "spec"
    assert meta.get("reason") == "lab_lecture_structure"


@pytest.mark.skipif(not LAB_FIXTURE.is_file(), reason="real-world lab fixture is local-only")
def test_lab_fixture_preserves_p0_content(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv("PPT_ACADEMIZER_SKIP_PP_REPAIR", "1")
    source = Presentation(str(LAB_FIXTURE))
    output, warnings, slide_count, meta = academize_pptx(
        LAB_FIXTURE,
        work_dir=tmp_path,
        quality_mode="unlimited",
    )
    result = Presentation(str(output))

    assert _VALIDATOR.validate_pptx(output) == []
    assert _text_outline_count(output) == 0
    assert _visible_no_fill_text_shape_line_count(output, (1, 3, 4, 5, 10, 11, 13)) == 0
    assert _visible_non_text_shape_line_count(output, 3) >= 5
    assert meta["pipeline"] == "spec_json"
    assert slide_count == 13
    assert len(result.slides) == 13
    assert sum(1 for slide in result.slides if _notes_text(slide)) == 13
    assert _notes_text(result.slides[0]) == _notes_text(source.slides[0])

    result_text = "\n".join(_slide_text(slide) for slide in result.slides)
    for token in ("apiVersion", "kind:", "metadata:", "containers:", "volumeMounts:"):
        assert token in result_text
    assert "Dashboard의 YAML Import 기능을 쓰면 리소스를 즉시 생성할 수 있다." in result_text
    assert "Deployment가 원하는 Pod 상태를 유지하게 합니다" in result_text

    slide8_yaml = _yaml_shape_texts(result.slides[7])
    slide9_yaml = _yaml_shape_texts(result.slides[8])
    assert len(slide8_yaml) >= 2
    assert len(slide9_yaml) >= 2
    assert not any("Dashboard의 YAML Import 기능" in text for text in slide8_yaml)
    assert not any("K8s Dashboard 실습 | 8" in text for text in slide8_yaml)
    assert not any("ConfigMap을 만들고 Deployment" in text for text in slide9_yaml)
    assert not any("K8s Dashboard 실습 | 9" in text for text in slide9_yaml)
    _assert_ordered("\n".join(slide8_yaml), ("# Namespace", "# Service", "# Ingress"))
    _assert_ordered(
        "\n".join(slide9_yaml),
        ("index.html: |", "    <!DOCTYPE html>", "    <html>", "volumeMounts:", "volumes:"),
    )

    speaker_slide = next(
        slide for slide in result.slides if "강사용 진행 멘트" in _slide_text(slide)
    )
    assert _severe_overlap_count(speaker_slide) == 0
    assert _placeholder_box(speaker_slide, 10)
    assert any(w["code"] == "PIPELINE_SPEC" for w in warnings)


@pytest.mark.skipif(not LAB_FIXTURE.is_file(), reason="real-world lab fixture is local-only")
def test_lab_fixture_uses_readable_header_and_body_geometry(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("PPT_ACADEMIZER_SKIP_PP_REPAIR", "1")
    output, _warnings, _slide_count, _meta = academize_pptx(
        LAB_FIXTURE,
        work_dir=tmp_path,
        quality_mode="unlimited",
    )
    result = Presentation(str(output))

    for slide_number in (2, 3, 4, 5, 11, 12):
        slide = result.slides[slide_number - 1]
        _title_left, title_top, _title_w, title_h = _placeholder_box(slide, 10)
        _sub_left, subtitle_top, _sub_w, subtitle_h = _placeholder_box(slide, 12)
        _body_left, body_top, _body_w, _body_h = _placeholder_box(slide, 13)

        assert title_top < subtitle_top
        assert title_top + title_h <= subtitle_top
        assert subtitle_top + subtitle_h <= body_top
        assert body_top >= 1_150_000

    speaker_slide = result.slides[11]
    speaker_title = next(
        (shape.text or "").strip()
        for shape in speaker_slide.placeholders
        if shape.placeholder_format.idx == 10
    )
    assert speaker_title == "강사용 진행 멘트"

    for slide_number in (8, 9):
        slide = result.slides[slide_number - 1]
        _title_left, title_top, title_w, _title_h = _placeholder_box(slide, 10)
        _body_left, body_top, body_w, _body_h = _placeholder_box(slide, 12)

        assert title_top < body_top
    assert title_w >= 9_000_000
    assert body_w >= 10_000_000


@pytest.mark.skipif(not LAB_FIXTURE.is_file(), reason="real-world lab fixture is local-only")
def test_lab_fixture_preserves_visual_heavy_shapes(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("PPT_ACADEMIZER_SKIP_PP_REPAIR", "1")
    output, _warnings, _slide_count, _meta = academize_pptx(
        LAB_FIXTURE,
        work_dir=tmp_path,
        quality_mode="unlimited",
    )
    result = Presentation(str(output))

    assert _non_placeholder_shape_count(result.slides[0]) >= 20
    for slide_number in (3, 4, 5, 10, 11, 13):
        slide = result.slides[slide_number - 1]
        assert _non_placeholder_shape_count(slide) >= 6
        for idx in (12, 13):
            left, top, width, height = _placeholder_box(slide, idx)
            assert left > result.slide_width
            assert top > result.slide_height
            assert width <= 60_000
            assert height <= 60_000
            assert not next(
                (shape.text or "").strip()
                for shape in slide.placeholders
                if shape.placeholder_format.idx == idx
            )
