"""Shared helpers for building decks from the academy 2026 template."""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.slide import SlidePart
from pptx.util import Inches, Pt

from scripts.shape_migrate_lib import copy_shape_hybrid, renumber_shape_ids

FONT_BODY = "프리젠테이션 4 Regular"
FONT_TITLE = "프리젠테이션 8 ExtraBold"

# Content slide title (ph10) — CONTRABASS v9 parity
CLR_DK1 = RGBColor(0x00, 0x00, 0x00)
LAYOUT_PH10_LEFT = 781_263
SLIDE_MARGIN_X = 380_000
SLIDE_MARGIN_BOTTOM = 520_000
BODY_AREA_LEFT = 514_984
BODY_TOP_MIN = 1_050_000
HEADER_BOTTOM_EMU = 1_150_000
GOVERNING_TOP = 757_189
GOVERNING_WIDTH = 5_760_404
TITLE_ROW_TOP = 370_849
LAB_GOVERNING_TOP = 720_000
LAB_BODY_TOP = 1_250_000
LAB_TITLE_HEIGHT = 300_000
LAB_GOVERNING_HEIGHT = 380_000
LAB_COLLAPSED_PLACEHOLDER_SIZE = 45_000
TITLE_PT = 18
CONTENT_TITLE_LAYOUTS = frozenset({"내지_거버닝 O", "1_내지_거버닝 X"})
TITLE_PH10_MAX_CHARS = 72
_MIN_BODY_PLACEHOLDER_W = int(Inches(2.0))
_MIN_BODY_PLACEHOLDER_H = int(Inches(0.8))
_CODE_BLOCK_TOKENS = (
    "apiVersion:",
    "kind:",
    "metadata:",
    "spec:",
    "containers:",
    "volumeMounts:",
    "volumes:",
    "ingressClassName:",
    "index.html: |",
)

# Template slide-layout guide roles (see design.md / layout XML placeholders).
#   표지 가이드     → 2_표지
#   간지1 가이드    → 목차 (CONTENTS)
#   간지2 가이드    → 간지 (챕터 타이틀 + 번호)
#   본문 거버닝 O   → 내지_거버닝 O  (ph10 제목, ph12 거버닝 메시지, ph13 본문)
#   본문 거버닝 X   → 1_내지_거버닝 X (ph10 제목, ph12 본문)

_LAYOUT_HANDLERS: dict[str, str] = {
    "2_표지": "_fill_cover",
    "목차": "_fill_contents_guide",
    "간지": "_fill_chapter_guide",
    "내지_거버닝 O": "_fill_body_governing_o",
    "1_내지_거버닝 X": "_fill_body_governing_x",
}

# (seed_geometry, layout_geometry) captured once per build — see build_from_json_specs.
_PLACEHOLDER_GEOM: tuple[dict, dict] | None = None

_MIN_PLACEHOLDER_W_EMU = int(Inches(0.5))
_MIN_PLACEHOLDER_H_EMU = int(Inches(0.15))


def delete_slide(prs: Presentation, index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    sld_id = sld_ids[index]
    prs.part.drop_rel(sld_id.rId)
    sld_id_lst.remove(sld_id)


def strip_all_slides(prs: Presentation) -> None:
    while len(prs.slides) > 0:
        delete_slide(prs, len(prs.slides) - 1)


def find_layout(prs: Presentation, name: str):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    names = [ly.name for master in prs.slide_masters for ly in master.slide_layouts]
    raise KeyError(f"Layout {name!r} not found. Available: {names}")


def _rewrite_relationship_ids(root, rid_map: dict[str, str]) -> None:
    for el in root.iter():
        for k, v in list(el.attrib.items()):
            if v in rid_map:
                el.set(k, rid_map[v])


def duplicate_slide_from_seed(prs: Presentation, seed_slide) -> object:
    pres_part = prs.part
    pkg = pres_part.package
    src_part = seed_slide.part
    partname = pkg.next_partname("/ppt/slides/slide%d.xml")
    new_elm = deepcopy(src_part._element)
    new_part = SlidePart(partname, CT.PML_SLIDE, pkg, new_elm)
    rid_map: dict[str, str] = {}
    for rel in list(src_part.rels.values()):
        if rel.reltype == RT.NOTES_SLIDE:
            continue
        old_rid = rel.rId
        if rel.is_external:
            new_rid = new_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = new_part.relate_to(rel.target_part, rel.reltype)
        rid_map[old_rid] = new_rid
    _rewrite_relationship_ids(new_elm, rid_map)
    r_id = pres_part.relate_to(new_part, RT.SLIDE)
    pres_part._element.get_or_add_sldIdLst().add_sldId(r_id)
    return new_part.slide


def _ensure_cloneable_placeholders(slide) -> None:
    """Add layout placeholders missing on the slide so they stay editable in Normal view."""
    existing = {ph.placeholder_format.idx for ph in slide.placeholders}
    for layout_ph in slide.slide_layout.iter_cloneable_placeholders():
        idx = layout_ph.placeholder_format.idx
        if idx not in existing:
            slide.shapes.clone_placeholder(layout_ph)
            existing.add(idx)


def _placeholder_idx(slide, idx: int):
    """Return placeholder by idx, cloning from layout when the slide copy is missing it."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    for layout_ph in slide.slide_layout.iter_cloneable_placeholders():
        if layout_ph.placeholder_format.idx == idx:
            cloned = slide.shapes.clone_placeholder(layout_ph)
            if cloned is not None:
                return cloned
    raise KeyError(f"Placeholder idx={idx} not found on {slide.slide_layout.name!r}")


def _presentation_for_slide(slide):
    return slide.part.package.presentation_part.presentation


def _placeholder_too_narrow(ph) -> bool:
    return int(ph.width or 0) < _MIN_BODY_PLACEHOLDER_W


def ensure_content_body_placeholder_geometry(slide) -> None:
    """Template starter slides collapse ph12/ph13 (≈0.13in wide) — expand to body band."""
    layout = slide.slide_layout.name
    if layout not in CONTENT_TITLE_LAYOUTS:
        return
    prs = _presentation_for_slide(slide)
    sw, sh = int(prs.slide_width), int(prs.slide_height)
    max_w = sw - BODY_AREA_LEFT - SLIDE_MARGIN_X

    if layout == "1_내지_거버닝 X":
        ph12 = _placeholder_idx(slide, 12)
        body_top = BODY_TOP_MIN
        body_h = sh - body_top - SLIDE_MARGIN_BOTTOM
        if _placeholder_too_narrow(ph12) or int(ph12.top or 0) < HEADER_BOTTOM_EMU:
            ph12.left = BODY_AREA_LEFT
            ph12.top = body_top
            ph12.width = max(max_w, _MIN_BODY_PLACEHOLDER_W)
            ph12.height = max(body_h, _MIN_BODY_PLACEHOLDER_H)
        return

    ph12 = _placeholder_idx(slide, 12)
    if _placeholder_too_narrow(ph12):
        ph12.left = BODY_AREA_LEFT
        ph12.top = GOVERNING_TOP
        ph12.width = min(max_w, GOVERNING_WIDTH)
        ph12.height = max(int(ph12.height or 0), int(Inches(0.45)))

    try:
        ph13 = _placeholder_idx(slide, 13)
    except KeyError:
        return
    if ph13 is None:
        return
    body_top = max(int(ph13.top or 0), BODY_TOP_MIN)
    body_h = sh - body_top - SLIDE_MARGIN_BOTTOM
    if _placeholder_too_narrow(ph13) or int(ph13.height or 0) < _MIN_BODY_PLACEHOLDER_H:
        ph13.left = BODY_AREA_LEFT
        ph13.top = body_top
        ph13.width = max(max_w, _MIN_BODY_PLACEHOLDER_W)
        ph13.height = max(body_h, _MIN_BODY_PLACEHOLDER_H)


def _prepare_slide_for_editing(slide) -> None:
    _restore_slide_placeholder_geometry(slide)
    _ensure_cloneable_placeholders(slide)
    ensure_content_body_placeholder_geometry(slide)


def layout_seed_slides(prs: Presentation) -> dict[str, object]:
    seeds: dict[str, object] = {}
    for slide in prs.slides:
        name = slide.slide_layout.name
        if name not in seeds:
            seeds[name] = slide
    return seeds


def set_tf_font(tf, *, title: bool) -> None:
    name = FONT_TITLE if title else FONT_BODY
    for p in tf.paragraphs:
        p.font.name = name


def configure_text_frame_for_wrap(tf) -> None:
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE


def _estimate_text_width_emu(
    text: str,
    font_pt: float,
    *,
    h_margin_emu: int,
    char_factor: float = 0.55,
    width_padding: float = 1.0,
) -> int:
    """Longest line length → box width (Korean/Latin mixed)."""
    lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
    if not lines:
        lines = [" "]
    longest = max(lines, key=len)
    char_w_emu = max(1, int(font_pt * 914400 / 72 * char_factor))
    raw = h_margin_emu + len(longest) * char_w_emu
    return int(raw * width_padding)


def _autofit_horizontal_textbox(
    shape,
    *,
    font_pt: int,
    min_width_in: float,
    max_width_in: float,
    char_factor: float = 0.55,
    width_padding: float = 1.1,
    word_wrap: bool = False,
) -> None:
    """높이 고정, 가장 긴 줄 기준으로 가로만 조절."""
    tf = shape.text_frame
    fixed_height = int(shape.height)
    tf.word_wrap = word_wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE

    h_margin = int(tf.margin_left) + int(tf.margin_right)
    new_w = _estimate_text_width_emu(
        tf.text,
        float(font_pt),
        h_margin_emu=h_margin,
        char_factor=char_factor,
        width_padding=width_padding,
    )
    shape.width = max(int(Inches(min_width_in)), min(new_w, int(Inches(max_width_in))))
    shape.height = fixed_height


def _autofit_cover_textbox(shape, *, title_font_pt: int) -> None:
    """표지 하단 제목 박스: 높이는 템플릿 고정, 가로만 글자 수에 맞춤."""
    _autofit_horizontal_textbox(
        shape,
        font_pt=title_font_pt,
        min_width_in=2.0,
        max_width_in=10.8,
        char_factor=0.55,
        width_padding=1.08,
        word_wrap=False,
    )


def _ensure_multiline_height(shape, *, font_pt: int, line_count: int) -> None:
    """목차 등 여러 줄: 레이아웃 XML 높이(얇은 띠)보다 본문 줄 수에 맞게 높이 확보."""
    tf = shape.text_frame
    margins = int(tf.margin_top) + int(tf.margin_bottom)
    line_h = int(font_pt * 914400 / 72 * 1.45)
    needed = margins + max(1, line_count) * line_h
    if int(shape.height) < needed:
        shape.height = needed


def _capture_seed_placeholder_geometry(prs: Presentation) -> dict[str, dict[int, tuple[int, int, int, int]]]:
    """Starter-slide positions from the template file (designer-tuned)."""
    out: dict[str, dict[int, tuple[int, int, int, int]]] = {}
    for slide in prs.slides:
        name = slide.slide_layout.name
        if name in out:
            continue
        out[name] = {
            ph.placeholder_format.idx: (int(ph.left), int(ph.top), int(ph.width), int(ph.height))
            for ph in slide.placeholders
        }
    return out


def _capture_layout_placeholder_geometry(prs: Presentation) -> dict[str, dict[int, tuple[int, int, int, int]]]:
    """Layout XML defaults — used only to fix collapsed starter placeholders."""
    out: dict[str, dict[int, tuple[int, int, int, int]]] = {}
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            out[layout.name] = {
                ph.placeholder_format.idx: (int(ph.left), int(ph.top), int(ph.width), int(ph.height))
                for ph in layout.placeholders
            }
    return out


def _merge_seed_layout_box(
    seed: tuple[int, int, int, int],
    layout: tuple[int, int, int, int],
) -> tuple[int, int, int, int]:
    """Keep starter ``left``/``top``; borrow ``width``/``height`` from layout when starter is collapsed."""
    sl, st, sw, sh = seed
    _, _, lw, lh = layout
    # Layout XML also has collapsed guide boxes — only borrow when layout is usable.
    if int(sw) < _MIN_PLACEHOLDER_W_EMU and int(lw) >= _MIN_PLACEHOLDER_W_EMU:
        w = lw
    elif int(sw) < _MIN_PLACEHOLDER_W_EMU:
        w = sw
    else:
        w = sw
    if int(sh) < _MIN_PLACEHOLDER_H_EMU and int(lh) >= _MIN_PLACEHOLDER_H_EMU:
        h = lh
    elif int(sh) < _MIN_PLACEHOLDER_H_EMU:
        h = sh
    else:
        h = sh
    return sl, st, w, h


def _restore_slide_placeholder_geometry(slide) -> None:
    """Apply template starter positions (not bare layout XML defaults)."""
    if _PLACEHOLDER_GEOM is None:
        return
    seed_geoms, layout_geoms = _PLACEHOLDER_GEOM
    name = slide.slide_layout.name
    seed = seed_geoms.get(name, {})
    layout = layout_geoms.get(name, {})
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx not in seed:
            continue
        box = seed[idx]
        if idx in layout:
            box = _merge_seed_layout_box(box, layout[idx])
        ph.left, ph.top, ph.width, ph.height = box


def _ph_by_idx(slide, idx: int):
    """Lookup placeholder; clone from layout when missing (``_prepare_slide_for_editing``)."""
    return _placeholder_idx(slide, idx)


def _non_placeholder_text_shapes(slide) -> list:
    return sorted(
        [sh for sh in slide.shapes if sh.has_text_frame and not sh.is_placeholder],
        key=lambda sh: (sh.top, sh.left),
    )


def _assign_plain(target, text: str, *, title: bool, title_font_pt: int, body_font_pt: int) -> None:
    if not getattr(target, "has_text_frame", False):
        raise ValueError(f"Shape {getattr(target, 'name', target)!r} has no text frame.")
    target.text = text
    tf = target.text_frame
    configure_text_frame_for_wrap(tf)
    set_tf_font(tf, title=title)
    for p in tf.paragraphs:
        p.font.size = Pt(title_font_pt if title else body_font_pt)


def _shape_by_name(slide, name: str):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise KeyError(f"Shape {name!r} not found on layout {slide.slide_layout.name!r}")


def _apply_image_placeholder(slide, item: dict) -> None:
    idx = int(item.get("ph", item.get("idx")))
    path = Path(item["path"]).expanduser()
    if not path.is_file():
        raise FileNotFoundError(f"Image not found: {path}")
    ph = _ph_by_idx(slide, idx)
    if not hasattr(ph, "insert_picture"):
        ph_type = ph.placeholder_format.type
        raise ValueError(
            f"Placeholder ph{idx} on {slide.slide_layout.name!r} cannot accept images "
            f"(type={ph_type})."
        )
    ph.insert_picture(str(path))


def _apply_named_shape(slide, item: dict, *, title_font_pt: int, body_font_pt: int) -> None:
    shape = _shape_by_name(slide, item["name"])
    if "text" not in item:
        return
    is_title = bool(item.get("title", False))
    font_pt = int(item.get("font_pt", title_font_pt if is_title else body_font_pt))
    _assign_plain(shape, item["text"], title=is_title, title_font_pt=font_pt, body_font_pt=body_font_pt)
    if item.get("autofit_cover"):
        _autofit_cover_textbox(shape, title_font_pt=font_pt)


def _apply_spec_overrides(
    slide,
    spec: dict,
    *,
    title_font_pt: int,
    body_font_pt: int,
) -> None:
    for key, value in (spec.get("placeholders") or {}).items():
        ph = _ph_by_idx(slide, int(key))
        is_title = ph.placeholder_format.type in (
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
        )
        _assign_plain(
            ph,
            str(value),
            title=is_title,
            title_font_pt=title_font_pt,
            body_font_pt=body_font_pt,
        )

    for item in spec.get("images") or []:
        _apply_image_placeholder(slide, item)

    for item in spec.get("shapes") or []:
        _apply_named_shape(slide, item, title_font_pt=title_font_pt, body_font_pt=body_font_pt)


def _fill_cover(slide, texts: list[str], *, title_font_pt: int, body_font_pt: int) -> None:
    shapes = _non_placeholder_text_shapes(slide)
    if not shapes:
        raise ValueError(f"Cover layout {slide.slide_layout.name!r} has no text shape.")
    # Primary title box: largest area among free text shapes (template guide).
    shape = max(shapes, key=lambda sh: int(sh.width) * int(sh.height))
    cover_text = texts[0]
    lines = [ln for ln in cover_text.split("\n") if ln.strip()]
    shape.text = cover_text
    tf = shape.text_frame
    configure_text_frame_for_wrap(tf)
    set_tf_font(tf, title=True)
    subtitle_pt = min(body_font_pt, 22)
    for i, p in enumerate(tf.paragraphs):
        if i == 0:
            p.font.size = Pt(title_font_pt)
        elif i == 1:
            p.font.size = Pt(subtitle_pt)
        else:
            p.font.size = Pt(body_font_pt)
    if len(lines) >= 2:
        _ensure_multiline_height(shape, font_pt=title_font_pt, line_count=len(lines))
        tf.word_wrap = True
        _autofit_horizontal_textbox(
            shape,
            font_pt=title_font_pt,
            min_width_in=2.0,
            max_width_in=10.8,
            char_factor=0.55,
            width_padding=1.08,
            word_wrap=True,
        )
    else:
        _autofit_cover_textbox(shape, title_font_pt=title_font_pt)


def _fill_contents_guide(slide, texts: list[str], *, title_font_pt: int, body_font_pt: int) -> None:
    """간지1 가이드 — ``목차`` 레이아웃, ph10 단일 목차 블록."""
    if len(texts) != 1:
        raise ValueError(f"목차 expects 1 text block, got {len(texts)}")
    font_pt = 18
    ph = _ph_by_idx(slide, 10)
    lines = [ln for ln in texts[0].split("\n") if ln.strip()]
    _ensure_multiline_height(ph, font_pt=font_pt, line_count=len(lines))
    ph.text = texts[0]
    configure_text_frame_for_wrap(ph.text_frame)
    set_tf_font(ph.text_frame, title=False)
    for p in ph.text_frame.paragraphs:
        p.font.size = Pt(font_pt)
    _autofit_horizontal_textbox(
        ph,
        font_pt=font_pt,
        min_width_in=4.2,
        max_width_in=11.5,
        char_factor=0.62,
        width_padding=1.14,
        word_wrap=False,
    )


def _fill_chapter_guide(slide, texts: list[str], *, title_font_pt: int, body_font_pt: int) -> None:
    """간지2 가이드 — ``간지`` 레이아웃, ph10 타이틀 + ph11 챕터 번호."""
    if len(texts) != 2:
        raise ValueError(f"간지 expects 2 text blocks [title, number], got {len(texts)}")
    _assign_plain(_ph_by_idx(slide, 10), texts[0], title=True, title_font_pt=36, body_font_pt=body_font_pt)
    _assign_plain(_ph_by_idx(slide, 11), texts[1], title=False, title_font_pt=title_font_pt, body_font_pt=20)


def _fill_body_governing_o(slide, texts: list[str], *, title_font_pt: int, body_font_pt: int) -> None:
    """본문 거버닝 O — ph10 제목, ph12 거버닝 메시지(2줄 이내), ph13 본문."""
    if len(texts) != 3:
        raise ValueError(f"내지_거버닝 O expects 3 text blocks [title, subtitle, body], got {len(texts)}")
    _assign_plain(_ph_by_idx(slide, 10), texts[0], title=True, title_font_pt=title_font_pt, body_font_pt=body_font_pt)
    _assign_plain(_ph_by_idx(slide, 12), texts[1], title=False, title_font_pt=title_font_pt, body_font_pt=14)
    _assign_plain(_ph_by_idx(slide, 13), texts[2], title=False, title_font_pt=title_font_pt, body_font_pt=body_font_pt)


def _fill_body_governing_x(slide, texts: list[str], *, title_font_pt: int, body_font_pt: int) -> None:
    """본문 거버닝 X — ph10 제목, ph12 본문(가이드 칸)."""
    if len(texts) != 2:
        raise ValueError(f"1_내지_거버닝 X expects 2 text blocks [title, body], got {len(texts)}")
    _assign_plain(_ph_by_idx(slide, 10), texts[0], title=True, title_font_pt=title_font_pt, body_font_pt=body_font_pt)
    _assign_plain(_ph_by_idx(slide, 12), texts[1], title=False, title_font_pt=title_font_pt, body_font_pt=body_font_pt)


def _is_lab_lecture_spec(spec: dict) -> bool:
    return (spec.get("_ingest") or {}).get("layout_profile") == "lab_lecture"


def _is_lab_visual_preservation_spec(spec: dict) -> bool:
    return (spec.get("_ingest") or {}).get("visual_preservation") == "native_shapes"


def _collapse_lab_body_placeholders(slide) -> None:
    prs = _presentation_for_slide(slide)
    left = int(prs.slide_width) + SLIDE_MARGIN_X
    top = int(prs.slide_height) + SLIDE_MARGIN_BOTTOM
    body_indices = (12, 13) if slide.slide_layout.name == "내지_거버닝 O" else (12,)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx not in body_indices:
            continue
        offset = LAB_COLLAPSED_PLACEHOLDER_SIZE if shape.placeholder_format.idx == 13 else 0
        shape.text = ""
        shape.left = left
        shape.top = top + offset
        shape.width = LAB_COLLAPSED_PLACEHOLDER_SIZE
        shape.height = LAB_COLLAPSED_PLACEHOLDER_SIZE
        shape.fill.background()
        _set_shape_line_no_fill(shape)


def _set_shape_line_no_fill(shape) -> None:
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ln = shape._element.find(f".//{{{a_ns}}}ln")
    if ln is None:
        return
    for child in list(ln):
        if child.tag in (
            f"{{{a_ns}}}noFill",
            f"{{{a_ns}}}solidFill",
            f"{{{a_ns}}}gradFill",
            f"{{{a_ns}}}pattFill",
            f"{{{a_ns}}}blipFill",
        ):
            ln.remove(child)
    no_fill = shape._element.makeelement(f"{{{a_ns}}}noFill")
    ln.insert(0, no_fill)


def _apply_lab_body_typography(shape, *, body_font_pt: int) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    configure_text_frame_for_wrap(tf)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for para in tf.paragraphs:
        para.font.name = FONT_BODY
        para.font.size = Pt(body_font_pt)
        para.line_spacing = 1.08
        para.space_after = Pt(3)


def _apply_lab_lecture_content_layout(slide, spec: dict, *, body_font_pt: int) -> None:
    if not _is_lab_lecture_spec(spec) or slide.slide_layout.name not in CONTENT_TITLE_LAYOUTS:
        return
    prs = _presentation_for_slide(slide)
    sw, sh = int(prs.slide_width), int(prs.slide_height)
    body_w = max(sw - BODY_AREA_LEFT - SLIDE_MARGIN_X, _MIN_BODY_PLACEHOLDER_W)
    body_h = max(sh - LAB_BODY_TOP - SLIDE_MARGIN_BOTTOM, _MIN_BODY_PLACEHOLDER_H)

    ph10 = _ph_by_idx(slide, 10)
    ph10.left = LAYOUT_PH10_LEFT
    ph10.top = TITLE_ROW_TOP
    ph10.width = min(sw - LAYOUT_PH10_LEFT - SLIDE_MARGIN_X, int(Inches(10.5)))
    ph10.height = max(int(ph10.height or 0), LAB_TITLE_HEIGHT)

    if _is_lab_visual_preservation_spec(spec):
        _collapse_lab_body_placeholders(slide)
        return

    if slide.slide_layout.name == "내지_거버닝 O":
        ph12 = _ph_by_idx(slide, 12)
        ph12.left = BODY_AREA_LEFT
        ph12.top = LAB_GOVERNING_TOP
        ph12.width = min(body_w, GOVERNING_WIDTH)
        ph12.height = max(int(ph12.height or 0), LAB_GOVERNING_HEIGHT)
        _apply_lab_body_typography(ph12, body_font_pt=14)

        ph13 = _ph_by_idx(slide, 13)
        ph13.left = BODY_AREA_LEFT
        ph13.top = LAB_BODY_TOP
        ph13.width = body_w
        ph13.height = body_h
        _apply_lab_body_typography(ph13, body_font_pt=max(14, body_font_pt - 1))
        return

    ph12 = _ph_by_idx(slide, 12)
    ph12.left = BODY_AREA_LEFT
    ph12.width = body_w
    ph12.height = max(sh - int(ph12.top or 0) - SLIDE_MARGIN_BOTTOM, _MIN_BODY_PLACEHOLDER_H)
    ingest = spec.get("_ingest") or {}
    if ingest.get("code_block") and ingest.get("code_block_shape_preservation"):
        ph12.top = LAB_GOVERNING_TOP
        ph12.height = LAB_GOVERNING_HEIGHT
        _apply_lab_body_typography(ph12, body_font_pt=14)
    elif not ingest.get("code_block"):
        ph12.top = LAB_BODY_TOP
        ph12.height = body_h
        _apply_lab_body_typography(ph12, body_font_pt=max(14, body_font_pt - 1))


def apply_lab_lecture_layouts(prs: Presentation, specs: list[dict]) -> None:
    for slide, spec in zip(prs.slides, specs):
        _apply_lab_lecture_content_layout(slide, spec, body_font_pt=16)


def fill_slide_from_spec(
    slide,
    spec: dict,
    *,
    title_font_pt: int = 28,
    body_font_pt: int = 16,
) -> None:
    """Fill slide from JSON spec; preserve unmentioned template shapes for editing."""
    _prepare_slide_for_editing(slide)

    texts = spec.get("texts")
    if texts is not None:
        if isinstance(texts, str):
            texts = [texts]
        texts = [t if isinstance(t, str) else str(t) for t in texts]
        layout_name = slide.slide_layout.name
        handler_name = _LAYOUT_HANDLERS.get(layout_name)
        if handler_name is None:
            raise ValueError(
                f"No placeholder guide for layout {layout_name!r}. "
                f"Supported: {sorted(_LAYOUT_HANDLERS)}"
            )
        handler = globals()[handler_name]
        handler(slide, texts, title_font_pt=title_font_pt, body_font_pt=body_font_pt)
        if layout_name in CONTENT_TITLE_LAYOUTS:
            rebalance_content_placeholders(slide)
            ensure_content_body_placeholder_geometry(slide)
            _apply_lab_lecture_content_layout(slide, spec, body_font_pt=body_font_pt)
            if (spec.get("_ingest") or {}).get("code_block"):
                _apply_code_block_text_style(slide)

    _apply_spec_overrides(slide, spec, title_font_pt=title_font_pt, body_font_pt=body_font_pt)


def _apply_code_block_text_style(slide) -> None:
    body_indices = (12, 13) if slide.slide_layout.name == "내지_거버닝 O" else (12,)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx not in body_indices:
            continue
        if not shape.has_text_frame:
            continue
        text = (shape.text or "").strip()
        if not text:
            continue
        shape.text_frame.word_wrap = True
        shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        for para in shape.text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            for run in para.runs:
                run.font.name = "Courier New"
                run.font.size = Pt(10)


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return (shape.text or "").replace("\x0b", "\n").strip()


def _is_lab_code_text(text: str) -> bool:
    return any(token in text for token in _CODE_BLOCK_TOKENS)


def _is_lab_footer_text(text: str) -> bool:
    compact = " ".join(text.split())
    if "|" not in compact:
        return False
    prefix, suffix = compact.rsplit("|", 1)
    return bool(prefix.strip()) and suffix.strip().isdigit()


def _source_lab_code_shapes(source_slide) -> list:
    code_shapes = []
    for shape in source_slide.shapes:
        text = _shape_text(shape)
        if text and _is_lab_code_text(text):
            code_shapes.append(shape)
    code_shapes.sort(key=lambda shape: (int(shape.top or 0), int(shape.left or 0)))
    return code_shapes


def _source_lab_prose_text(source_slide, code_shapes: list) -> str:
    code_ids = {id(shape) for shape in code_shapes}
    candidates: list[tuple[int, int, str]] = []
    for shape in source_slide.shapes:
        if id(shape) in code_ids:
            continue
        text = _shape_text(shape)
        if not text or _is_lab_code_text(text) or _is_lab_footer_text(text):
            continue
        top = int(shape.top or 0)
        if top <= TITLE_ROW_TOP + LAB_TITLE_HEIGHT:
            continue
        candidates.append((top, int(shape.left or 0), text))
    return "\n".join(text for _top, _left, text in sorted(candidates)).strip()


def _apply_lab_code_placeholder_text(dst_slide, source_slide, code_shapes: list) -> None:
    prose = _source_lab_prose_text(source_slide, code_shapes)
    try:
        body = _ph_by_idx(dst_slide, 13 if dst_slide.slide_layout.name == "내지_거버닝 O" else 12)
    except KeyError:
        return
    body.text = prose
    if prose:
        _apply_lab_body_typography(body, body_font_pt=14)
    else:
        body.fill.background()
        _set_shape_line_no_fill(body)


def apply_lab_code_block_shapes(
    prs: Presentation,
    specs: list[dict],
    source_prs: Presentation,
) -> int:
    applied = 0
    for slide, spec in zip(prs.slides, specs):
        ingest = spec.get("_ingest") or {}
        if not (_is_lab_lecture_spec(spec) and ingest.get("code_block")):
            continue
        source_slide = _source_slide_for_spec(source_prs, spec)
        if source_slide is None:
            continue
        code_shapes = _source_lab_code_shapes(source_slide)
        if len(code_shapes) < 2:
            continue
        _apply_lab_code_placeholder_text(slide, source_slide, code_shapes)
        for shape in code_shapes:
            before = len(slide.shapes)
            copy_shape_hybrid(slide, shape)
            if len(slide.shapes) > before:
                applied += 1
        if code_shapes:
            renumber_shape_ids(slide)
    return applied


def fill_slide_text_placeholders(
    slide,
    texts: list[str],
    *,
    title_font_pt: int = 28,
    body_font_pt: int = 16,
) -> None:
    """Map JSON ``texts`` to template placeholders per layout guide."""
    fill_slide_from_spec(slide, {"texts": texts}, title_font_pt=title_font_pt, body_font_pt=body_font_pt)


def apply_kubernetes_sample_five_slides(prs: Presentation) -> None:
    global _PLACEHOLDER_GEOM
    _PLACEHOLDER_GEOM = (_capture_seed_placeholder_geometry(prs), _capture_layout_placeholder_geometry(prs))
    try:
        _apply_kubernetes_sample_five_slides(prs)
    finally:
        _PLACEHOLDER_GEOM = None


def _apply_kubernetes_sample_five_slides(prs: Presentation) -> None:
    fill_slide_text_placeholders(
        prs.slides[0],
        ["Kubernetes 입문 실습\n아카데미 샘플 강의안 (자동 생성)"],
    )
    fill_slide_text_placeholders(
        prs.slides[1],
        [
            "1. 개요 — 컨테이너와 클러스터\n"
            "2. Deployment — 원하는 상태 선언\n"
            "3. Service / Ingress — 노출과 라우팅\n"
            "4. 실습 체크리스트\n"
            "5. 정리 및 참고",
        ],
    )
    fill_slide_text_placeholders(prs.slides[2], ["개요", "1."])
    fill_slide_text_placeholders(
        prs.slides[3],
        [
            "Kubernetes 한눈에 보기",
            "핵심 개념",
            "• 컨테이너: 애플리케이션과 라이브러리를 묶은 실행 단위\n"
            "• Pod: 컨테이너를 감싼 최소 배포 단위\n"
            "• Deployment: 원하는 Pod 수·이미지 버전을 선언하고 유지",
        ],
    )
    fill_slide_text_placeholders(
        prs.slides[4],
        [
            "Deployment → Service → Ingress",
            "Nginx 이미지로 Deployment를 만들고, Service로 안정적인 엔드포인트를 만든 뒤, "
            "Ingress로 외부 HTTP 접속을 구성합니다. 각 단계 후 "
            "“Welcome to nginx!” 페이지로 동작을 확인합니다.",
        ],
    )

    while len(prs.slides) > 5:
        delete_slide(prs, len(prs.slides) - 1)


def apply_background_images(prs: Presentation, specs: Iterable[dict]) -> int:
    """Place exported slide PNGs on ``1_내지_거버닝 X`` slides (Google export fallback)."""
    applied = 0
    for slide, spec in zip(prs.slides, specs):
        rel = spec.get("background_image")
        if not rel or slide.slide_layout.name not in ("1_내지_거버닝 X", "내지_거버닝 O"):
            continue
        path = Path(rel)
        if not path.is_file():
            continue
        sw, sh = int(prs.slide_width), int(prs.slide_height)
        margin_x = int(Inches(0.5))
        top = int(Inches(1.35))
        pic = slide.shapes.add_picture(
            str(path),
            margin_x,
            top,
            width=sw - 2 * margin_x,
            height=sh - top - int(Inches(0.5)),
        )
        pic.name = f"ingest_bg_{applied}"
        applied += 1
    return applied


def _should_preserve_visual_shapes(spec: dict) -> bool:
    return _is_lab_visual_preservation_spec(spec)


def _source_slide_for_spec(source_prs: Presentation, spec: dict):
    index = (spec.get("_ingest") or {}).get("source_slide_index")
    if not isinstance(index, int):
        return None
    if index < 0 or index >= len(source_prs.slides):
        return None
    return source_prs.slides[index]


def _remove_non_placeholder_shapes(slide) -> None:
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        shape._element.getparent().remove(shape._element)


def _clear_content_body_placeholders(slide) -> None:
    _collapse_lab_body_placeholders(slide)


def _shape_starts_below_header(shape) -> bool:
    return int(shape.top or 0) >= LAB_BODY_TOP - 180_000


def _copy_visual_shapes(src_slide, dst_slide, *, cover: bool) -> int:
    copied = 0
    for shape in src_slide.shapes:
        if shape.is_placeholder:
            continue
        if not cover and not _shape_starts_below_header(shape):
            continue
        before = len(dst_slide.shapes)
        copy_shape_hybrid(dst_slide, shape)
        if len(dst_slide.shapes) > before:
            copied += 1
    if copied:
        renumber_shape_ids(dst_slide)
    return copied


def apply_lab_visual_shapes(
    prs: Presentation,
    specs: list[dict],
    source_prs: Presentation,
) -> int:
    applied = 0
    for slide, spec in zip(prs.slides, specs):
        if not _should_preserve_visual_shapes(spec):
            continue
        source_slide = _source_slide_for_spec(source_prs, spec)
        if source_slide is None:
            continue
        cover = slide.slide_layout.name == "2_표지"
        _remove_non_placeholder_shapes(slide)
        _clear_content_body_placeholders(slide)
        applied += _copy_visual_shapes(source_slide, slide, cover=cover)
    return applied


def build_from_json_specs(prs: Presentation, specs: Iterable[dict]) -> None:
    global _PLACEHOLDER_GEOM
    specs_list = list(specs)
    seeds = layout_seed_slides(prs)
    initial_part_ids = {id(s.part) for s in prs.slides}
    _PLACEHOLDER_GEOM = (_capture_seed_placeholder_geometry(prs), _capture_layout_placeholder_geometry(prs))

    try:
        _build_from_json_specs_loop(prs, specs_list, seeds, initial_part_ids)
        polish_academy_presentation(prs)
        apply_lab_lecture_layouts(prs, specs_list)
    finally:
        _PLACEHOLDER_GEOM = None


def _build_from_json_specs_loop(
    prs: Presentation,
    specs_list: list[dict],
    seeds: dict,
    initial_part_ids: set[int],
) -> None:
    for spec in specs_list:
        layout_name = spec["layout"]
        if (
            spec.get("texts") is None
            and not spec.get("placeholders")
            and not spec.get("images")
            and not spec.get("shapes")
        ):
            raise ValueError(
                f"Slide spec for {layout_name!r} needs 'texts' and/or "
                "'placeholders' / 'images' / 'shapes'."
            )
        find_layout(prs, layout_name)
        try:
            seed = seeds[layout_name]
        except KeyError as e:
            available = sorted(seeds.keys())
            raise KeyError(
                f"No starter slide in this template for layout {layout_name!r}. "
                f"Template includes starters for: {available}"
            ) from e
        slide = duplicate_slide_from_seed(prs, seed)
        fill_slide_from_spec(slide, spec)

    to_remove = [i for i, s in enumerate(prs.slides) if id(s.part) in initial_part_ids]
    for i in reversed(to_remove):
        delete_slide(prs, i)


# Backward-compatible alias for tests/scripts that imported sync.
def sync_placeholder_geometry_from_layout(slide) -> None:
    _restore_slide_placeholder_geometry(slide)


def normalize_slide_title(title: str) -> str:
    return " ".join((title or "").split())


def _content_ph(slide, idx: int):
    """Placeholder lookup without re-applying seed geometry."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(f"Placeholder idx={idx} not found on {slide.slide_layout.name!r}")


def rebalance_content_placeholders(slide) -> None:
    """Move overflow out of ph10 into body placeholder (spec ingest safety net)."""
    if slide.slide_layout.name not in CONTENT_TITLE_LAYOUTS:
        return
    try:
        ph10 = _ph_by_idx(slide, 10)
    except KeyError:
        return
    raw = (ph10.text or "").replace("\x0b", "\n").strip()
    if not raw:
        return
    lines = [ln.strip() for ln in raw.split("\n") if ln.strip()]
    title_line = lines[0] if lines else raw
    overflow = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""
    if len(title_line) > TITLE_PH10_MAX_CHARS:
        cut = title_line[:TITLE_PH10_MAX_CHARS]
        if " " in cut:
            cut = cut.rsplit(" ", 1)[0]
        overflow = (title_line[len(cut) :].strip() + "\n" + overflow).strip()
        title_line = cut.strip()
    if not overflow and len(title_line) <= TITLE_PH10_MAX_CHARS:
        assign_content_slide_title(ph10, title_line)
        return
    assign_content_slide_title(ph10, title_line)
    body_idx = 13 if slide.slide_layout.name == "내지_거버닝 O" else 12
    try:
        body_ph = _ph_by_idx(slide, body_idx)
    except KeyError:
        return
    existing = (body_ph.text or "").strip()
    merged = overflow
    if existing and existing not in (" ", "거버닝 메시지…"):
        merged = f"{overflow}\n\n{existing}".strip() if overflow else existing
    if merged:
        _assign_plain(
            body_ph,
            merged,
            title=False,
            title_font_pt=14,
            body_font_pt=16,
        )


def apply_content_header_geometry(slide, slide_width: int | None = None) -> None:
    """Official template ph10/ph12 header positions (not full-slide PaaS width)."""
    ensure_content_body_placeholder_geometry(slide)
    ph10 = _content_ph(slide, 10)
    ph10.left = LAYOUT_PH10_LEFT
    if slide.slide_layout.name == "내지_거버닝 O":
        ph12 = _placeholder_idx(slide, 12)
        ph10.top = ph12.top
    else:
        ph10.top = TITLE_ROW_TOP
    if slide_width:
        max_w = slide_width - LAYOUT_PH10_LEFT - SLIDE_MARGIN_X
        ph10.width = min(int(max_w), int(Inches(10.5)))


def assign_content_slide_title(ph10, title: str) -> None:
    """ph10: one line, 18pt runs, left align, vertical middle (v9 CONTRABASS)."""
    title = normalize_slide_title(title)
    tf = ph10.text_frame
    configure_text_frame_for_wrap(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ph10.text = ""
    while len(tf.paragraphs) > 1:
        tf._element.remove(tf.paragraphs[-1]._p)

    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    para.clear()

    if ": " in title:
        head, tail = title.split(": ", 1)
        r_ko = para.add_run()
        r_ko.text = f"{head}: "
        r_ko.font.name = FONT_TITLE
        r_ko.font.size = Pt(TITLE_PT)
        r_ko.font.color.rgb = CLR_DK1
        r_en = para.add_run()
        r_en.text = tail
        r_en.font.name = FONT_TITLE
        r_en.font.size = Pt(TITLE_PT)
        r_en.font.color.rgb = CLR_DK1
    else:
        run = para.add_run()
        run.text = title
        run.font.name = FONT_TITLE
        run.font.size = Pt(TITLE_PT)
        run.font.color.rgb = CLR_DK1


def apply_slide_title_layout(slide, slide_width: int, title: str) -> None:
    apply_content_header_geometry(slide, slide_width)
    ph10 = _content_ph(slide, 10)
    assign_content_slide_title(ph10, title)
    max_w = slide_width - LAYOUT_PH10_LEFT - SLIDE_MARGIN_X
    ph10.width = min(int(max_w), int(Inches(10.5)))


def hide_empty_governing_placeholder(slide) -> None:
    """Clear template governing guide on ph12; never remove ph13 (본문) from the slide."""
    if slide.slide_layout.name != "내지_거버닝 O":
        return
    guide_hint = "거버닝 메시지"
    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 12:
            continue
        text = (shape.text or "").strip()
        if not text or guide_hint in text:
            shape.text = ""
        return


def polish_academy_presentation(prs: Presentation) -> None:
    """Apply content title styling to all 본문 layout slides."""
    global _PLACEHOLDER_GEOM
    prev = _PLACEHOLDER_GEOM
    try:
        if _PLACEHOLDER_GEOM is None:
            _PLACEHOLDER_GEOM = (
                _capture_seed_placeholder_geometry(prs),
                _capture_layout_placeholder_geometry(prs),
            )
        sw = int(prs.slide_width)
        for slide in prs.slides:
            if slide.slide_layout.name not in CONTENT_TITLE_LAYOUTS:
                continue
            title = ""
            for sh in slide.placeholders:
                if sh.placeholder_format.idx == 10:
                    title = normalize_slide_title(sh.text or "")
            ensure_content_body_placeholder_geometry(slide)
            if title:
                apply_slide_title_layout(slide, sw, title)
            hide_empty_governing_placeholder(slide)
    finally:
        _PLACEHOLDER_GEOM = prev


def fix_open_in_slide_view(path: Path) -> None:
    """Open in normal slide view, not slide master (ppt/viewProps.xml)."""
    import re
    import zipfile

    path = Path(path)
    tmp = path.with_suffix(".view.zip")
    with zipfile.ZipFile(path, "r") as zin:
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "ppt/viewProps.xml":
                    xml = data.decode("utf-8")
                    xml = re.sub(
                        r'lastView="(?:sldMasterView|slideMasterView)"',
                        'lastView="sldView"',
                        xml,
                    )
                    data = xml.encode("utf-8")
                zout.writestr(item, data)
    tmp.replace(path)


def save_academy_deck(prs: Presentation, path: Path, specs: list[dict] | None = None) -> None:
    """Save with title polish + repair-safe round-trip + slide view fix."""
    path = Path(path)
    polish_academy_presentation(prs)
    if specs is not None:
        apply_lab_lecture_layouts(prs, specs)
    tmp = path.with_suffix(".tmp.pptx")
    prs.save(str(tmp))
    Presentation(str(tmp)).save(str(path))
    tmp.unlink(missing_ok=True)
    fix_open_in_slide_view(path)
    finalize_academy_package(path)


def finalize_academy_package(path: Path) -> int:
    """Strip Mac PowerPoint 'Repair' triggers (fonts/changesInfo, SVG→PNG, slide renumber)."""
    path = Path(path)
    from scripts.sanitize_pptx_package import finalize_pptx_package

    svg_n = finalize_pptx_package(path)
    fix_open_in_slide_view(path)
    return svg_n
