#!/usr/bin/env python3
"""Quick check: placeholder widths and text on selected slides (1-based index)."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

_SLIVER = 280000


def main() -> int:
    path = Path(sys.argv[1])
    indices = [int(x) for x in sys.argv[2:]] if len(sys.argv) > 2 else list(range(1, len(Presentation(path).slides) + 1))
    prs = Presentation(path)
    for i in indices:
        if i < 1 or i > len(prs.slides):
            continue
        slide = prs.slides[i - 1]
        print(f"\n=== slide {i} layout={slide.slide_layout.name!r} ===")
        for ph in sorted(slide.placeholders, key=lambda p: p.placeholder_format.idx):
            w = int(ph.width)
            flag = " SLIVER" if w < _SLIVER else ""
            txt = (ph.text or "").replace("\n", " | ")[:80]
            print(f"  ph{ph.placeholder_format.idx}: w={w/914400:.2f}in{flag}  {txt!r}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
