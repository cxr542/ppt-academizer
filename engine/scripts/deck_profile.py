#!/usr/bin/env python3
"""Detect which academize pipeline fits a source deck (academy-design §5 vs §7)."""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from scripts.deck_migrate_config import config_for_kind, detect_deck_kind
from scripts.pptx_ingest import is_image_only_slide, slide_text_blocks

Profile = str  # "spec" | "migrate_cmp" | "google_image"

_AI_FILENAME_RE = re.compile(
    r"chatgpt|openai|gpt[-_]?4|gemini|google\s*ai|copilot|claude|anthropic|"
    r"gamma\.?app|slidesai|slide\.ai|beautiful\.ai|tome\.app|"
    r"ai[-_ ]?(deck|ppt|slide|export|generated)|generated[-_ ]?by",
    re.I,
)


@dataclass(frozen=True)
class DeckStructure:
    slide_count: int
    shape_count: int
    shapes_per_slide: float
    chart_count: int
    table_count: int
    group_count: int
    picture_count: int
    google_image_ratio: float
    avg_text_blocks: float
    total_text_blocks: int
    text_slide_ratio: float


def _walk_shapes(shapes, counts: dict[str, int]) -> None:
    for sh in shapes:
        st = sh.shape_type
        if st == MSO_SHAPE_TYPE.GROUP:
            counts["groups"] += 1
            _walk_shapes(sh.shapes, counts)
            continue
        if st == MSO_SHAPE_TYPE.CHART:
            counts["charts"] += 1
        elif st == MSO_SHAPE_TYPE.TABLE:
            counts["tables"] += 1
        elif st == MSO_SHAPE_TYPE.PICTURE:
            counts["pictures"] += 1
        counts["shapes"] += 1


def analyze_deck_structure(prs: Presentation) -> DeckStructure:
    """Content signals for §5 vs §7 routing (not filename)."""
    slide_count = len(prs.slides)
    if slide_count == 0:
        return DeckStructure(0, 0, 0.0, 0, 0, 0, 0, 0.0, 0.0, 0, 0.0)

    counts = {"shapes": 0, "charts": 0, "tables": 0, "groups": 0, "pictures": 0}
    image_slides = 0
    text_slides = 0
    total_blocks = 0

    for slide in prs.slides:
        _walk_shapes(slide.shapes, counts)
        blocks = slide_text_blocks(slide)
        if blocks:
            text_slides += 1
            total_blocks += len(blocks)
        if is_image_only_slide(slide):
            image_slides += 1

    shape_count = counts["shapes"]
    return DeckStructure(
        slide_count=slide_count,
        shape_count=shape_count,
        shapes_per_slide=shape_count / slide_count,
        chart_count=counts["charts"],
        table_count=counts["tables"],
        group_count=counts["groups"],
        picture_count=counts["pictures"],
        google_image_ratio=image_slides / slide_count,
        avg_text_blocks=total_blocks / slide_count,
        total_text_blocks=total_blocks,
        text_slide_ratio=text_slides / slide_count,
    )


def _structure_to_meta(st: DeckStructure, extra: dict | None = None) -> dict:
    meta = {
        "slide_count": st.slide_count,
        "google_image_ratio": round(st.google_image_ratio, 3),
        "avg_text_blocks": round(st.avg_text_blocks, 2),
        "shape_count": st.shape_count,
        "shapes_per_slide": round(st.shapes_per_slide, 2),
        "total_text_blocks": st.total_text_blocks,
        "chart_count": st.chart_count,
        "table_count": st.table_count,
        "group_count": st.group_count,
        "picture_count": st.picture_count,
        "text_slide_ratio": round(st.text_slide_ratio, 3),
    }
    if extra:
        meta.update(extra)
    return meta


def is_partner_shape_heavy(st: DeckStructure) -> bool:
    """Partner / CONTRABASS-style: many shapes, diagrams, pictures — use §7 migrate."""
    if st.slide_count < 8:
        return False
    if st.google_image_ratio >= 0.45:
        return False
    if st.shapes_per_slide >= 4.0:
        return True
    if st.shape_count >= st.slide_count * 3 and st.avg_text_blocks >= 1.2:
        return True
    if st.shapes_per_slide >= 2.5 and (
        st.table_count >= 1
        or st.chart_count >= 1
        or st.group_count >= 2
        or st.picture_count >= st.slide_count * 2
    ):
        return True
    return False


def ai_filename_hint(filename: str) -> bool:
    return bool(_AI_FILENAME_RE.search(filename or ""))


def is_ai_freeform_text_deck(st: DeckStructure) -> bool:
    """ChatGPT/Copilot-style: blank slides + text boxes, few diagrams."""
    if st.slide_count < 2 or st.slide_count > 80:
        return False
    if st.google_image_ratio >= 0.35:
        return False
    if st.chart_count or st.table_count:
        return False
    if st.shapes_per_slide > 7.0:
        return False
    if st.text_slide_ratio < 0.8:
        return False
    if st.avg_text_blocks < 1.0:
        return False
    if st.total_text_blocks < max(2, int(st.slide_count * 0.9)):
        return False
    if st.shapes_per_slide >= 3.5:
        return False
    return True


def is_text_lecture_deck(st: DeckStructure) -> bool:
    """Simple text / Google export — use §5 spec."""
    if st.google_image_ratio >= 0.4 and st.total_text_blocks <= st.slide_count:
        return True
    if (
        st.slide_count <= 20
        and st.shapes_per_slide < 2.0
        and st.chart_count == 0
        and st.table_count == 0
        and st.group_count == 0
        and st.total_text_blocks >= max(2, st.slide_count // 3)
    ):
        return True
    return False


def detect_deck_profile(
    source: Path | Presentation,
    *,
    filename_hint: str | None = None,
) -> tuple[Profile, dict]:
    """Return (profile, metadata) for routing academize pipeline."""
    if isinstance(source, (str, Path)):
        path = Path(source)
        prs = Presentation(str(path))
        name = path.name.lower()
    else:
        prs = source
        name = (filename_hint or "").lower()

    st = analyze_deck_structure(prs)
    deck_kind = detect_deck_kind(name)

    if st.google_image_ratio >= 0.5 and st.total_text_blocks == 0:
        return "google_image", _structure_to_meta(st, {"reason": "majority_google_image_slides"})

    if is_partner_shape_heavy(st):
        cfg = config_for_kind(deck_kind)
        return "migrate_cmp", _structure_to_meta(
            st,
            {
                "deck_kind": deck_kind,
                "auto_toc_after_cover": cfg.auto_toc_after_cover,
                "reason": "partner_shape_heavy",
            },
        )

    if is_text_lecture_deck(st):
        return "spec", _structure_to_meta(st, {"reason": "text_lecture_structure"})

    ai_hint = ai_filename_hint(name)
    if is_ai_freeform_text_deck(st) or (
        ai_hint
        and st.text_slide_ratio >= 0.65
        and st.google_image_ratio < 0.45
        and st.shapes_per_slide < 3.5
    ):
        return "spec", _structure_to_meta(
            st,
            {
                "reason": "ai_freeform_export",
                "ai_filename_hint": ai_hint,
            },
        )

    if st.total_text_blocks >= max(3, st.slide_count // 2) and st.google_image_ratio < 0.85:
        if st.shapes_per_slide >= 3.0 and st.slide_count >= 8:
            cfg = config_for_kind(deck_kind)
            return "migrate_cmp", _structure_to_meta(
                st,
                {
                    "deck_kind": deck_kind,
                    "auto_toc_after_cover": cfg.auto_toc_after_cover,
                    "reason": "dense_shapes_fallback",
                },
            )
        return "spec", _structure_to_meta(st, {"reason": "extractable_text_blocks"})

    return "spec", _structure_to_meta(st, {"reason": "default_spec_heuristic"})
