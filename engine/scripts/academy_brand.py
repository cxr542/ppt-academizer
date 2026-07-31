"""Academy brand pack: logos, theme colors, icon catalog (from template + elements)."""
from __future__ import annotations

import json
import os
import re
from functools import lru_cache
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

# McKinsey / partner decks often use cyan accents; map to academy accent1.
_NON_BRAND_ACCENTS = frozenset(
    {
        "00A3E0",
        "00B0F0",
        "00AEEF",
        "05A3E0",
    }
)

_ICON_MIN = 180_000
_ICON_MAX = 2_400_000
_LABEL_MAX_CHARS = 48
_NEAR_EMU = 1_800_000


def _hex_to_rgb(value: str) -> RGBColor:
    h = value.strip().lstrip("#").upper()
    if len(h) != 6:
        raise ValueError(f"invalid hex color: {value!r}")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def resolve_brand_dir(template_path: Path | None = None) -> Path | None:
    """Locate brand/ directory next to TEMPLATE_PPTX or via BRAND_DIR."""
    env = os.environ.get("BRAND_DIR", "").strip()
    if env:
        p = Path(env).expanduser()
        if p.is_dir() and (p / "colors.json").is_file():
            return p
    if template_path is not None:
        cand = Path(template_path).expanduser().resolve().parent / "brand"
        if cand.is_dir() and (cand / "colors.json").is_file():
            return cand
    local = Path.home() / ".config" / "ppt-academizer-render" / "brand"
    if local.is_dir() and (local / "colors.json").is_file():
        return local
    return None


def _display_label(text: str) -> str:
    return (text or "").replace("\x0b", " ").replace("\n", " ").strip()


def _ascii_key(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", "", (text or "").lower())


def _space_key(text: str) -> str:
    return re.sub(r"\s+", " ", _display_label(text).lower())


@lru_cache(maxsize=4)
def load_brand_pack(brand_dir: str) -> dict:
    root = Path(brand_dir)
    colors = json.loads((root / "colors.json").read_text(encoding="utf-8"))
    placements = {}
    if (root / "placements.json").is_file():
        placements = json.loads((root / "placements.json").read_text(encoding="utf-8"))
    icons = {"icons": []}
    if (root / "icon-catalog.json").is_file():
        icons = json.loads((root / "icon-catalog.json").read_text(encoding="utf-8"))
    aliases = {}
    if (root / "icon-aliases.json").is_file():
        aliases = (
            json.loads((root / "icon-aliases.json").read_text(encoding="utf-8")).get(
                "aliases"
            )
            or {}
        )
    layout_patterns = {}
    if (root / "layout-patterns.json").is_file():
        layout_patterns = (
            json.loads((root / "layout-patterns.json").read_text(encoding="utf-8")).get(
                "patterns"
            )
            or {}
        )
    theme = colors.get("theme") or {}
    by_space: dict[str, dict] = {}
    by_ascii: dict[str, dict] = {}
    for item in icons.get("icons") or []:
        label = str(item.get("label") or "")
        by_space[_space_key(label)] = item
        ak = _ascii_key(label)
        if len(ak) >= 3:
            by_ascii[ak] = item
    alias_map: dict[str, str] = {}
    for src, dst in aliases.items():
        alias_map[_space_key(src)] = str(dst)
        ak = _ascii_key(src)
        if len(ak) >= 2:
            alias_map[ak] = str(dst)
    return {
        "dir": root,
        "colors": colors,
        "theme": theme,
        "placements": placements,
        "icons": icons.get("icons") or [],
        "by_space": by_space,
        "by_ascii": by_ascii,
        "alias_map": alias_map,
        "layout_patterns": layout_patterns,
    }


def brand_rgb(pack: dict, role_or_hex: str, default: str = "#000000") -> RGBColor:
    theme = pack.get("theme") or {}
    roles = (pack.get("colors") or {}).get("roles") or {}
    key = roles.get(role_or_hex, role_or_hex)
    hex_v = theme.get(key) or theme.get(role_or_hex) or default
    if isinstance(hex_v, str) and hex_v.startswith("#"):
        return _hex_to_rgb(hex_v)
    return _hex_to_rgb(default)


def _icon_path(pack: dict, item: dict) -> Path | None:
    root = pack["dir"]
    rel = item.get("file")
    if not rel:
        return None
    path = root / rel
    if path.is_file():
        return path
    # Prefer PNG when catalog still points at SVG.
    if rel.endswith(".svg"):
        png = root / (rel[:-4] + ".png")
        if png.is_file():
            return png
    return None


def lookup_icon(pack: dict, label: str) -> Path | None:
    """Resolve a short label to a brand icon file (PNG preferred)."""
    raw = _display_label(label)
    if not raw or len(raw) > _LABEL_MAX_CHARS:
        return None
    space = _space_key(raw)
    ascii_k = _ascii_key(raw)
    if not space and len(ascii_k) < 3:
        return None

    by_space = pack.get("by_space") or {}
    by_ascii = pack.get("by_ascii") or {}
    alias_map = pack.get("alias_map") or {}

    # Alias → canonical catalog label
    canonical = alias_map.get(space) or alias_map.get(ascii_k)
    if canonical:
        item = by_space.get(_space_key(canonical)) or by_ascii.get(_ascii_key(canonical))
        if item:
            return _icon_path(pack, item)

    item = by_space.get(space)
    if item:
        return _icon_path(pack, item)
    if len(ascii_k) >= 3:
        item = by_ascii.get(ascii_k)
        if item:
            return _icon_path(pack, item)
        # Prefix/equality only — never empty-string or long-prose substring.
        for key, cand in by_ascii.items():
            if len(key) < 4:
                continue
            if ascii_k == key or (
                min(len(ascii_k), len(key)) >= 5
                and (ascii_k.startswith(key) or key.startswith(ascii_k))
            ):
                return _icon_path(pack, cand)
    return None


def _is_icon_sized(width: int, height: int) -> bool:
    if width < _ICON_MIN or height < _ICON_MIN:
        return False
    if width > _ICON_MAX or height > _ICON_MAX:
        return False
    ratio = width / max(height, 1)
    return 0.45 <= ratio <= 2.2


def _slide_has_picture_near(slide, left: int, top: int, *, tol: int = 200_000) -> bool:
    for sh in slide.shapes:
        if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if abs(int(sh.left or 0) - left) <= tol and abs(int(sh.top or 0) - top) <= tol:
            return True
    return False


def ensure_brand_pictures(slide, pack: dict, kind: str) -> int:
    """Stamp brand logos onto the slide when missing (layout-independent)."""
    placements = (pack.get("placements") or {}).get(kind) or {}
    if not placements:
        return 0
    root = pack["dir"]
    added = 0
    for _name, spec in placements.items():
        rel = spec.get("file")
        if not rel:
            continue
        path = root / rel
        if not path.is_file():
            continue
        left = int(spec["left"])
        top = int(spec["top"])
        width = int(spec["width"])
        height = int(spec["height"])
        if _slide_has_picture_near(slide, left, top):
            continue
        slide.shapes.add_picture(str(path), Emu(left), Emu(top), Emu(width), Emu(height))
        added += 1
    return added


def replace_matched_icons(slide, pack: dict) -> list[str]:
    """Replace icon-sized pictures whose nearby short label matches the catalog."""
    labels: list[tuple[int, int, int, int, str]] = []
    pictures = []
    for sh in slide.shapes:
        if sh.is_placeholder:
            continue
        L, T = int(sh.left or 0), int(sh.top or 0)
        W, H = int(sh.width or 0), int(sh.height or 0)
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and _is_icon_sized(W, H):
            pictures.append(sh)
            continue
        if not getattr(sh, "has_text_frame", False):
            continue
        text = _display_label(sh.text or "")
        if not text or len(text) > _LABEL_MAX_CHARS:
            continue
        first = text.split(" ")[0] if " " in text and len(text) > 40 else text
        # Keep full short labels; trim only obviously long lines.
        label = text if len(text) <= _LABEL_MAX_CHARS else first
        if lookup_icon(pack, label):
            labels.append((L, T, W, H, label))

    replaced: list[str] = []
    for pic in pictures:
        L, T = int(pic.left or 0), int(pic.top or 0)
        W, H = int(pic.width or 0), int(pic.height or 0)
        best_label = None
        best_d = 10**18
        for tl, tt, tw, th, label in labels:
            # Prefer label under / beside the icon.
            cx, cy = L + W // 2, T + H // 2
            lx, ly = tl + tw // 2, tt + th // 2
            if tt + th < T - 200_000:
                continue
            d = abs(lx - cx) + abs(ly - (T + H))
            if d <= _NEAR_EMU and d < best_d:
                best_d = d
                best_label = label
        if not best_label:
            # Fallback: shape name tokens (rare).
            name = (pic.name or "").replace("그래픽", "").strip()
            if name:
                best_label = name if lookup_icon(pack, name) else None
        if not best_label:
            continue
        icon = lookup_icon(pack, best_label)
        if not icon:
            continue
        try:
            pic._element.getparent().remove(pic._element)
            slide.shapes.add_picture(str(icon), Emu(L), Emu(T), Emu(W), Emu(H))
            replaced.append(best_label)
        except Exception:
            continue
    return replaced


def _rgb_hex_of_fill(shape) -> str | None:
    try:
        if shape.fill.type != MSO_FILL.SOLID:
            return None
        return str(shape.fill.fore_color.rgb).upper()
    except Exception:
        return None


def normalize_brand_accents(slide, pack: dict) -> int:
    """Remap common non-academy cyan accents to theme accent1."""
    accent = brand_rgb(pack, "accent", "#006DFF")
    changed = 0

    def walk(shapes) -> None:
        nonlocal changed
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes)
                continue
            fill_hex = _rgb_hex_of_fill(shape)
            if fill_hex in _NON_BRAND_ACCENTS:
                try:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = accent
                    changed += 1
                except Exception:
                    pass
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        rh = str(run.font.color.rgb).upper()
                    except Exception:
                        continue
                    if rh in _NON_BRAND_ACCENTS:
                        run.font.color.rgb = accent
                        changed += 1

    walk(slide.shapes)
    return changed
