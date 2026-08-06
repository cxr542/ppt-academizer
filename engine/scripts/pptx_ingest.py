#!/usr/bin/env python3
"""Extract text and assets from source .pptx (including Google Slides image slides)."""

from __future__ import annotations

from pathlib import Path

from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE

IMAGE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"


def iter_shapes(shapes, depth: int = 0):
    """Yield all shapes, recursing into groups. (Raw relative coordinates)."""
    for sh in shapes:
        yield sh, depth
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(sh.shapes, depth + 1)


def iter_shapes_abs(shapes, depth: int = 0, parent_tx: tuple[float, float, float, float, float, float] | None = None):
    """Yield (shape, depth, abs_l, abs_t, abs_w, abs_h), resolving group scale/offsets."""
    if parent_tx is None:
        parent_tx = (0.0, 0.0, 1.0, 1.0, 0.0, 0.0)
    g_off_x, g_off_y, scale_x, scale_y, ch_off_x, ch_off_y = parent_tx

    for sh in shapes:
        raw_l = float(getattr(sh, "left", 0) or 0)
        raw_t = float(getattr(sh, "top", 0) or 0)
        raw_w = float(getattr(sh, "width", 0) or 0)
        raw_h = float(getattr(sh, "height", 0) or 0)

        abs_l = g_off_x + (raw_l - ch_off_x) * scale_x
        abs_t = g_off_y + (raw_t - ch_off_y) * scale_y
        abs_w = raw_w * scale_x
        abs_h = raw_h * scale_y

        yield sh, depth, abs_l, abs_t, abs_w, abs_h

        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            xfrm = sh._element.xfrm
            new_g_off_x, new_g_off_y = abs_l, abs_t
            new_scale_x, new_scale_y = scale_x, scale_y
            new_ch_off_x, new_ch_off_y = 0.0, 0.0
            if xfrm is not None:
                if xfrm.chExt is not None and xfrm.ext is not None:
                    if getattr(xfrm.chExt, "cx", 0) and getattr(xfrm.ext, "cx", 0):
                        new_scale_x = scale_x * (xfrm.ext.cx / xfrm.chExt.cx)
                    if getattr(xfrm.chExt, "cy", 0) and getattr(xfrm.ext, "cy", 0):
                        new_scale_y = scale_y * (xfrm.ext.cy / xfrm.chExt.cy)
                if xfrm.chOff is not None:
                    new_ch_off_x = float(getattr(xfrm.chOff, "x", 0) or 0)
                    new_ch_off_y = float(getattr(xfrm.chOff, "y", 0) or 0)
            yield from iter_shapes_abs(
                sh.shapes, depth + 1,
                (new_g_off_x, new_g_off_y, new_scale_x, new_scale_y, new_ch_off_x, new_ch_off_y)
            )


def _extract_font_size(shape) -> int | None:
    try:
        if shape.has_text_frame and shape.text_frame.paragraphs:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font and r.font.size:
                        return int(r.font.size)
    except Exception:
        pass
    return None

def _is_bold(shape) -> bool:
    try:
        if shape.has_text_frame and shape.text_frame.paragraphs:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font and r.font.bold:
                        return True
    except Exception:
        pass
    return False

def slide_text_blocks(slide) -> list[dict]:
    """All text-bearing shapes with absolute coordinates and font hints."""
    blocks: list[dict] = []
    for sh, depth, abs_l, abs_t, abs_w, abs_h in iter_shapes_abs(slide.shapes):
        if not getattr(sh, "has_text_frame", False):
            continue
        t = (sh.text or "").strip()
        if not t:
            continue
        blocks.append(
            {
                "top": abs_t / 914400,
                "left": abs_l / 914400,
                "width": abs_w / 914400,
                "height": abs_h / 914400,
                "text": t,
                "len": len(t),
                "depth": depth,
                "font_size": _extract_font_size(sh),
                "is_bold": _is_bold(sh),
                "is_placeholder": sh.is_placeholder,
            }
        )
    blocks.sort(key=lambda b: (b["top"], b["left"]))
    return blocks


def slide_notes_text(slide) -> str:
    """Speaker notes plain text (if any)."""
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        return (tf.text or "").strip()
    except Exception:
        return ""


def is_image_only_slide(slide) -> bool:
    """Google Slides .pptx export: full-slide background picture, no shapes."""
    if slide_text_blocks(slide):
        return False
    if len(slide.shapes) > 0:
        # Allow only empty group shell
        has_real = False
        for sh, _ in iter_shapes(slide.shapes):
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue
            has_real = True
            break
        if has_real:
            return False
    try:
        return slide.background.fill.type == MSO_FILL.PICTURE
    except Exception:
        return False


def slide_background_image_part(slide):
    """Return ImagePart for slide background, or None."""
    for rel in slide.part.rels.values():
        if rel.reltype == IMAGE_REL and rel.target_ref.startswith("../media/"):
            # Prefer background blip (first image rel is background for Google export)
            return rel.target_part
    return None


def export_slide_background_image(slide, dest: Path) -> Path | None:
    """Write slide background image to ``dest``; return path or None."""
    part = slide_background_image_part(slide)
    if part is None or not getattr(part, "blob", None):
        return None
    dest = Path(dest)
    dest.parent.mkdir(parents=True, exist_ok=True)
    ext = ".png"
    ref = ""
    for rel in slide.part.rels.values():
        if rel.reltype == IMAGE_REL:
            ref = rel.target_ref
            break
    if ref.lower().endswith(".jpg") or ref.lower().endswith(".jpeg"):
        ext = ".jpg"
    if dest.suffix.lower() not in (".png", ".jpg", ".jpeg"):
        dest = dest.with_suffix(ext)
    dest.write_bytes(part.blob)
    return dest
