"""Snap McKinsey-style equal card rows onto academy 2/3/4-column layout patterns."""
from __future__ import annotations

import json
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE

import scripts.shape_migrate_lib as bac

_MIN_CARD_W = 2_000_000
_MIN_CARD_H = 900_000
_TOP_TOL = 250_000
_HEIGHT_TOL = 500_000
_WIDTH_TOL_RATIO = 0.22


def _hex_to_rgb(value: str) -> RGBColor:
    h = value.strip().lstrip("#").upper()
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def load_layout_patterns(brand_dir: Path | str | None) -> dict:
    if not brand_dir:
        return {}
    path = Path(brand_dir) / "layout-patterns.json"
    if not path.is_file():
        return {}
    data = json.loads(path.read_text(encoding="utf-8"))
    return data.get("patterns") or {}


def _solid_rgb(shape) -> str | None:
    try:
        if shape.fill.type != MSO_FILL.SOLID:
            return None
        return str(shape.fill.fore_color.rgb).upper()
    except Exception:
        return None


def _is_light_fill(rgb: str | None) -> bool:
    """Soft white/gray/pale-blue card fills (not navy / brand cyan)."""
    if not rgb or len(rgb) < 6:
        return False
    r, g, b = int(rgb[0:2], 16), int(rgb[2:4], 16), int(rgb[4:6], 16)
    return min(r, g, b) >= 0xC0 or ((r + g + b) >= 0x2D0 and min(r, g, b) >= 0xA8)


def _is_full_bleed_panel(shape, slide_width: int) -> bool:
    return int(shape.width or 0) >= int(slide_width * 0.75)


def _is_card_panel(shape, slide_width: int) -> bool:
    if getattr(shape, "is_placeholder", False):
        return False
    if shape.shape_type not in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX):
        return False
    top = int(shape.top or 0)
    w, h = int(shape.width or 0), int(shape.height or 0)
    if top < bac.HEADER_BOTTOM_EMU:
        return False
    if w < _MIN_CARD_W or h < _MIN_CARD_H:
        return False
    if _is_full_bleed_panel(shape, slide_width):
        return False
    # Thin accent rails are not cards.
    if w < 400_000:
        return False
    if _solid_rgb(shape) is None:
        return False
    return True


def _row_is_soft_card_pattern(row: list) -> bool:
    """Require mostly light fills so designed dark chrome is not snapped."""
    n = len(row)
    if n not in (2, 3, 4):
        return False
    light = sum(1 for sh in row if _is_light_fill(_solid_rgb(sh)))
    # 2-col: both soft. 3/4-col: allow one dark accent card in the row.
    need = n if n == 2 else n - 1
    return light >= need


def detect_card_column_count(slide, slide_width: int) -> int:
    """Return 2/3/4 when a single equal-width card row is present, else 0."""
    panels = [sh for sh in slide.shapes if _is_card_panel(sh, slide_width)]
    if len(panels) < 2:
        return 0
    panels.sort(key=lambda sh: (int(sh.top or 0), int(sh.left or 0)))
    # Take the densest top-aligned soft card row (upper row wins ties).
    best_row: list = []
    for seed in panels:
        row = [
            sh
            for sh in panels
            if abs(int(sh.top or 0) - int(seed.top or 0)) <= _TOP_TOL
            and abs(int(sh.height or 0) - int(seed.height or 0)) <= _HEIGHT_TOL
        ]
        if len(row) not in (2, 3, 4) or not _row_is_soft_card_pattern(row):
            continue
        if len(row) > len(best_row):
            best_row = row
    if len(best_row) not in (2, 3, 4):
        return 0
    best_row.sort(key=lambda sh: int(sh.left or 0))
    widths = [int(sh.width or 0) for sh in best_row]
    avg = sum(widths) / len(widths)
    if any(abs(w - avg) / max(avg, 1) > _WIDTH_TOL_RATIO for w in widths):
        return 0
    # Stacked card rows need a dedicated pattern; skip single-row snap.
    row_top = int(best_row[0].top or 0)
    row_elems = {sh._element for sh in best_row}
    for sh in panels:
        if sh._element in row_elems:
            continue
        if int(sh.top or 0) > row_top + _TOP_TOL:
            if abs(int(sh.height or 0) - int(best_row[0].height or 0)) <= _HEIGHT_TOL:
                return 0
    return len(best_row)


def _card_row_panels(slide, slide_width: int, n: int) -> list:
    panels = [sh for sh in slide.shapes if _is_card_panel(sh, slide_width)]
    panels.sort(key=lambda sh: (int(sh.top or 0), int(sh.left or 0)))
    for seed in panels:
        row = [
            sh
            for sh in panels
            if abs(int(sh.top or 0) - int(seed.top or 0)) <= _TOP_TOL
            and abs(int(sh.height or 0) - int(seed.height or 0)) <= _HEIGHT_TOL
        ]
        if len(row) == n and _row_is_soft_card_pattern(row):
            row.sort(key=lambda sh: int(sh.left or 0))
            return row
    return []


def _map_rect(
    left: int,
    top: int,
    width: int,
    height: int,
    old: tuple[int, int, int, int],
    new: tuple[int, int, int, int],
) -> tuple[int, int, int, int]:
    ox, oy, ow, oh = old
    nx, ny, nw, nh = new
    sx = nw / max(ow, 1)
    sy = nh / max(oh, 1)
    return (
        int(nx + (left - ox) * sx),
        int(ny + (top - oy) * sy),
        max(int(width * sx), bac.MIN_SHAPE_W),
        max(int(height * sy), bac.MIN_SHAPE_H),
    )


def apply_card_column_pattern(
    slide,
    patterns: dict,
    slide_width: int,
    *,
    column_count: int | None = None,
) -> str | None:
    """Snap equal card row + in-card content onto academy pattern. Returns pattern name."""
    n = column_count or detect_card_column_count(slide, slide_width)
    if n not in (2, 3, 4):
        return None
    pattern = patterns.get(str(n)) or patterns.get(n)
    if not pattern:
        return None
    slots = pattern.get("cards") or []
    if len(slots) != n:
        return None
    panels = _card_row_panels(slide, slide_width, n)
    if len(panels) != n:
        return None

    old_rects = [
        (int(sh.left or 0), int(sh.top or 0), int(sh.width or 0), int(sh.height or 0))
        for sh in panels
    ]
    # python-pptx may wrap the same shape in new proxies — key on XML element.
    panel_elems = {sh._element for sh in panels}

    # Hide wide soft backgrounds that fight the academy card chrome.
    for sh in list(slide.shapes):
        if sh._element in panel_elems or getattr(sh, "is_placeholder", False):
            continue
        if not _is_full_bleed_panel(sh, slide_width):
            continue
        if int(sh.top or 0) < bac.HEADER_BOTTOM_EMU:
            continue
        if _solid_rgb(sh) is None:
            continue
        try:
            sh._element.getparent().remove(sh._element)
        except Exception:
            pass

    for sh, slot in zip(panels, slots):
        sh.left = int(slot["left"])
        sh.top = int(slot["top"])
        sh.width = int(slot["width"])
        sh.height = int(slot["height"])
        try:
            sh.fill.solid()
            sh.fill.fore_color.rgb = _hex_to_rgb(slot.get("fill") or "#E6F1FF")
        except Exception:
            pass

    new_rects = [
        (int(s["left"]), int(s["top"]), int(s["width"]), int(s["height"])) for s in slots
    ]

    max_card_w = max(ow for _, _, ow, _ in old_rects)
    for sh in slide.shapes:
        if sh._element in panel_elems or getattr(sh, "is_placeholder", False):
            continue
        L, T = int(sh.left or 0), int(sh.top or 0)
        W, H = int(sh.width or 0), int(sh.height or 0)
        if T < bac.HEADER_BOTTOM_EMU:
            continue
        # Left full-height accent rail — keep, don't force into a card.
        if W < 250_000 and H > 2_000_000:
            continue
        # Full-width hero / titles spanning the row must stay above the cards.
        if W > max_card_w * 1.25:
            continue
        cx = L + W // 2
        cy = T + H // 2
        col = None
        for i, (ox, oy, ow, oh) in enumerate(old_rects):
            # Only remap shapes that lived inside a card (avoid pulling hero copy above).
            if ox - 80_000 <= cx <= ox + ow + 80_000 and oy - 80_000 <= cy <= oy + oh + 80_000:
                col = i
                break
        if col is None:
            continue
        nl, nt, nw, nh = _map_rect(L, T, W, H, old_rects[col], new_rects[col])
        # Keep short labels from stretching to nearly the full card height.
        if (
            getattr(sh, "has_text_frame", False)
            and (sh.text or "").strip()
            and H < old_rects[col][3] * 0.35
            and nh > int(new_rects[col][3] * 0.45)
        ):
            nh = max(
                int(H * (new_rects[col][3] / max(old_rects[col][3], 1))),
                bac.MIN_SHAPE_H,
            )
        try:
            sh.left, sh.top, sh.width, sh.height = nl, nt, nw, nh
        except Exception:
            continue

    # Tall academy slots can rise under / through full-width hero copy — park above.
    pattern_top = min(int(s["top"]) for s in slots)
    pattern_h = max(int(s["height"]) for s in slots)
    hero_floor = bac.HEADER_BOTTOM_EMU + 40_000
    band_bottom = pattern_top + int(pattern_h * 0.4)
    for sh in slide.shapes:
        if sh._element in panel_elems or getattr(sh, "is_placeholder", False):
            continue
        try:
            T, W, H = int(sh.top), int(sh.width), int(sh.height)
        except Exception:
            continue
        if W <= max_card_w * 1.25:
            continue
        # Intersects the upper band of the snapped card row.
        if T + H <= pattern_top or T >= band_bottom:
            continue
        avail = pattern_top - hero_floor - 80_000
        if avail < bac.MIN_SHAPE_H:
            continue
        try:
            if H > avail:
                sh.height = avail
            sh.top = hero_floor
        except Exception:
            continue

    return str(pattern.get("name") or f"레이아웃_{n}단")
