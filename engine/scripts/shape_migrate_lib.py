#!/usr/bin/env python3
"""Hybrid shape copy for academy deck migration (pictures, lines, outlines, connectors)."""
from __future__ import annotations

import io
import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

ACADEMY = Path(
    "/Users/yhkim/Desktop/[파트너사] 02. 클라우드 구현기술(PaaS)_아카데미디자인.pptx"
)
BACKUP = Path(
    "/Users/yhkim/Desktop/[파트너사] 02. 클라우드 구현기술(PaaS)_수정필요_backup.pptx"
)
OUTPUT = Path(
    "/Users/yhkim/Desktop/[파트너사] 02. 클라우드 구현기술(PaaS)_아카데미적용_20260518_최종_v16.pptx"
)

# Flowchart accent boxes are small; terminal/console bars are full-width.
FLOWCHART_BOX_MIN_W = 800_000
FLOWCHART_BOX_MAX_W = 2_000_000

# Google Slides export: terminal groups share child row coords (2150309); add group offset when c_top < g_top.
TERMINAL_TEMPLATE_GROUP_TOP = 1_965_680
TERMINAL_TEMPLATE_CHILD_TOP = 2_150_309
TERMINAL_TEMPLATE_CHILD_LEFT = 460_375
TERMINAL_ROW_DY = TERMINAL_TEMPLATE_CHILD_TOP - TERMINAL_TEMPLATE_GROUP_TOP
GOVERNING_PLACEHOLDER_HINT = "거버닝 메시지"

TERMINAL_FILL_RGB = RGBColor(0x2B, 0x2B, 0x2B)

OFFICE_EXT_URIS = ("FF2B5EF4", "DCECCB84", "BB962C8B", "C183EC19")

HEADER_BOTTOM_EMU = 1_150_000
SKIP_SHAPE_NAMES = {"Google Shape;520;p15"}
NUM_LEFT, NUM_TOP = 515_938, 360_540
NUM_WIDTH = 400_000  # two-digit step numbers (10, 11, …)
# Academy master: title x=781263 (= num left + narrow num width + 143497 gap).
# Do not push title right when num column widens for two digits.
ACADEMY_TITLE_LEFT = 781_263
TITLE_LEFT = ACADEMY_TITLE_LEFT
TITLE_TOP = 370_849
TITLE_WIDTH = 9_948_172  # keeps right edge; was 9660000 when title started at 1059435
GOV_LEFT, GOV_TOP, GOV_WIDTH = 514_984, 757_189, 5_760_404
GOV_MAX_LEN = 42
BODY_TOP_MIN = 1_050_000
RIGHT_COL_MIN = 7_000_000
CONTENT_LAYOUTS = {"내지_거버닝 O", "1_내지_거버닝 X"}
SLIDE_MARGIN_X = 380_000
SLIDE_MARGIN_BOTTOM = 520_000
MIN_SHAPE_W = 180_000
MIN_SHAPE_H = 80_000


def norm(text: str) -> str:
    return " ".join(text.split())


def text_similar(a: str, b: str) -> bool:
    if not a or not b:
        return False
    a, b = norm(a), norm(b)
    if a == b:
        return True
    n = min(28, len(a), len(b))
    return a[:n] == b[:n] or a in b or b in a


def set_placeholder(slide, idx: int, text: str) -> None:
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            shape.text = text
            return


def section_num_from_slide(src_slide) -> str:
    for shape in src_slide.shapes:
        if "Google Shape;520" in shape.name:
            t = shape.text.strip()
            if t.isdigit():
                return t
    return ""


def extract_header(src_slide):
    title = governing = None
    for shape in src_slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text.strip()
        if not text:
            continue
        name = shape.name
        if name == "제목 1" or (
            shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE
        ):
            title = text.split("\n")[0]
        elif name == "텍스트 개체 틀 3":
            governing = text.split("\n")[0]
        elif name.startswith("제목") and title is None and len(text) < 80:
            title = text.split("\n")[0]
    return title, governing


def is_body_shape(shape, title: str | None, governing: str | None) -> bool:
    if shape.name in SKIP_SHAPE_NAMES:
        return False
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        return False

    st = shape.shape_type
    top = int(shape.top or 0)
    w, h = int(shape.width or 0), int(shape.height or 0)

    # Pictures / lines / groups / tables: do not drop large diagrams in the header band.
    if st in (
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.LINE,
        MSO_SHAPE_TYPE.GROUP,
        MSO_SHAPE_TYPE.TABLE,
    ):
        if st == MSO_SHAPE_TYPE.PICTURE and top and top < HEADER_BOTTOM_EMU:
            return w >= 1_200_000 or h >= 900_000
        return True

    # Google-export thin vertical rules beside the title strip.
    if st == MSO_SHAPE_TYPE.AUTO_SHAPE and w < 250_000 and h > 200_000:
        return True

    if top and top < HEADER_BOTTOM_EMU:
        return False

    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:
                if text.isdigit() and len(text) <= 3:
                    return False
                first = text.split("\n")[0]
                if title and first == title:
                    return False
                if governing and first == governing:
                    return False
        return True

    if shape.has_text_frame:
        text = shape.text.strip()
        if not text:
            return False
        if text.isdigit() and len(text) <= 3:
            return False
        first = text.split("\n")[0]
        if title and first == title:
            return False
        if governing and first == governing:
            return False
    return True


def _sanitize_shape_xml(element) -> None:
    remove = []
    for el in element.iter():
        tag = el.tag.split("}")[-1] if "}" in el.tag else el.tag
        if tag in ("creationId", "custDataLst"):
            remove.append(el)
            continue
        if tag == "pPr":
            for attr in ("eaLnBrk", "hangingPunct", "latinLnBrk"):
                el.attrib.pop(attr, None)
            continue
        if tag == "cNvPr":
            el.attrib.pop("descr", None)
            el.attrib.pop("title", None)
            continue
        if tag == "modId":
            el.attrib.pop("val", None)
            continue
        if tag == "extLst" and len(el) == 0:
            remove.append(el)
            continue
        if tag == "ext":
            uri = el.get("uri") or ""
            if any(marker in uri for marker in OFFICE_EXT_URIS):
                remove.append(el)
    for el in remove:
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)


def _offset_shape_xml(element, dx: int, dy: int) -> None:
    if not dx and not dy:
        return
    tag_off = "{http://schemas.openxmlformats.org/drawingml/2006/main}off"
    for off in element.iter(tag_off):
        off.set("x", str(int(off.get("x", 0)) + dx))
        off.set("y", str(int(off.get("y", 0)) + dy))


def _cnv_id(shape) -> str:
    for el in shape._element.iter():
        if el.tag.split("}")[-1] == "cNvPr" and el.get("id") is not None:
            return str(el.get("id"))
    return ""


def _is_cxn_sp(shape) -> bool:
    return shape._element.tag.split("}")[-1] == "cxnSp"


def _has_p_style(shape) -> bool:
    return any(el.tag.split("}")[-1] == "style" for el in shape._element.iter())


def _is_flowchart_accent_box(shape) -> bool:
    """Small p:style rectangles (not full-width terminal bars)."""
    if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
        return False
    if not _has_p_style(shape):
        return False
    if shape.has_text_frame and shape.text.strip():
        return False
    w = int(shape.width or 0)
    return FLOWCHART_BOX_MIN_W <= w < FLOWCHART_BOX_MAX_W


def _is_spurious_overlay_box(shape) -> bool:
    """Tiny p:style rects over screenshots (e.g. 477k EMU on Killercoda slide)."""
    if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
        return False
    if not _has_p_style(shape):
        return False
    if shape.has_text_frame and shape.text.strip():
        return False
    w, h = int(shape.width or 0), int(shape.height or 0)
    return w < FLOWCHART_BOX_MIN_W and h < 1_200_000


def _group_child_position(group, child, off_l: int, off_t: int) -> tuple[int, int]:
    """Resolve Google-export group child coordinates to slide-absolute EMU."""
    g_l, g_t = int(group.left or 0), int(group.top or 0)
    c_l, c_t = int(child.left or 0), int(child.top or 0)
    if c_t >= g_t:
        return off_l + c_l, off_t + c_t
    # Reused template row inside lower groups: offset from group anchor.
    rel_l = c_l - TERMINAL_TEMPLATE_CHILD_LEFT
    rel_t = c_t - TERMINAL_TEMPLATE_CHILD_TOP
    return (
        off_l + g_l + TERMINAL_TEMPLATE_CHILD_LEFT + rel_l,
        off_t + g_t + TERMINAL_ROW_DY + rel_t,
    )


def clone_connector_xml(shape, slide, off_l: int = 0, off_t: int = 0) -> None:
    """Keep p:style + stCxn/endCxn — bent arrows need them in PowerPoint."""
    newel = deepcopy(shape._element)
    _sanitize_shape_xml(newel)
    _offset_shape_xml(newel, off_l, off_t)
    slide.shapes._spTree.insert_element_before(newel, "p:extLst")


def _is_graphic_frame(shape) -> bool:
    tag = shape._element.tag
    return tag.endswith("}graphicFrame") or tag == "graphicFrame"


def clone_shape_xml(shape, slide, off_l: int = 0, off_t: int = 0) -> None:
  newel = deepcopy(shape._element)
  _sanitize_shape_xml(newel)
  _offset_shape_xml(newel, off_l, off_t)
  slide.shapes._spTree.insert_element_before(newel, "p:extLst")


def _copy_fill(dst_shape, src_shape) -> None:
    try:
        src, dst = src_shape.fill, dst_shape.fill
        if src.type == MSO_FILL.SOLID:
            dst.solid()
            sc = src.fore_color
            if sc.type == MSO_COLOR_TYPE.SCHEME:
                dst.fore_color.theme_color = sc.theme_color
                if sc.brightness:
                    dst.fore_color.brightness = sc.brightness
            elif sc.rgb:
                dst.fore_color.rgb = sc.rgb
        elif src.type == MSO_FILL.BACKGROUND:
            dst.background()
    except Exception:
        pass


def _copy_line(dst_shape, src_shape) -> None:
    try:
        sl, dl = src_shape.line, dst_shape.line
        if sl.fill.type == MSO_FILL.SOLID:
            dl.fill.solid()
            sc = sl.fill.fore_color
            if sc.rgb:
                dl.fill.fore_color.rgb = sc.rgb
            elif sc.type == MSO_COLOR_TYPE.SCHEME:
                dl.fill.fore_color.theme_color = sc.theme_color
                if sc.brightness:
                    dl.fill.fore_color.brightness = sc.brightness
        elif sl.fill.type == MSO_FILL.BACKGROUND:
            dl.fill.background()
        if sl.width is not None:
            dl.width = sl.width
        if sl.dash_style is not None:
            dl.dash_style = sl.dash_style
    except Exception:
        pass


def add_picture_safe(slide, shape, left=None, top=None) -> None:
    try:
        blob = shape.image.blob
    except Exception:
        return
    slide.shapes.add_picture(
        io.BytesIO(blob),
        left if left is not None else shape.left,
        top if top is not None else shape.top,
        width=shape.width,
        height=shape.height,
    )


def _copy_font_from_source(dst_shape, src_shape) -> None:
    if not (dst_shape.has_text_frame and src_shape.has_text_frame):
        return
    try:
        src_paras = src_shape.text_frame.paragraphs
        dst_paras = dst_shape.text_frame.paragraphs
        for si, sp in enumerate(src_paras):
            if si >= len(dst_paras):
                break
            if not sp.runs or not dst_paras[si].runs:
                continue
            sf = sp.runs[0].font
            for run in dst_paras[si].runs:
                if sf.size:
                    run.font.size = sf.size
                if sf.name:
                    run.font.name = sf.name
                if sf.bold is not None:
                    run.font.bold = sf.bold
                try:
                    if sf.color.type == MSO_COLOR_TYPE.RGB and sf.color.rgb:
                        run.font.color.rgb = sf.color.rgb
                    elif sf.color.type == MSO_COLOR_TYPE.SCHEME:
                        run.font.color.theme_color = sf.color.theme_color
                except Exception:
                    pass
    except Exception:
        pass


def copy_textbox_native(slide, shape, left=None, top=None) -> None:
    box = slide.shapes.add_textbox(
        left if left is not None else shape.left,
        top if top is not None else shape.top,
        shape.width,
        shape.height,
    )
    box.text_frame.word_wrap = True
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    if shape.has_text_frame:
        box.text_frame.text = shape.text
        _copy_font_from_source(box, shape)


def _apply_terminal_fill(dst_shape, src_shape) -> None:
    """Google-export terminal boxes use tx1 (black on academy); use dark gray."""
    if src_shape.has_text_frame and src_shape.text.strip():
        return
    try:
        if src_shape.fill.type != MSO_FILL.SOLID:
            return
        fc = src_shape.fill.fore_color
        if fc.type == MSO_COLOR_TYPE.SCHEME:
            dst_shape.fill.solid()
            dst_shape.fill.fore_color.rgb = TERMINAL_FILL_RGB
    except Exception:
        pass


def copy_accent_box_native(slide, shape, left, top) -> None:
    """Flowchart boxes that rely on p:style lnRef/fillRef (accent1 outline)."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, shape.width, shape.height)
    box.fill.solid()
    box.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    box.fill.fore_color.brightness = 0.35
    box.line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    box.line.width = Pt(1.75)


def copy_autoshape_native(slide, shape, left, top) -> None:
    try:
        auto_type = shape.auto_shape_type
    except (ValueError, AttributeError):
        auto_type = MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(auto_type, left, top, shape.width, shape.height)
    _copy_fill(box, shape)
    _apply_terminal_fill(box, shape)
    _copy_line(box, shape)
    if shape.has_text_frame:
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        if shape.text.strip():
            tf.text = shape.text
            _copy_font_from_source(box, shape)


def _strip_connector_style_element(conn) -> None:
    """p:style on programmatic cxnSp triggers Mac PowerPoint content removal."""
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    for el in list(conn._element.iter()):
        if el.tag == f"{{{p_ns}}}style":
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)


def copy_line_native(slide, shape, left=None, top=None) -> None:
    """Straight connector from shape bounds (safe for Google-export cxnSp)."""
    x1 = left if left is not None else shape.left
    y1 = top if top is not None else shape.top
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        x1,
        y1,
        x1 + shape.width,
        y1 + shape.height,
    )
    _copy_line(conn, shape)
    try:
        conn.line.fill.solid()
        conn.line.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    except Exception:
        pass
    _strip_connector_style_element(conn)


def copy_shape_hybrid(
    dst_slide,
    shape,
    off_l: int = 0,
    off_t: int = 0,
    id_remap: dict[str, str] | None = None,
) -> None:
    """Native text/pics/terminals; small accent boxes native; connectors XML + cxn ids."""
    left, top = int(shape.left or 0) + off_l, int(shape.top or 0) + off_t
    old_id = _cnv_id(shape)
    n_before = len(dst_slide.shapes)
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            c_left, c_top = _group_child_position(shape, child, off_l, off_t)
            copy_shape_hybrid(
                dst_slide,
                child,
                c_left - int(child.left or 0),
                c_top - int(child.top or 0),
                id_remap,
            )
        return
    if _is_spurious_overlay_box(shape):
        return
    # Never XML-clone cxnSp: Google exports often have broken stCxn/endCxn → PP repair dialog.
    if _is_cxn_sp(shape) or st == MSO_SHAPE_TYPE.LINE:
        copy_line_native(dst_slide, shape, left, top)
    elif st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        if _is_flowchart_accent_box(shape):
            copy_accent_box_native(dst_slide, shape, left, top)
        else:
            # Native rebuild keeps PowerPoint from showing "repair" on p:style XML clones.
            copy_autoshape_native(dst_slide, shape, left, top)
    elif st == MSO_SHAPE_TYPE.PICTURE:
        add_picture_safe(dst_slide, shape, left, top)
        if len(dst_slide.shapes) <= n_before:
            clone_shape_xml(shape, dst_slide, off_l, off_t)
    elif getattr(shape, "has_chart", False) and shape.has_chart:
        clone_shape_xml(shape, dst_slide, off_l, off_t)
    elif st == MSO_SHAPE_TYPE.CHART:
        clone_shape_xml(shape, dst_slide, off_l, off_t)
    elif _is_graphic_frame(shape):
        clone_shape_xml(shape, dst_slide, off_l, off_t)
    elif st == MSO_SHAPE_TYPE.TABLE:
        clone_shape_xml(shape, dst_slide, off_l, off_t)
    elif st == MSO_SHAPE_TYPE.TEXT_BOX or shape.has_text_frame:
        copy_textbox_native(dst_slide, shape, left, top)
    else:
        if _has_p_style(shape) or _is_cxn_sp(shape):
            clone_shape_xml(shape, dst_slide, off_l, off_t)
        else:
            copy_autoshape_native(dst_slide, shape, left, top)

    if id_remap is not None and old_id:
        if len(dst_slide.shapes) > n_before:
            id_remap[old_id] = _cnv_id(dst_slide.shapes[n_before])
        elif _is_cxn_sp(shape):
            id_remap[old_id] = old_id


def copy_body_shapes(src_slide, dst_slide, title, governing) -> dict[str, str]:
    id_remap: dict[str, str] = {}
    for shape in src_slide.shapes:
        if not is_body_shape(shape, title, governing):
            continue
        copy_shape_hybrid(dst_slide, shape, id_remap=id_remap)
    return id_remap


def migrate_content_slide(ref_slide, src_slide) -> None:
    title, governing = extract_header(src_slide)

    for shape in list(ref_slide.shapes):
        if not shape.is_placeholder:
            shape._element.getparent().remove(shape._element)

    set_placeholder(ref_slide, 10, title or "")
    set_placeholder(ref_slide, 12, section_num_from_slide(src_slide))
    gov = governing or ""
    if len(gov) <= GOV_MAX_LEN:
        set_placeholder(ref_slide, 13, gov)

    apply_header_geometry(ref_slide)
    id_remap = copy_body_shapes(src_slide, ref_slide, title, governing)
    renumber_map = renumber_shape_ids(ref_slide)
    fix_orphan_connector_refs(ref_slide, id_remap, renumber_map)


def migrate_section_slide(ref_slide, src_slide) -> None:
    for shape in list(ref_slide.shapes):
        if not shape.is_placeholder:
            shape._element.getparent().remove(shape._element)

    id_remap: dict[str, str] = {}
    for shape in src_slide.shapes:
        if shape.is_placeholder:
            continue
        copy_shape_hybrid(ref_slide, shape, id_remap=id_remap)
    renumber_map = renumber_shape_ids(ref_slide)
    fix_orphan_connector_refs(ref_slide, id_remap, renumber_map)


def apply_header_geometry(slide) -> None:
    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        if idx == 12:
            shape.left, shape.top, shape.width = NUM_LEFT, NUM_TOP, NUM_WIDTH
        elif idx == 10:
            shape.left, shape.top, shape.width = TITLE_LEFT, TITLE_TOP, TITLE_WIDTH
        elif idx == 13:
            shape.left, shape.top, shape.width = GOV_LEFT, GOV_TOP, GOV_WIDTH
            if GOVERNING_PLACEHOLDER_HINT in (shape.text or ""):
                shape.text = ""


def polish_slide(slide, *, skip_body_dedup: bool = False) -> None:
    if slide.slide_layout.name not in CONTENT_LAYOUTS:
        return

    apply_header_geometry(slide)

    if skip_body_dedup:
        return

    items = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()
        if not t:
            continue
        items.append(
            {
                "shape": shape,
                "norm": norm(t),
                "left": int(shape.left or 0),
                "top": int(shape.top or 0),
                "is_ph": shape.is_placeholder,
                "ph_idx": shape.placeholder_format.idx if shape.is_placeholder else None,
            }
        )

    for item in items:
        if item["ph_idx"] == 13:
            gov = item["norm"]
            if len(gov) > GOV_MAX_LEN:
                item["shape"].text = ""
                continue
            for other in items:
                if other["is_ph"] or other["shape"] is item["shape"]:
                    continue
                if other["left"] >= RIGHT_COL_MIN and text_similar(gov, other["norm"]):
                    item["shape"].text = ""
                    break

    rights = [
        x
        for x in items
        if not x["is_ph"] and x["left"] >= RIGHT_COL_MIN and x["top"] >= BODY_TOP_MIN
    ]
    if not rights:
        return
    for item in items:
        if item["is_ph"] or item["left"] >= RIGHT_COL_MIN or item["top"] < BODY_TOP_MIN:
            continue
        for r in rights:
            if text_similar(item["norm"], r["norm"]):
                item["shape"]._element.getparent().remove(item["shape"]._element)
                break


def fix_orphan_connector_refs(
    slide, backup_id_remap: dict[str, str], renumber_map: dict[str, str]
) -> None:
    """Rewire stCxn/endCxn when native picture/text replaced a cloned shape id."""
    sp_tree = slide.shapes._spTree
    valid = {
        el.get("id")
        for el in sp_tree.iter()
        if el.tag.split("}")[-1] == "cNvPr" and el.get("id")
    }

    def resolve(ref: str) -> str | None:
        dst = backup_id_remap.get(ref, ref)
        final = renumber_map.get(dst, dst)
        return final if final in valid else None

    for el in sp_tree.iter():
        tag = el.tag.split("}")[-1]
        if tag not in ("stCxn", "endCxn"):
            continue
        ref = el.get("id")
        if not ref or ref in valid:
            continue
        fixed = resolve(ref)
        if fixed:
            el.set("id", fixed)


def renumber_shape_ids(slide) -> dict[str, str]:
    """Assign unique cNvPr ids in tree order; rewire connector references."""
    sp_tree = slide.shapes._spTree
    cnvprs = [
        el
        for el in sp_tree.iter()
        if el.tag.split("}")[-1] == "cNvPr" and "id" in el.attrib
    ]
    id_map: dict[str, str] = {}
    seen: set[str] = set()
    next_id = 2
    for el in cnvprs:
        old = str(el.get("id", ""))
        while str(next_id) in seen:
            next_id += 1
        new = str(next_id)
        if old != new:
            id_map[old] = new
            el.set("id", new)
        else:
            id_map[old] = old
        seen.add(new)
        next_id += 1

    for el in sp_tree.iter():
        tag = el.tag.split("}")[-1]
        if tag in ("stCxn", "endCxn"):
            ref = el.get("id")
            if ref in id_map:
                el.set("id", id_map[ref])
    return id_map


def _scrub_slide_xml(data: bytes) -> bytes:
    from lxml import etree

    root = etree.fromstring(data)
    remove: list = []
    for el in root.iter():
        tag = el.tag.split("}")[-1]
        if tag == "creationId":
            remove.append(el)
        elif tag == "spAutoFit":
            remove.append(el)
        elif tag == "ext":
            uri = el.get("uri") or ""
            if any(marker in uri for marker in OFFICE_EXT_URIS):
                remove.append(el)
    for el in remove:
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)
    return etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True
    )


def _scrub_slide_parts(path: Path) -> None:
    import re
    import zipfile

    tmp = path.with_suffix(".scrub.zip")
    with zipfile.ZipFile(path, "r") as zin:
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if re.match(r"ppt/slides/slide\d+\.xml$", item.filename):
                    data = _scrub_slide_xml(data)
                zout.writestr(item, data)
    tmp.replace(path)


def _disable_shape_autofit(shape) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE


def fit_body_shapes(slide, slide_width: int, slide_height: int) -> None:
    """Text-only width clamp + disable spAutoFit (never squash box heights)."""
    if slide.slide_layout.name not in CONTENT_LAYOUTS:
        return

    max_right = int(slide_width) - SLIDE_MARGIN_X

    for shape in slide.shapes:
        if shape.is_placeholder:
            continue
        top = int(shape.top or 0)
        if top < HEADER_BOTTOM_EMU:
            continue
        if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
            continue
        if not shape.has_text_frame:
            continue

        _disable_shape_autofit(shape)

        left = int(shape.left or 0)
        width = int(shape.width or 0)
        right = left + width
        if right > max_right:
            overflow = right - max_right
            new_w = width - overflow
            if new_w >= MIN_SHAPE_W:
                shape.width = new_w
            else:
                shape.left = max(NUM_LEFT, left - overflow)
                shape.width = max(MIN_SHAPE_W, width - overflow)


def _fix_open_in_slide_view(path: Path) -> None:
    import re
    import zipfile

    tmp = path.with_suffix(".view.zip")
    with zipfile.ZipFile(path, "r") as zin:
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "ppt/viewProps.xml":
                    xml = data.decode("utf-8")
                    xml = re.sub(
                        r'lastView="(?:sldMasterView|slideMasterView)"',
                        'lastView="sldView"',
                        xml,
                    )
                    data = xml.encode("utf-8")
                zout.writestr(item, data)
    tmp.replace(path)


def build_pairs():
    pairs: list[tuple[int, int, str]] = []
    pairs.append((2, 0, "section"))
    for i in range(25):
        pairs.append((3 + i, 1 + i, "content"))
    pairs.append((28, 26, "section"))
    for i in range(20):
        pairs.append((29 + i, 27 + i, "content"))
    return pairs


def main() -> None:
    if not ACADEMY.exists():
        raise SystemExit(f"Missing template: {ACADEMY}")
    if not BACKUP.exists():
        raise SystemExit(f"Missing backup: {BACKUP}")

    if OUTPUT.exists():
        backup = OUTPUT.with_name(OUTPUT.stem + "_손상백업.pptx")
        shutil.copy2(OUTPUT, backup)
        print(f"Backed up previous file to: {backup.name}")

    shutil.copy2(ACADEMY, OUTPUT)
    prs = Presentation(str(OUTPUT))
    src = Presentation(str(BACKUP))

    for ref_idx, src_idx, kind in build_pairs():
        ref_slide = prs.slides[ref_idx]
        src_slide = src.slides[src_idx]
        if kind == "section":
            migrate_section_slide(ref_slide, src_slide)
        else:
            migrate_content_slide(ref_slide, src_slide)

    for slide in prs.slides:
        polish_slide(slide)
        fit_body_shapes(slide, prs.slide_width, prs.slide_height)

    # Round-trip through a temp file so OPC parts are rewritten cleanly.
    tmp = OUTPUT.with_suffix(".tmp.pptx")
    prs.save(str(tmp))
    Presentation(str(tmp)).save(str(OUTPUT))
    tmp.unlink(missing_ok=True)
    _scrub_slide_parts(OUTPUT)
    _fix_open_in_slide_view(OUTPUT)

    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
