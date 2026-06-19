#!/usr/bin/env python3
"""Convert a legacy partner PPT (내지A/간지 layouts) into academy 2026 JSON specs."""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation

from scripts.front_matter import extract_front_matter
from scripts.pptx_ingest import (
    export_slide_background_image,
    is_image_only_slide,
    iter_shapes,
    slide_notes_text,
    slide_text_blocks,
)

TEXTS_LEN_BY_LAYOUT = {
    "2_표지": 1,
    "목차": 1,
    "간지": 2,
    "내지_거버닝 O": 3,
    "1_내지_거버닝 X": 2,
    "내지_참고": 2,
}


def validate_spec(spec: dict, *, slide_index: int | None = None) -> list[str]:
    """Return validation messages for academy-design §5 texts mapping."""
    issues: list[dict] = []
    layout = spec.get("layout", "")
    texts = spec.get("texts")
    if texts is None:
        return ["missing texts"]
    expected = TEXTS_LEN_BY_LAYOUT.get(layout)
    if expected is not None and len(texts) != expected:
        prefix = f"slide {slide_index}: " if slide_index is not None else ""
        issues.append(f"{prefix}layout {layout!r} expects {expected} texts, got {len(texts)}")
    return issues

ROOT = Path(__file__).resolve().parent.parent


def _is_page_number(text: str) -> bool:
    return bool(re.fullmatch(r"\d{1,2}", text.strip()))


def _is_series_banner(text: str) -> bool:
    return "무작정 따라하기" in text or text in ("도커", "쿠버네티스")


# Title band (inches) — academy ph10 header row, not left sidebar labels.
TITLE_TOP_MAX_IN = 1.55
TITLE_LINE_MAX_CHARS = 72
TITLE_SIDEBAR_LEFT_MAX_IN = 0.32
GOVERNING_MAX_CHARS = 48


def _first_line(text: str) -> str:
    return text.replace("\x0b", " ").split("\n")[0].strip()


def _split_block_lines(text: str) -> tuple[str, str]:
    """First line → ph10 title; remainder → body (same text box or long one-liner)."""
    lines = [ln.strip() for ln in text.replace("\x0b", "\n").split("\n") if ln.strip()]
    if not lines:
        return "", ""
    if len(lines) > 1:
        return lines[0][:TITLE_LINE_MAX_CHARS], "\n".join(lines[1:])
    one = lines[0]
    if len(one) <= TITLE_LINE_MAX_CHARS:
        return one, ""
    for sep in (". ", "。 ", ": ", " - "):
        if sep in one[: TITLE_LINE_MAX_CHARS + 20]:
            head, _, tail = one.partition(sep)
            if len(head) <= TITLE_LINE_MAX_CHARS and tail.strip():
                return head.strip(), tail.strip()
    cut = one[:TITLE_LINE_MAX_CHARS]
    if " " in cut:
        cut = cut.rsplit(" ", 1)[0]
    return cut.strip(), one[len(cut) :].strip()


def _is_left_body_column(block: dict) -> bool:
    """Far-left header block with long copy — body column, never ph10 title."""
    left = float(block.get("left", 0) or 0)
    return left < TITLE_SIDEBAR_LEFT_MAX_IN and len(_first_line(block["text"])) > 42


def _is_sidebar_label_block(block: dict) -> bool:
    """Far-left short governing label (e.g. 시대 요약), not slide title."""
    if _is_left_body_column(block):
        return False
    left = float(block.get("left", 0) or 0)
    line = _first_line(block["text"])
    return left < TITLE_SIDEBAR_LEFT_MAX_IN and len(line) < 28


def _pick_title_block(title_band: list[dict]) -> dict | None:
    if not title_band:
        return None
    pool = [b for b in title_band if not _is_left_body_column(b)]
    if not pool:
        pool = title_band
    return min(pool, key=lambda b: (b["top"], len(_first_line(b["text"]))))


def _pick_title_and_body_from_blocks(usable: list[dict]) -> tuple[str, str, str]:
    """Map ingest blocks → ph10 title / ph12 governing / ph13·ph12 body."""
    sorted_u = sorted(usable, key=lambda b: (b["top"], b.get("left", 0.0)))
    title_band = [b for b in sorted_u if b["top"] <= TITLE_TOP_MAX_IN]
    lower_band = [b for b in sorted_u if b["top"] > TITLE_TOP_MAX_IN]

    subtitle = ""
    sidebar = [b for b in title_band if _is_sidebar_label_block(b)]
    if sidebar:
        subtitle = _first_line(sidebar[0]["text"])[:GOVERNING_MAX_CHARS]

    title_text = ""
    body_extra = ""
    title_source = _pick_title_block(title_band)
    if title_source:
        title_text, body_extra = _split_block_lines(title_source["text"])
    elif sorted_u:
        title_text, body_extra = _split_block_lines(sorted_u[0]["text"])

    body_parts: list[str] = []
    if body_extra:
        body_parts.append(body_extra)
    for b in lower_band:
        t = b["text"].strip()
        if t and _first_line(t) != title_text:
            body_parts.append(t)
    for b in title_band:
        if b is title_source or b in sidebar:
            continue
        t = b["text"].strip()
        if t and _first_line(t) != title_text and _first_line(t) != subtitle:
            body_parts.append(t)

    if not body_parts:
        for b in sorted(sorted_u, key=lambda x: -x["len"]):
            if b is title_source:
                continue
            extra_title, extra_body = _split_block_lines(b["text"])
            if extra_body:
                if not title_text:
                    title_text = extra_title
                body_parts.append(extra_body)
                break

    body = "\n\n".join(body_parts).strip()
    if not subtitle:
        for b in sorted_u:
            if b is title_source:
                continue
            t = _first_line(b["text"])
            if (
                t != title_text
                and t not in body
                and GOVERNING_MAX_CHARS >= len(t) >= 4
                and len(t) < 28
                and (_is_sidebar_label_block(b) or b["top"] <= TITLE_TOP_MAX_IN)
            ):
                subtitle = t
                break
    if subtitle.strip() == title_text.strip():
        subtitle = ""
    if _should_promote_speaker_script_heading(title_text, subtitle, body):
        body = f"{title_text}\n\n{body}".strip()
        title_text, subtitle = subtitle, title_text
    return title_text, body, subtitle


def _classify_category(text: str) -> bool:
    return text in {
        "실습 시나리오",
        "소개",
        "테스트",
        "사전 준비",
        "Docker cli",
        "kubectl",
    } or text.startswith("간단 사용법")


CODE_BLOCK_TOKENS = (
    "apiVersion:",
    "kind:",
    "metadata:",
    "spec:",
    "containers:",
    "volumeMounts:",
    "kubectl",
    "ingressClassName:",
)


def _contains_code_block(text: str) -> bool:
    return any(token in text for token in CODE_BLOCK_TOKENS)


def _is_visual_heavy_slide(slide, *, code_block: bool = False) -> bool:
    if code_block:
        return False
    text_shapes = 0
    visual_shapes = 0
    for shape, depth in iter_shapes(slide.shapes):
        if depth > 0:
            visual_shapes += 1
        if shape.is_placeholder:
            continue
        if shape.has_text_frame and (shape.text or "").strip():
            text_shapes += 1
        visual_shapes += 1
    return visual_shapes >= 10 and text_shapes >= 5


def _is_speaker_script_heading(text: str) -> bool:
    compact = " ".join((text or "").split())
    return ("강사용" in compact or "강사" in compact) and (
        "멘트" in compact or "스크립트" in compact or "진행" in compact
    )


def _should_promote_speaker_script_heading(title: str, subtitle: str, body: str) -> bool:
    quoted_lines = sum(1 for line in body.splitlines() if line.strip().startswith(("“", '"')))
    return _is_speaker_script_heading(subtitle) and len(title.strip()) <= 24 and quoted_lines >= 2


def _attach_ingest_metadata(spec: dict, slide, slide_no: int, *, code_block: bool = False) -> dict:
    ing = dict(spec.get("_ingest") or {})
    ing["source_slide_index"] = slide_no - 1
    notes = slide_notes_text(slide)
    if notes:
        ing["speaker_notes"] = notes
    if code_block:
        ing["code_block"] = True
        ing["code_block_shape_preservation"] = "source_text_boxes"
        ing["warning"] = "YAML/code block detected; verify indentation and wrapping."
    if _is_visual_heavy_slide(slide, code_block=code_block):
        ing["visual_preservation"] = "native_shapes"
    spec["_ingest"] = ing
    return spec


def _spec_from_image_slide(
    slide,
    slide_no: int,
    *,
    assets_dir: Path | None,
    deck_title: str,
) -> dict:
    """Google Slides export: content is a full-slide background image."""
    notes = slide_notes_text(slide)
    title = notes if notes and not _is_page_number(notes) else f"슬라이드 {slide_no}"
    if slide_no == 1 and deck_title:
        title = deck_title
    body = " "
    if notes and len(notes) > 60:
        lines = [ln.strip() for ln in notes.splitlines() if ln.strip()]
        if lines:
            title = lines[0][:120]
            body = "\n".join(lines[1:]) if len(lines) > 1 else notes
    elif notes and len(notes) > 20:
        body = notes

    spec: dict = {
        "layout": "1_내지_거버닝 X",
        "texts": [title, body],
        "_ingest": {"kind": "google_image_slide"},
    }

    if assets_dir is not None:
        assets_dir.mkdir(parents=True, exist_ok=True)
        dest = assets_dir / f"slide-{slide_no:02d}.png"
        exported = export_slide_background_image(slide, dest)
        if exported:
            spec["background_image"] = str(exported.resolve())
        else:
            spec["_ingest"]["warning"] = "background_image_missing"
    if (spec["texts"][1] or "").strip() in ("", " "):
        spec["_ingest"]["warning"] = (
            "본문이 슬라이드 배경 이미지에만 있습니다. "
            "발표자 노트를 넣거나 §5 JSON으로 편집 가능한 텍스트를 준비하세요."
        )

    return spec


def _slide_to_spec(
    slide,
    slide_no: int,
    chapter_no: int,
    *,
    assets_dir: Path | None = None,
    deck_title: str = "",
) -> dict:
    if is_image_only_slide(slide):
        spec = _spec_from_image_slide(slide, slide_no, assets_dir=assets_dir, deck_title=deck_title)
        spec["_ingest"]["layout_reason"] = "google_image_slide"
        return _attach_ingest_metadata(spec, slide, slide_no)

    layout_name = slide.slide_layout.name
    blocks = slide_text_blocks(slide)

    if not blocks:
        return _attach_ingest_metadata(
            {
                "layout": "내지_거버닝 O",
                "texts": [f"슬라이드 {slide_no}", "", ""],
                "_ingest": {"layout_reason": "empty_slide"},
            },
            slide,
            slide_no,
        )

    # Chapter divider (간지)
    if layout_name == "간지" or (len(blocks) <= 2 and all(b["len"] < 30 for b in blocks)):
        title = max(blocks, key=lambda b: b["len"])["text"]
        return _attach_ingest_metadata(
            {
                "layout": "간지",
                "texts": [title, f"{chapter_no}."],
                "_ingest": {"layout_reason": "layout_name_간지" if layout_name == "간지" else "heuristic_section"},
            },
            slide,
            slide_no,
        )

    usable = [
        b
        for b in blocks
        if not _is_page_number(b["text"])
        and not _is_series_banner(b["text"])
        and b["len"] > 1
    ]
    if not usable:
        usable = blocks
    has_code_block = any(_contains_code_block(b["text"]) for b in usable)

    if 1 <= len(usable) <= 6:
        title_text, body, subtitle = _pick_title_and_body_from_blocks(usable)
        if body:
            if subtitle:
                return _attach_ingest_metadata(
                    {
                        "layout": "내지_거버닝 O",
                        "texts": [title_text, subtitle, body],
                        "_ingest": {"layout_reason": "title_subtitle_body_stack"},
                    },
                    slide,
                    slide_no,
                    code_block=has_code_block,
                )
            return _attach_ingest_metadata(
                {
                    "layout": "1_내지_거버닝 X",
                    "texts": [title_text, body],
                    "_ingest": {"layout_reason": "title_body_stack"},
                },
                slide,
                slide_no,
                code_block=has_code_block,
            )

    categories = [b for b in usable if _classify_category(b["text"]) or 0.45 <= b["top"] <= 0.75]
    titles = [
        b
        for b in usable
        if 0.85 <= b["top"] <= TITLE_TOP_MAX_IN
        and b not in categories
        and not _is_sidebar_label_block(b)
        and len(_first_line(b["text"])) <= TITLE_LINE_MAX_CHARS
    ]
    bodies = [b for b in usable if b not in categories and b not in titles]

    if titles:
        titles.sort(key=lambda b: (b["top"], len(_first_line(b["text"]))))
    title_text = (
        _first_line(titles[0]["text"])
        if titles
        else (_first_line(usable[0]["text"])[:TITLE_LINE_MAX_CHARS] if usable else f"슬라이드 {slide_no}")
    )
    if categories and title_text == categories[0]["text"] and len(titles) > 0:
        title_text = titles[0]["text"]
    elif categories and title_text == categories[0]["text"] and len(usable) > 1:
        for b in usable:
            if b["text"] != title_text and b["len"] > len(title_text) * 0.5:
                title_text = b["text"]
                break
    if len(titles) > 1:
        extra_title = "\n".join(t["text"] for t in titles[1:])
        bodies = [
            {"text": extra_title, "top": 1.6, "left": 0.0, "len": len(extra_title)}
        ] + [
            {"text": b["text"], "top": b["top"], "left": b.get("left", 0.0), "len": b["len"]}
            for b in bodies
        ]

    subtitle = categories[0]["text"] if categories else ""
    if not subtitle and usable:
        for b in usable:
            if b["text"] != title_text and b["len"] < 24:
                subtitle = b["text"]
                break
    if not subtitle:
        subtitle_candidates = [
            b
            for b in usable
            if b["text"] != title_text
            and 0.70 <= b["top"] <= 1.15
            and len(_first_line(b["text"])) <= 120
        ]
        if subtitle_candidates:
            subtitle = _first_line(subtitle_candidates[0]["text"])
    if subtitle.strip() == title_text.strip():
        subtitle = ""

    body_parts: list[str] = []
    for b in sorted(bodies, key=lambda x: (x["top"], x.get("left", 0.0))):
        t = b["text"].strip()
        if not t or t == title_text or t == subtitle:
            continue
        if t in body_parts:
            continue
        body_parts.append(t)
    body = "\n\n".join(body_parts).strip()

    if not body and len(usable) > 1:
        rest = [b["text"] for b in usable if b["text"] not in (title_text, subtitle)]
        body = "\n\n".join(rest)

    if not subtitle or (len(body) > 400 and "\n\n" not in body[:200]):
        merged = body or subtitle or title_text
        return _attach_ingest_metadata(
            {
                "layout": "1_내지_거버닝 X",
                "texts": [title_text, merged],
                "_ingest": {"layout_reason": "long_body_or_no_governing"},
            },
            slide,
            slide_no,
            code_block=has_code_block,
        )

    if _should_promote_speaker_script_heading(title_text, subtitle, body):
        body = f"{title_text}\n\n{body}".strip()
        title_text, subtitle = subtitle, title_text

    return _attach_ingest_metadata(
        {
            "layout": "내지_거버닝 O",
            "texts": [title_text, subtitle or " ", body or " "],
            "_ingest": {"layout_reason": "title_subtitle_body"},
        },
        slide,
        slide_no,
        code_block=has_code_block,
    )


def convert_presentation(
    src: Path,
    *,
    deck_title: str,
    deck_subtitle: str,
    assets_dir: Path | None = None,
    include_default_front_matter: bool = True,
    default_toc_text: str | None = None,
    front_matter_mode: str = "auto",
) -> list[dict]:
    """Convert source slides to academy JSON specs.

    front_matter_mode:
      - ``auto`` (default when include_default_front_matter): extract cover/TOC from source
      - ``generic``: legacy PaaS-style cover + hardcoded Docker/K8s TOC (deprecated)
      - ``none``: no extra slides; convert all source slides only
    """
    prs = Presentation(str(src))
    specs: list[dict] = []
    skip_indices: set[int] = set()

    if include_default_front_matter:
        if front_matter_mode == "generic":
            toc = default_toc_text or (
                "1. 컨테이너 실습 (Docker)\n"
                "2. Kubernetes 실습 (K8s)\n"
                "3. 실습 시나리오·CLI/GUI 따라하기"
            )
            specs.extend(
                [
                    {"layout": "2_표지", "texts": [f"{deck_title}\n{deck_subtitle}"]},
                    {"layout": "목차", "texts": [toc]},
                ]
            )
        elif front_matter_mode == "auto":
            fm_specs, skip_indices = extract_front_matter(
                prs, deck_title=deck_title, deck_subtitle=deck_subtitle
            )
            specs.extend(fm_specs)
        else:
            specs.append({"layout": "2_표지", "texts": [f"{deck_title}\n{deck_subtitle}"]})

    chapter = 0
    for i, slide in enumerate(prs.slides, start=1):
        if (i - 1) in skip_indices:
            continue
        if slide.slide_layout.name == "간지":
            chapter += 1
        spec = _slide_to_spec(
            slide,
            i,
            chapter,
            assets_dir=assets_dir,
            deck_title=deck_title if i == 1 else "",
        )
        specs.append(spec)
    return specs


def analyze_specs(specs: list[dict]) -> list[dict]:
    """Per-spec summary for /analyze API and smoke validation."""
    rows: list[dict] = []
    for idx, spec in enumerate(specs, start=1):
        ing = spec.get("_ingest") or {}
        texts = spec.get("texts") or []
        preview = [str(t)[:80] + ("…" if len(str(t)) > 80 else "") for t in texts[:3]]
        rows.append(
            {
                "slide": idx,
                "layout": spec.get("layout"),
                "texts_count": len(texts),
                "texts_preview": preview,
                "layout_reason": ing.get("layout_reason"),
                "ingest_kind": ing.get("kind"),
                "background_image": bool(spec.get("background_image")),
                "validation": validate_spec(spec, slide_index=idx),
            }
        )
    return rows


def main() -> int:
    parser = argparse.ArgumentParser(description="Convert legacy deck to academy JSON")
    parser.add_argument("source", type=Path, help="Source .pptx path")
    parser.add_argument(
        "--out-json",
        type=Path,
        default=ROOT / "docs/examples/paas-partner-converted.json",
    )
    parser.add_argument("--assets-dir", type=Path, default=None, help="Extract slide images here")
    parser.add_argument("--title", default="[파트너사] 02. 클라우드 구현기술(PaaS)")
    parser.add_argument("--subtitle", default="클라우드 구현기술 실습")
    parser.add_argument("--no-front-matter", action="store_true")
    args = parser.parse_args()

    if not args.source.is_file():
        print(f"Not found: {args.source}", file=sys.stderr)
        return 2

    specs = convert_presentation(
        args.source,
        deck_title=args.title,
        deck_subtitle=args.subtitle,
        assets_dir=args.assets_dir,
        include_default_front_matter=not args.no_front_matter,
    )
    args.out_json.parent.mkdir(parents=True, exist_ok=True)
    args.out_json.write_text(json.dumps(specs, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(f"Wrote {len(specs)} slide specs → {args.out_json}")
    if args.assets_dir:
        print(f"Assets: {args.assets_dir}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
