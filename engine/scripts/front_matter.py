#!/usr/bin/env python3
"""Extract cover / TOC from source deck for spec pipeline (no hardcoded PaaS TOC)."""

from __future__ import annotations

import re

from pptx import Presentation

from scripts.pptx_ingest import iter_shapes, slide_notes_text, slide_text_blocks


_CHAPTER_LINE_RE = re.compile(r"^\d+[\.\)]\s+\S")
# 본문 불릿(• 1976년: …)은 목차가 아님 — 챕터 색인만 인식
_EVENT_BULLET_RE = re.compile(r"^\d{4}\s*년|[;:]")


def _numbered_toc_lines(text: str) -> list[str] | None:
    """Chapter index: at least two lines like ``1. Section title`` (not • event bullets)."""
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    if len(lines) < 2:
        return None
    chapter_lines = [ln for ln in lines if _CHAPTER_LINE_RE.match(ln)]
    if len(chapter_lines) >= 2:
        return chapter_lines
    # Optional: short • lines without year/colon detail (rare outline-only TOC)
    bullet_lines = [
        ln
        for ln in lines
        if ln.startswith("•")
        and len(ln) < 56
        and not _EVENT_BULLET_RE.search(ln.lstrip("•").strip())
    ]
    if len(bullet_lines) >= 3:
        return bullet_lines
    return None


def _is_cover_like(blocks: list[dict]) -> bool:
    if not blocks:
        return False
    if len(blocks) <= 4 and blocks[0]["len"] >= 4:
        return True
    if len(blocks) <= 16 and blocks[0].get("top", 9.0) < 1.2:
        merged = "\n".join(str(block["text"]) for block in blocks).lower()
        if "dashboard" in merged and ("실습" in merged or "nginx" in merged):
            return True
    return len(blocks) <= 2


def _is_toc_like(blocks: list[dict]) -> bool:
    if not blocks:
        return False
    if any(b["len"] > 120 for b in blocks):
        return False
    merged = "\n".join(b["text"] for b in blocks)
    return _numbered_toc_lines(merged) is not None


def _is_cover_subtitle_candidate(text: str) -> bool:
    """Return True when a line looks like a subtitle, not body copy."""
    t = (text or "").strip()
    if not t:
        return False
    if len(t) > 40:
        return False
    if t.endswith((".", "!", "?")):
        return False
    if "다루기까지" in t or "되기까지" in t or "제조사가" in t:
        return False
    return True


def _compose_cover_text(blocks: list[dict]) -> str:
    """Title + optional short subtitle (first line of 2nd block only)."""
    if not blocks:
        return ""
    title = blocks[0]["text"].strip()
    if len(blocks) < 2:
        return title
    lines = [ln.strip() for ln in blocks[1]["text"].splitlines() if ln.strip()]
    if lines and _is_cover_subtitle_candidate(lines[0]):
        return f"{title}\n{lines[0]}"
    return title


def _attach_source_notes(spec: dict, slide, source_index: int) -> dict:
    ing = dict(spec.get("_ingest") or {})
    ing["source_slide_index"] = source_index
    notes = slide_notes_text(slide)
    if notes:
        ing["speaker_notes"] = notes
    if _is_visual_heavy_slide(slide):
        ing["visual_preservation"] = "native_shapes"
    spec["_ingest"] = ing
    return spec


def _is_visual_heavy_slide(slide) -> bool:
    text_shapes = 0
    visual_shapes = 0
    for shape, depth in iter_shapes(slide.shapes):
        if depth > 0:
            visual_shapes += 1
        if shape.is_placeholder:
            continue
        if shape.has_text_frame and (shape.text or "").strip():
            text_shapes += 1
        visual_shapes += 1
    return visual_shapes >= 10 and text_shapes >= 5


def extract_front_matter(
    prs: Presentation,
    *,
    deck_title: str,
    deck_subtitle: str,
) -> tuple[list[dict], set[int]]:
    """Build cover/toc specs and source slide indices to skip when converting body."""
    specs: list[dict] = []
    skip: set[int] = set()

    if not prs.slides:
        specs.append({"layout": "2_표지", "texts": [f"{deck_title}\n{deck_subtitle}"]})
        return specs, skip

    s0 = prs.slides[0]
    b0 = slide_text_blocks(s0)
    layout0 = s0.slide_layout.name

    if layout0 in ("2_표지", "표지") or _is_cover_like(b0):
        if b0:
            cover = _compose_cover_text(b0)
            specs.append(
                _attach_source_notes(
                    {"layout": "2_표지", "texts": [cover], "_ingest": {"kind": "source_cover"}},
                    s0,
                    0,
                )
            )
        else:
            specs.append(
                _attach_source_notes(
                    {
                        "layout": "2_표지",
                        "texts": [f"{deck_title}\n{deck_subtitle}"],
                        "_ingest": {"kind": "generated_cover"},
                    },
                    s0,
                    0,
                )
            )
        skip.add(0)

    if len(prs.slides) > 1:
        s1 = prs.slides[1]
        b1 = slide_text_blocks(s1)
        layout1 = s1.slide_layout.name
        merged1 = "\n".join(x["text"] for x in b1)

        if layout1 in ("목차", "TOC") or _is_toc_like(b1):
            toc_lines = _numbered_toc_lines(merged1) or [x["text"] for x in b1]
            toc_text = "\n".join(toc_lines)
            specs.append(
                _attach_source_notes(
                    {"layout": "목차", "texts": [toc_text], "_ingest": {"kind": "source_toc"}},
                    s1,
                    1,
                )
            )
            skip.add(1)

    if not specs:
        specs.append(
            {
                "layout": "2_표지",
                "texts": [f"{deck_title}\n{deck_subtitle}"],
                "_ingest": {"kind": "generated_cover"},
            }
        )

    return specs, skip
