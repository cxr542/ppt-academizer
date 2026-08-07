#!/usr/bin/env python3
"""Migrate CMP source deck into OKESTRO academy layout (academy-design.md §7)."""
from __future__ import annotations

import argparse
import importlib.util
import re
import shutil
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PPT_TEST = Path(__file__).resolve().parent.parent
if str(PPT_TEST) not in sys.path:
    sys.path.insert(0, str(PPT_TEST))

from scripts.academy_brand import (  # noqa: E402
    brand_rgb,
    ensure_brand_pictures,
    load_brand_pack,
    normalize_brand_accents,
    replace_matched_icons,
    resolve_brand_dir,
)
from scripts.academy_layout_patterns import (  # noqa: E402
    apply_card_column_pattern,
    detect_card_column_count,
)
from scripts.academy_template import resolve_academy_template_path  # noqa: E402
from scripts.deck_migrate_config import (  # noqa: E402
    DeckMigrateConfig,
    migrate_config_for_source,
)
from scripts.pptx_ingest import iter_shapes, iter_shapes_abs, slide_notes_text  # noqa: E402
from scripts.llm_classifier import analyze_slide_with_llm  # noqa: E402

import scripts.academy_deck_build_lib as adl  # noqa: E402

from scripts.academy_deck_build_lib import (  # noqa: E402
    FONT_BODY,
    FONT_TITLE,
    _assign_plain,
    _autofit_horizontal_textbox,
    _capture_layout_placeholder_geometry,
    _capture_seed_placeholder_geometry,
    _fill_chapter_guide,
    _fill_contents_guide,
    _fill_cover,
    _ph_by_idx,
    _prepare_slide_for_editing,
    _restore_slide_placeholder_geometry,
    configure_text_frame_for_wrap,
    delete_slide,
    duplicate_slide_from_seed,
    layout_seed_slides,
    save_academy_deck,
    set_tf_font,
)

# design.md — OKESTRO theme text colors
CLR_DK1 = RGBColor(0x00, 0x00, 0x00)
CLR_DK2 = RGBColor(0x44, 0x54, 0x6A)
CLR_ACCENT1 = RGBColor(0x00, 0x6D, 0xFF)

_SHAPE_LIB = PPT_TEST / "scripts" / "shape_migrate_lib.py"
_SPEC = importlib.util.spec_from_file_location("bac", _SHAPE_LIB)
bac = importlib.util.module_from_spec(_SPEC)
assert _SPEC.loader is not None
_SPEC.loader.exec_module(bac)

# CMP diagrams use large exported PNG panels — not PaaS gray card backgrounds.
DROP_CARD_PANEL_PICTURES = False

DEFAULT_SOURCE = Path.home() / "Desktop" / "클라우드 구현기술(CMP)_v1.0_수정요청.pptx"

# Legacy aliases (CMP defaults — prefer DeckMigrateConfig).
SKIP_SRC_INDICES = frozenset({8})
CMP_TOC_SLIDE_RANGE = range(2, 8)

_HANGUL = re.compile(r"[\uAC00-\uD7A3]")
_STEP_NUMBER_LINE = re.compile(r"^\d{1,3}$")

# Academy template layouts — used to detect double-migration of an output deck.
ACADEMY_LAYOUT_MARKERS = frozenset(
    {"2_표지", "내지_거버닝 O", "1_내지_거버닝 X", "간지", "목차"}
)

# Official layout coords (내지_거버닝 O) — do not use PaaS-wide title width.
LAYOUT_PH10_LEFT = 781_263
CONTENT_BODY_LEFT = 962_025
CONTENT_COLUMN_GAP = 360_000
CONTENT_RIGHT_COL_RATIO = 0.47  # right column starts ~47% of slide width
TITLE_PT = 18
TITLE_MAX_CHARS = 72
TITLE_SIDEBAR_LEFT_MAX_IN = 0.32
BODY_SHAPE_PT = 12
CAPTION_SHAPE_PT = 12

LAYOUT_COVER = "2_표지"
LAYOUT_TOC = "목차"  # 간지1
LAYOUT_SECTION = "간지"  # 간지2
LAYOUT_CONTENT = "내지_거버닝 O"

HEADER_TITLE_TOP_MAX = 900_000
SECTION_PH10_MIN_WIDTH = int(Inches(3.6))
# Radial hub diagrams (CMP 6-tech slide): not 2-column body — skip column relayout.
RADIAL_MIN_BODY_LABELS = 5
FULL_BLEED_W_RATIO = 0.88
FULL_BLEED_H_RATIO = 0.55
TOP_BAND_PICTURE_RATIO = 0.45


def uses_academy_layouts(prs: Presentation) -> bool:
    """Source deck is authored on OKESTRO academy masters (allowed to migrate)."""
    names = {s.slide_layout.name for s in prs.slides}
    return len(names & ACADEMY_LAYOUT_MARKERS) >= 2


def is_likely_reacademize_output(prs: Presentation, *, filename_hint: str = "") -> bool:
    """True only for ppt-academizer *output* re-uploaded — not academy-template originals."""
    low = (filename_hint or "").lower().replace(" ", "-")
    for marker in (
        "academy-output",
        "academy-cmp-",
        "academy-cmp-academy",
        "_ppfixed",
        "_repaired",
    ):
        if marker in low:
            return True

    if not uses_academy_layouts(prs):
        return False

    dup_num_headers = 0
    content_slides = 0
    for slide in prs.slides:
        if slide.slide_layout.name != LAYOUT_CONTENT:
            continue
        content_slides += 1
        t10 = t12 = ""
        for sh in slide.placeholders:
            idx = sh.placeholder_format.idx
            if idx == 10:
                t10 = (sh.text or "").strip()
            elif idx == 12:
                t12 = (sh.text or "").strip()
        if t10 and t10 == t12 and _is_step_number_line(t10):
            dup_num_headers += 1
    if content_slides and dup_num_headers / content_slides >= 0.45:
        return True
    return False


# Back-compat alias (narrower than v1.3.4 guard).
is_academy_template_deck = is_likely_reacademize_output


def _is_step_number_line(line: str) -> bool:
    """Google/partner step column or re-migrated ph12 — not a slide title."""
    return bool(_STEP_NUMBER_LINE.match((line or "").strip()))


def _hangul_ratio(text: str) -> float:
    compact = re.sub(r"\s+", "", text or "")
    if not compact:
        return 0.0
    return len(_HANGUL.findall(compact)) / len(compact)


def _slide_has_picture(slide) -> bool:
    return any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in slide.shapes)


def is_picture_only_slide(slide) -> bool:
    """True when the slide has pictures and no non-empty text (full-frame visual)."""
    has_pic = False
    for sh in slide.shapes:
        if sh.has_text_frame and (sh.text or "").strip():
            return False
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            has_pic = True
    return has_pic


_TOC_TITLE_MARKERS = frozenset(
    {
        "목차",
        "CONTENTS",
        "Contents",
        "Table of Contents",
        "TABLE OF CONTENTS",
    }
)


def classify_slide(
    slide,
    index: int,
    part_cover_indices: frozenset[int] | None = None,
) -> str:
    covers = part_cover_indices if part_cover_indices is not None else frozenset({0, 13})
    if index in covers:
        return "cover"

    layout_info = analyze_slide_layout(slide)
    layout = layout_info["layout"]

    if layout == "empty" and _slide_has_picture(slide):
        return "content"
    return layout


_GOV_METAPHOR = re.compile(r"(처럼|듯이|하듯|듯,|듯 )")
_GOV_SKIP_PREFIXES = ("개념 정의", "가치:", "가치 ：", "가치 ")


def _normalize_header_line(text: str) -> str:
    return text.replace("\x0b", " ").replace("\n", " ").strip()


def _pick_content_governing(
    texts: list[tuple[int, str]], title: str | None
) -> str | None:
    """Pick ph13 governing: metaphor line, else McKinsey-style header kicker."""
    gov_max = int(getattr(bac, "GOV_MAX_LEN", 96))
    ranked: list[tuple[int, int, str]] = []
    header_kickers: list[str] = []
    for top, text in texts:
        line = _normalize_header_line(text)
        if not line or line == title or _is_step_number_line(line):
            continue
        if any(line.startswith(p) for p in _GOV_SKIP_PREFIXES):
            continue
        if len(line) < 12 or len(line) > gov_max:
            continue
        if top <= HEADER_TITLE_TOP_MAX:
            header_kickers.append(line)
        score = 0
        if _GOV_METAPHOR.search(line):
            score += 10
        if line.endswith(("습니다.", "합니다.", "합니다")) or "습니다" in line:
            score += 2
        if score:
            ranked.append((score, len(line), line))
    if ranked:
        ranked.sort(key=lambda c: (-c[0], -c[1]))
        best_score, _len, best = ranked[0]
        if best_score >= 10:
            return best
    # Action-title decks: second header-band line is the governing kicker.
    return header_kickers[0] if header_kickers else None


def extract_header(src_slide, kind: str) -> tuple[str | None, str | None]:
    layout_info = analyze_slide_with_llm(src_slide)
    title = layout_info.get("title")
    gov = layout_info.get("governing")

    if not title and layout_info.get("blocks"):
        title = _normalize_header_line(layout_info["blocks"][0]["text"])

    if title:
        title = _normalize_header_line(title)
    if gov:
        gov = _normalize_header_line(gov)

    return title, gov


def build_slide_plan(
    src: Presentation,
    cfg: DeckMigrateConfig | None = None,
) -> list[tuple[int, str]]:
    if cfg is None:
        cfg = migrate_config_for_source("cmp-partner.pptx")
    plan: list[tuple[int, str]] = []
    for i, slide in enumerate(src.slides):
        if i in cfg.skip_src_indices:
            continue
        if i == 0:
            plan.append((i, "cover"))
            continue
        kind = classify_slide(slide, i, cfg.part_cover_indices)
        if kind == "empty":
            plan.append((i, "content"))
        else:
            plan.append((i, kind))
    return plan


def expected_output_slide_count(plan: list[tuple[int, str]], cfg: DeckMigrateConfig) -> int:
    extra_toc = sum(1 for _, k in plan if k == "cover" and cfg.auto_toc_after_cover)
    return len(plan) + extra_toc


def collect_section_titles(
    src: Presentation,
    after_src_idx: int,
    *,
    part_cover_indices: frozenset[int] | None = None,
) -> list[str]:
    """Section titles until the next part cover (간지1 목차용)."""
    covers = part_cover_indices or frozenset({0, 13})
    titles: list[str] = []
    for i in range(after_src_idx, len(src.slides)):
        if i != after_src_idx and classify_slide(src.slides[i], i, covers) == "cover":
            break
        if classify_slide(src.slides[i], i, covers) == "section":
            t, _ = extract_header(src.slides[i], "section")
            if t:
                titles.append(t)
    return titles


def is_full_bleed_picture(shape, slide_width: int, slide_height: int) -> bool:
    """Drop slide-wide export backgrounds only — keep right-column hero diagrams."""
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    w, h = int(shape.width or 0), int(shape.height or 0)
    top, left = int(shape.top or 0), int(shape.left or 0)
    if w >= slide_width * FULL_BLEED_W_RATIO and h >= slide_height * FULL_BLEED_H_RATIO:
        return True
    # Wide top banner (nearly full slide width), not ~50% column art.
    if top < 200_000 and w >= slide_width * 0.75:
        return True
    # Top-left logo strip from Google export.
    if top < 200_000 and left < slide_width * 0.08 and w >= slide_width * 0.35:
        return True
    return False


def is_card_panel_picture(shape, slide_height: int) -> bool:
    """Rounded gray card panels exported as large pictures — drop on white academy slides."""
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    w, h = int(shape.width or 0), int(shape.height or 0)
    top = int(shape.top or 0)
    if top < bac.HEADER_BOTTOM_EMU:
        return False
    if w >= 2_800_000 and h >= 1_800_000:
        return True
    return False


def _walk_shapes(container):
    """Depth-first over slide or group shapes."""
    for shape in container.shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape)


def is_slide_background_picture(shape, slide_width: int, slide_height: int) -> bool:
    """Only drop true slide backgrounds (§6.7) — not chart/diagram images."""
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    w, h = int(shape.width or 0), int(shape.height or 0)
    top, left = int(shape.top or 0), int(shape.left or 0)
    if w >= slide_width * 0.97 and h >= slide_height * 0.93:
        return True
    if top < 200_000 and w >= slide_width * 0.80:
        return True
    if top < 200_000 and left < slide_width * 0.08 and w >= slide_width * 0.35:
        return True
    return False


def _is_diagram_raster_shape(shape, slide_width: int, slide_height: int) -> bool:
    if getattr(shape, "has_chart", False) and shape.has_chart:
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        return True
    if bac._is_graphic_frame(shape):
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return any(
            _is_diagram_raster_shape(child, slide_width, slide_height)
            for child in shape.shapes
        )
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return not is_slide_background_picture(shape, slide_width, slide_height)
    return False


def count_diagram_rasters(slide, slide_width: int, slide_height: int) -> int:
    return sum(
        1
        for shp in _walk_shapes(slide)
        if _is_diagram_raster_shape(shp, slide_width, slide_height)
    )


def is_image_primary_slide(src_slide, slide_width: int, slide_height: int) -> bool:
    """§6.7 — slide is mostly diagram/chart pictures with little body text."""
    rasters = count_diagram_rasters(src_slide, slide_width, slide_height)
    if rasters < 1:
        return False
    body_texts = 0
    for sh in _walk_shapes(src_slide):
        if not sh.has_text_frame or not (sh.text or "").strip():
            continue
        if int(sh.top or 0) <= bac.HEADER_BOTTOM_EMU:
            continue
        body_texts += 1
    return body_texts <= 2


def is_raster_heavy_slide(src_slide, slide_width: int, slide_height: int) -> bool:
    """Permissive §6.7 copy when diagrams dominate (incl. slides 11, 18, 23, 24)."""
    rasters = count_diagram_rasters(src_slide, slide_width, slide_height)
    if rasters == 0:
        return False
    if rasters >= 2:
        return True
    return is_image_primary_slide(src_slide, slide_width, slide_height)


_STOCK_PLACEHOLDER_DIMS = frozenset(
    {
        (400, 400),
        (600, 400),
        (400, 600),
        (640, 480),
        (480, 640),
        (800, 600),
        (600, 800),
    }
)


def is_stock_placeholder_picture(shape) -> bool:
    """Detect Google-export stubs (solid black / classic 600×400 placeholders)."""
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    try:
        from PIL import Image
        import io

        blob = shape.image.blob
        im = Image.open(io.BytesIO(blob)).convert("RGB")
    except Exception:
        return False
    w, h = im.size
    if w <= 64 and h <= 64:
        return False
    sample = im.resize((48, 48))
    colors = len(set(sample.getdata()))
    if colors <= 4:
        return True
    if (w, h) in _STOCK_PLACEHOLDER_DIMS and colors < 200 and len(blob) < 40_000:
        return True
    return False


def should_drop_body_picture(
    shape,
    slide_width: int,
    slide_height: int,
    *,
    permissive_raster: bool,
    keep_full_frame_image: bool = False,
) -> bool:
    """§6.7 — on diagram-heavy slides only drop slide backgrounds.

    ``keep_full_frame_image``: picture-only slides where the full-frame image *is*
    the content (e.g. VMware winback chart page) — do not drop it as a background.
    """
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    if is_stock_placeholder_picture(shape):
        return True
    is_bg = is_slide_background_picture(shape, slide_width, slide_height)
    is_bleed = is_full_bleed_picture(shape, slide_width, slide_height)
    if keep_full_frame_image and (is_bg or is_bleed):
        return False
    if permissive_raster:
        return is_bg
    if is_bg or is_bleed:
        return True
    if DROP_CARD_PANEL_PICTURES and is_card_panel_picture(shape, slide_height):
        return True
    return False


def _is_top_level_raster_candidate(shape) -> bool:
    if shape.shape_type in (
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.GROUP,
        MSO_SHAPE_TYPE.CHART,
    ):
        return True
    if getattr(shape, "has_chart", False) and shape.has_chart:
        return True
    return bac._is_graphic_frame(shape)


def _dst_has_similar_bbox(
    dst_slide,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    tol: int = 120_000,
) -> bool:
    for sh in dst_slide.shapes:
        if sh.is_placeholder:
            continue
        if int(sh.left or 0) - tol <= left <= int(sh.left or 0) + tol:
            if int(sh.top or 0) - tol <= top <= int(sh.top or 0) + tol:
                if abs(int(sh.width or 0) - width) < tol * 2:
                    return True
    return False


def sync_raster_diagram_assets(
    src_slide,
    dst_slide,
    title,
    governing,
    slide_width: int,
    slide_height: int,
    *,
    skip_connectors: bool = False,
    force_permissive: bool = False,
    src_classify_empty: bool = False,
) -> int:
    """§6.7 — copy missing chart/diagram images when dst has fewer rasters than src."""
    src_n = count_diagram_rasters(src_slide, slide_width, slide_height)
    dst_n = count_diagram_rasters(dst_slide, slide_width, slide_height)
    if src_n == 0:
        return 0
    keep_full = is_picture_only_slide(src_slide)
    permissive = (
        force_permissive
        or src_classify_empty
        or keep_full
        or is_raster_heavy_slide(src_slide, slide_width, slide_height)
        or is_image_primary_slide(src_slide, slide_width, slide_height)
    )
    # Picture-only full-frame slides report src_n=0 (image counted as background).
    if not keep_full and dst_n >= src_n and not (src_classify_empty and src_n > 0):
        return 0

    id_remap: dict[str, str] = {}
    copied = 0
    for shape in src_slide.shapes:
        if not _is_top_level_raster_candidate(shape):
            continue
        if not bac.is_body_shape(shape, title, governing):
            continue
        if should_drop_body_picture(
            shape,
            slide_width,
            slide_height,
            permissive_raster=permissive,
            keep_full_frame_image=keep_full,
        ):
            continue
        if skip_connectors and (
            bac._is_cxn_sp(shape) or shape.shape_type == MSO_SHAPE_TYPE.LINE
        ):
            continue
        left, top = int(shape.left or 0), int(shape.top or 0)
        w, h = int(shape.width or 0), int(shape.height or 0)
        if _dst_has_similar_bbox(dst_slide, left, top, w, h):
            continue
        n_before = len(dst_slide.shapes)
        bac.copy_shape_hybrid(dst_slide, shape, id_remap=id_remap)
        if len(dst_slide.shapes) > n_before:
            copied += 1
    if copied:
        renumber_map = bac.renumber_shape_ids(dst_slide)
        bac.fix_orphan_connector_refs(dst_slide, id_remap, renumber_map)
    return copied


def is_radial_diagram_slide(slide, slide_width: int, slide_height: int) -> bool:
    """Hub / 6-around-center layouts (not 2-column text). CMP slide 2 is the canonical case."""
    body: list[tuple[int, int]] = []
    for sh in slide.shapes:
        if not sh.has_text_frame or not (sh.text or "").strip():
            continue
        top = int(sh.top or 0)
        if top <= bac.HEADER_BOTTOM_EMU:
            continue
        body.append((int(sh.left or 0), top))
    if len(body) < RADIAL_MIN_BODY_LABELS:
        return False
    lefts = [p[0] for p in body]
    tops = [p[1] for p in body]
    if max(lefts) - min(lefts) < slide_width * 0.35:
        return False
    if max(tops) - min(tops) < slide_height * 0.22:
        return False
    mid_x = slide_width / 2
    bands = {"left": 0, "right": 0, "center": 0}
    for left, _top in body:
        if left < mid_x - slide_width * 0.12:
            bands["left"] += 1
        elif left > mid_x + slide_width * 0.12:
            bands["right"] += 1
        else:
            bands["center"] += 1
    return bands["left"] >= 2 and bands["right"] >= 1 and bands["center"] >= 1


def copy_body_shapes_filtered(
    src_slide,
    dst_slide,
    title,
    governing,
    slide_width: int,
    slide_height: int,
    *,
    skip_connectors: bool = False,
) -> dict[str, str]:
    permissive = is_raster_heavy_slide(src_slide, slide_width, slide_height)
    keep_full = is_picture_only_slide(src_slide)
    id_remap: dict[str, str] = {}
    for shape in src_slide.shapes:
        if not bac.is_body_shape(shape, title, governing):
            continue
        if should_drop_body_picture(
            shape,
            slide_width,
            slide_height,
            permissive_raster=permissive,
            keep_full_frame_image=keep_full,
        ):
            continue
        if skip_connectors and (
            bac._is_cxn_sp(shape) or shape.shape_type == MSO_SHAPE_TYPE.LINE
        ):
            continue
        bac.copy_shape_hybrid(dst_slide, shape, id_remap=id_remap)
    return id_remap


def _luminance(rgb: RGBColor) -> float:
    return 0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]


def academy_color_for_run(run, *, title: bool = False) -> None:
    """Map Google-export light-on-dark text to design.md colors on white slides."""
    if title:
        run.font.color.rgb = CLR_DK1
        return
    try:
        c = run.font.color
        if c.type == MSO_COLOR_TYPE.RGB and c.rgb is not None:
            rgb = c.rgb
            lum = _luminance(rgb)
            if lum >= 220:
                run.font.color.rgb = CLR_DK1
            elif lum >= 140:
                run.font.color.rgb = CLR_DK2
            return
    except Exception:
        pass
    run.font.color.rgb = CLR_DK1


def apply_academy_colors_to_text_frame(tf, *, title: bool = False) -> None:
    configure_text_frame_for_wrap(tf)
    set_tf_font(tf, title=title)
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = FONT_TITLE if title else FONT_BODY
            academy_color_for_run(run, title=title)


def apply_academy_fonts_and_colors(slide, *, placeholders_only: bool = False) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if placeholders_only and not shape.is_placeholder:
            continue
        is_title_ph = shape.is_placeholder and shape.placeholder_format.idx == 10
        is_num_ph = shape.is_placeholder and shape.placeholder_format.idx == 12
        if is_title_ph:
            continue
        apply_academy_colors_to_text_frame(
            shape.text_frame,
            title=False,
        )
        if is_num_ph:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = CLR_DK2


def _apply_body_text_frame_style(tf, *, caption: bool = False) -> None:
    pt = CAPTION_SHAPE_PT if caption else BODY_SHAPE_PT
    font_name = FONT_TITLE if caption else FONT_BODY
    configure_text_frame_for_wrap(tf)
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = font_name
            run.font.size = Pt(pt)
            run.font.color.rgb = CLR_DK1
            if caption:
                run.font.bold = True


def has_designed_chrome_layout(src_slide) -> bool:
    """McKinsey-style solid cards / accent bars — keep source type hierarchy."""
    solid_panels = 0
    accent_bars = 0
    for shape in src_slide.shapes:
        if getattr(shape, "is_placeholder", False):
            continue
        try:
            if shape.fill.type != MSO_FILL.SOLID:
                continue
        except Exception:
            continue
        w, h = int(shape.width or 0), int(shape.height or 0)
        text = shape.text.strip() if shape.has_text_frame else ""
        if text:
            continue
        if w >= 1_000_000 and h >= 1_000_000:
            solid_panels += 1
        elif w < 250_000 and h > 200_000:
            accent_bars += 1
    return solid_panels >= 1 or accent_bars >= 2


def apply_academy_body_shape_typography(slide) -> None:
    """§6.6 — non-placeholder body shapes and shapes inside GROUPs."""

    def walk(shapes) -> None:
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes)
                continue
            if shape.is_placeholder:
                idx = shape.placeholder_format.idx
                if idx in (10, 12, 13):
                    continue
            if not shape.has_text_frame:
                continue
            text = (shape.text or "").strip()
            if not text:
                continue
            first_line = text.replace("\x0b", " ").split("\n")[0]
            caption = len(first_line) <= 40 or (len(text) <= 80 and text.count("\n") < 2)
            _apply_body_text_frame_style(shape.text_frame, caption=caption)

    walk(slide.shapes)


def hide_empty_governing_placeholder(slide) -> None:
    """Remove template '거버닝 메시지…' guide when there is no governing text."""
    for shape in list(slide.placeholders):
        if shape.placeholder_format.idx != 13:
            continue
        if (shape.text or "").strip():
            return
        adl.collapse_content_placeholders(slide, indices=(13,))
        return


def apply_content_header_geometry(slide, slide_width: int | None = None) -> None:
    """Official template placeholder positions (not PaaS-wide title width)."""
    _restore_slide_placeholder_geometry(slide)
    ph10 = _ph_by_idx(slide, 10)
    ph10.left = LAYOUT_PH10_LEFT
    ph10.top = adl.TITLE_ROW_TOP
    ph10.height = max(int(ph10.height or 0), adl.LAB_TITLE_HEIGHT)
    if slide_width:
        max_w = slide_width - LAYOUT_PH10_LEFT - bac.SLIDE_MARGIN_X
        ph10.width = min(int(max_w), int(Inches(10.5)))
    elif int(ph10.width) > int(Inches(5.5)):
        ph10.width = int(Inches(3.0))


def normalize_slide_title(title: str) -> str:
    return " ".join((title or "").split())


def assign_content_slide_title(ph10, title: str) -> None:
    """Title: one line; Korean/English 18pt; left-aligned, vertically centered in box."""
    title = normalize_slide_title(title)
    tf = ph10.text_frame
    configure_text_frame_for_wrap(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ph10.text = ""
    while len(tf.paragraphs) > 1:
        tf._element.remove(tf.paragraphs[-1]._p)

    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    para.clear()

    if ": " in title:
        head, tail = title.split(": ", 1)
        if _HANGUL.search(head) and not _HANGUL.search(tail):
            run = para.add_run()
            run.text = title
            run.font.name = FONT_TITLE
            run.font.size = Pt(TITLE_PT)
            run.font.color.rgb = CLR_DK1
        else:
            r_ko = para.add_run()
            r_ko.text = f"{head}: "
            r_ko.font.name = FONT_TITLE
            r_ko.font.size = Pt(TITLE_PT)
            r_ko.font.color.rgb = CLR_DK1
            r_en = para.add_run()
            r_en.text = tail
            r_en.font.name = FONT_TITLE
            r_en.font.size = Pt(TITLE_PT)
            r_en.font.color.rgb = CLR_DK1
    else:
        run = para.add_run()
        run.text = title
        run.font.name = FONT_TITLE
        run.font.size = Pt(TITLE_PT)
        run.font.color.rgb = CLR_DK1


def _content_ph(slide, idx: int):
    """Placeholder lookup without re-applying seed geometry (avoids undoing width tweaks)."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(f"Placeholder idx={idx} not found on {slide.slide_layout.name!r}")


def apply_slide_title_layout(slide, slide_width: int, title: str) -> None:
    apply_content_header_geometry(slide, slide_width)
    ph10 = _content_ph(slide, 10)
    assign_content_slide_title(ph10, title)
    if slide_width:
        max_w = slide_width - LAYOUT_PH10_LEFT - bac.SLIDE_MARGIN_X
        ph10.width = min(int(max_w), int(Inches(10.5)))


def _body_text_lefts(slide) -> list[int]:
    lefts: list[int] = []
    for sh in slide.shapes:
        if sh.is_placeholder or not sh.has_text_frame:
            continue
        if not (sh.text or "").strip():
            continue
        if int(sh.top or 0) < bac.HEADER_BOTTOM_EMU:
            continue
        lefts.append(int(sh.left or 0))
    return lefts


def count_body_text_columns(slide, slide_width: int, *, gap_ratio: float = 0.08) -> int:
    """Cluster body textboxes by horizontal gap (3-card rows → 3 columns)."""
    lefts = sorted(_body_text_lefts(slide))
    if not lefts:
        return 0
    gap = max(int(slide_width * gap_ratio), 600_000)
    clusters = 1
    prev = lefts[0]
    for left in lefts[1:]:
        if left - prev >= gap:
            clusters += 1
        prev = left
    return clusters


def has_hero_body_picture(slide, slide_width: int, slide_height: int | None = None) -> bool:
    """Large body raster (left/right hero) — do not crush text into a 2-col grid."""
    height = slide_height
    if height is None:
        try:
            height = int(slide.part.package.presentation_part.presentation.slide_height)
        except Exception:
            height = int(slide_width * 9 / 16)
    min_w = int(slide_width * 0.22)
    min_h = int(height * 0.25)
    for sh in slide.shapes:
        if sh.is_placeholder or sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if int(sh.top or 0) < bac.HEADER_BOTTOM_EMU:
            continue
        if int(sh.width or 0) >= min_w and int(sh.height or 0) >= min_h:
            return True
    return False


def relayout_content_columns(slide, slide_width: int) -> None:
    """Two-column body: move right column left and stretch both columns to margins.

    Skip multi-column card rows (3+ horizontal clusters) — forcing 2-col stacks
    titles/bodies on top of each other (VMware winback slide 3).
    Skip hero-image + text splits (VMware slides 8/12) — mid-threshold otherwise
    pulls right-column copy onto the image.
    """
    if is_speaker_script_slide(slide):
        return
    if count_body_text_columns(slide, slide_width) >= 3:
        return
    if has_hero_body_picture(slide, slide_width):
        return
    max_right = slide_width - bac.SLIDE_MARGIN_X
    mid_threshold = int(slide_width * 0.44)
    right_left = int(slide_width * CONTENT_RIGHT_COL_RATIO)
    left_width = right_left - CONTENT_BODY_LEFT - CONTENT_COLUMN_GAP
    if left_width < bac.MIN_SHAPE_W:
        return

    bodies = [
        sh
        for sh in slide.shapes
        if not sh.is_placeholder
        and sh.has_text_frame
        and int(sh.top or 0) >= bac.HEADER_BOTTOM_EMU
        and (sh.text or "").strip()
    ]
    left_shapes = [sh for sh in bodies if int(sh.left or 0) < mid_threshold]
    right_shapes = [sh for sh in bodies if int(sh.left or 0) >= mid_threshold]
    if not right_shapes:
        return

    for sh in left_shapes:
        sh.left = CONTENT_BODY_LEFT
        sh.width = left_width
    right_width = max_right - right_left
    for sh in right_shapes:
        sh.left = right_left
        sh.width = right_width


def is_speaker_script_slide(slide) -> bool:
    texts = [
        (shape.text or "").strip()
        for shape in slide.shapes
        if shape.has_text_frame and (shape.text or "").strip()
    ]
    merged = "\n".join(texts)
    if "강사용" in merged or "강의 포인트" in merged:
        return True
    quoted_lines = sum(1 for text in texts if "“" in text or '"' in text)
    short_labels = sum(1 for text in texts if len(text) <= 16 and "\n" not in text)
    return quoted_lines >= 3 and short_labels >= 3


def copy_speaker_notes(src_slide, dst_slide) -> None:
    notes = slide_notes_text(src_slide)
    if notes:
        dst_slide.notes_slide.notes_text_frame.text = notes


def apply_brand_pack_to_slide(
    slide, pack: dict | None, *, kind: str, icon_hits: list[str] | None = None
) -> None:
    """Stamp brand logos, remap accents, and swap matched icons."""
    if not pack:
        return
    ensure_brand_pictures(slide, pack, kind)
    if kind == "content":
        normalize_brand_accents(slide, pack)
        replaced = replace_matched_icons(slide, pack)
        if icon_hits is not None and replaced:
            icon_hits.extend(replaced)


def fill_cover_from_source(slide, src_slide, cfg: DeckMigrateConfig) -> None:
    title, governing = extract_header(src_slide, "cover")
    cover_text = title or "강의"
    if cfg.cover_subtitle:
        cover_text = f"{cover_text}\n{cfg.cover_subtitle}"
    if governing:
        cover_text = f"{cover_text}\n{governing}"
    _fill_cover(slide, [cover_text], title_font_pt=36, body_font_pt=16)


def apply_academy_table_style(table, *, header_rows: int = 1) -> None:
    """§6.5 — editable cells, academy fonts/colors."""
    for r_idx, row in enumerate(table.rows):
        is_header = r_idx < header_rows
        for cell in row.cells:
            tf = cell.text_frame
            if not tf:
                continue
            apply_academy_colors_to_text_frame(tf, title=is_header)
            if is_header:
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.bold = True
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xE7, 0xE6, 0xE6)
                except Exception:
                    pass


def apply_academy_tables_on_slide(slide) -> None:
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            apply_academy_table_style(shape.table)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if child.shape_type == MSO_SHAPE_TYPE.TABLE:
                    apply_academy_table_style(child.table)


def _count_charts(slide) -> int:
    n = 0
    for sh in slide.shapes:
        if getattr(sh, "has_chart", False) and sh.has_chart:
            n += 1
        elif sh.shape_type == MSO_SHAPE_TYPE.CHART:
            n += 1
    return n


def fill_toc(slide, section_titles: list[str]) -> None:
    lines = [f"• {t}" for t in section_titles] if section_titles else ["• (목차)"]
    _fill_contents_guide(slide, ["\n".join(lines)], title_font_pt=28, body_font_pt=18)


_TOC_SKIP_LINES = frozenset(
    {
        "okestro confidential",
        "confidential",
        "internal use only",
    }
)


def extract_toc_entries(src_slide) -> list[str]:
    """Build TOC bullets from a source 목차 slide (number + title pairs)."""
    texts = [
        (int(sh.top or 0), int(sh.left or 0), sh.text.strip().replace("\x0b", " "))
        for sh in src_slide.shapes
        if sh.has_text_frame and sh.text.strip()
    ]
    texts.sort(key=lambda row: (row[0], row[1]))
    entries: list[str] = []
    i = 0
    while i < len(texts):
        line = texts[i][2].split("\n")[0].strip()
        if line in _TOC_TITLE_MARKERS or line.lower() in _TOC_SKIP_LINES:
            i += 1
            continue
        if re.fullmatch(r"\d{1,2}", line) and i + 1 < len(texts):
            title = texts[i + 1][2].split("\n")[0].strip()
            if (
                title
                and title not in _TOC_TITLE_MARKERS
                and title.lower() not in _TOC_SKIP_LINES
            ):
                entries.append(f"{line}. {title}")
                i += 2
                # Skip one subtitle line under the section title when present.
                if i < len(texts):
                    nxt = texts[i][2].split("\n")[0].strip()
                    if (
                        nxt
                        and not re.fullmatch(r"\d{1,2}", nxt)
                        and nxt not in _TOC_TITLE_MARKERS
                        and nxt.lower() not in _TOC_SKIP_LINES
                        and len(nxt) <= 40
                    ):
                        i += 1
                continue
        if (
            line
            and not re.fullmatch(r"\d{1,2}", line)
            and line.lower() not in _TOC_SKIP_LINES
        ):
            entries.append(line)
        i += 1
    return entries


def fill_toc_from_source(slide, src_slide) -> None:
    entries = extract_toc_entries(src_slide)
    fill_toc(slide, entries)


def fill_section_from_source(slide, src_slide, chapter_num: int) -> None:
    title, _ = extract_header(src_slide, "section")
    _fill_chapter_guide(
        slide,
        [title or "", str(chapter_num)],
        title_font_pt=36,
        body_font_pt=16,
    )
    ph10 = _ph_by_idx(slide, 10)
    ph10.width = max(int(ph10.width), SECTION_PH10_MIN_WIDTH)
    _autofit_horizontal_textbox(
        ph10,
        font_pt=36,
        min_width_in=3.6,
        max_width_in=11.0,
        char_factor=0.58,
        width_padding=1.12,
    )


def _presentation_dims(slide) -> tuple[int, int]:
    prs = slide.part.package.presentation_part.presentation
    return int(prs.slide_width), int(prs.slide_height)


def needs_canvas_rescale(src_sw: int, src_sh: int, dst_sw: int, dst_sh: int) -> bool:
    """True when source slide size differs enough that 1:1 coords will overflow."""
    if src_sw <= 0 or src_sh <= 0 or dst_sw <= 0 or dst_sh <= 0:
        return False
    return (
        abs(src_sw - dst_sw) / max(src_sw, dst_sw) > 0.08
        or abs(src_sh - dst_sh) / max(src_sh, dst_sh) > 0.08
    )


def rescale_body_shapes_to_canvas(
    slide,
    src_sw: int,
    src_sh: int,
    dst_sw: int,
    dst_sh: int,
) -> float:
    """Map shapes from source canvas into the academy content band. Returns scale."""
    content_top = bac.HEADER_BOTTOM_EMU
    content_bottom = dst_sh - bac.SLIDE_MARGIN_BOTTOM
    avail_w = max(dst_sw - 2 * bac.SLIDE_MARGIN_X, 1)
    avail_h = max(content_bottom - content_top, 1)
    scale = min(avail_w / src_sw, avail_h / src_sh)
    offset_x = bac.SLIDE_MARGIN_X + (avail_w - src_sw * scale) / 2
    # Top-align under header so UI mockups keep reading order.
    offset_y = float(content_top)
    for sh in slide.shapes:
        if sh.is_placeholder:
            continue
        try:
            sh.left = int(offset_x + int(sh.left or 0) * scale)
            sh.top = int(offset_y + int(sh.top or 0) * scale)
            sh.width = max(int(int(sh.width or 0) * scale), bac.MIN_SHAPE_W)
            sh.height = max(int(int(sh.height or 0) * scale), bac.MIN_SHAPE_H)
        except Exception:
            continue
    return float(scale)


def migrate_content_body(
    slide,
    src_slide,
    step_num: str,
    slide_width: int,
    slide_height: int,
    *,
    src_classify_empty: bool = False,
    brand_pack: dict | None = None,
) -> tuple[bool, bool, int, bool, bool, str | None]:
    """Returns (radial, skip_relayout, rasters, canvas_rescaled, preserve_style, pattern)."""
    title, governing = extract_header(src_slide, "content")
    gov = governing or ""
    radial = is_radial_diagram_slide(src_slide, slide_width, slide_height)
    src_sw, src_sh = _presentation_dims(src_slide)
    canvas_mismatch = needs_canvas_rescale(src_sw, src_sh, slide_width, slide_height)
    designed_chrome = has_designed_chrome_layout(src_slide)
    force_raster = src_classify_empty or count_diagram_rasters(
        src_slide, slide_width, slide_height
    ) >= 1
    # Keep McKinsey hierarchy/accent colors; forced 12pt black flattens designed decks.
    preserve_body_style = canvas_mismatch or designed_chrome

    apply_content_header_geometry(slide, slide_width)

    _assign_plain(
        _ph_by_idx(slide, 12),
        step_num,
        title=False,
        title_font_pt=14,
        body_font_pt=14,
    )
    if gov and len(gov) <= bac.GOV_MAX_LEN:
        _assign_plain(
            _ph_by_idx(slide, 13),
            gov,
            title=False,
            title_font_pt=14,
            body_font_pt=14,
        )
    else:
        hide_empty_governing_placeholder(slide)

    id_remap = copy_body_shapes_filtered(
        src_slide,
        slide,
        title,
        governing,
        slide_width,
        slide_height,
        skip_connectors=radial,
    )
    restored = sync_raster_diagram_assets(
        src_slide,
        slide,
        title,
        governing,
        slide_width,
        slide_height,
        skip_connectors=radial,
        force_permissive=force_raster,
        src_classify_empty=src_classify_empty,
    )
    renumber_map = bac.renumber_shape_ids(slide)
    if not radial:
        bac.fix_orphan_connector_refs(slide, id_remap, renumber_map)
    canvas_rescaled = False
    if canvas_mismatch:
        rescale_body_shapes_to_canvas(
            slide, src_sw, src_sh, slide_width, slide_height
        )
        canvas_rescaled = True

    pattern_name = None
    patterns = (brand_pack or {}).get("layout_patterns") or {}
    # Equal card rows beat radial-diagram heuristics (2-col label grids false-positive).
    if patterns and not canvas_rescaled:
        n = detect_card_column_count(slide, slide_width)
        if n:
            pattern_name = apply_card_column_pattern(
                slide, patterns, slide_width, column_count=n
            )
            if pattern_name:
                radial = False

    # Use *source* geometry — placeholder heroes may be dropped before dst check.
    skip_relayout = (
        radial
        or canvas_mismatch
        or designed_chrome
        or bool(pattern_name)
        or has_hero_body_picture(src_slide, slide_width, slide_height)
        or count_body_text_columns(src_slide, slide_width) >= 3
    )

    apply_slide_title_layout(slide, slide_width, title or "")
    if not skip_relayout:
        relayout_content_columns(slide, slide_width)
    apply_academy_fonts_and_colors(slide, placeholders_only=preserve_body_style)
    if not preserve_body_style:
        apply_academy_body_shape_typography(slide)
    apply_academy_tables_on_slide(slide)
    return (
        radial,
        skip_relayout,
        restored,
        canvas_rescaled,
        preserve_body_style,
        pattern_name,
    )


def add_toc_after_cover(
    prs: Presentation,
    seeds: dict,
    src: Presentation,
    cover_src_idx: int,
    cfg: DeckMigrateConfig,
) -> None:
    if not cfg.auto_toc_after_cover:
        return
    titles: list[str] = []
    if cfg.toc_slide_range:
        for i in cfg.toc_slide_range:
            if i >= len(src.slides):
                break
            t, _ = extract_header(src.slides[i], "content")
            if t:
                titles.append(t.replace("\x0b", " ").strip())
    if not titles:
        titles = collect_section_titles(
            src, cover_src_idx + 1, part_cover_indices=cfg.part_cover_indices
        )
    if not titles:
        return
    slide = duplicate_slide_from_seed(prs, seeds[LAYOUT_TOC])
    _prepare_slide_for_editing(slide)
    fill_toc(slide, titles)


def migrate_cmp_deck(
    source: Path,
    output: Path,
    *,
    template_path: Path | None = None,
    skip_powerpoint_repair: bool = True,
    skip_ooxml_repair: bool = False,
) -> tuple[Path, list[dict], int]:
    """§7 migrate: seed clone + shape transplant. Returns (output_path, warnings)."""
    source = Path(source).resolve()
    output = Path(output).resolve()
    template = Path(template_path).resolve() if template_path else resolve_academy_template_path()

    if not source.is_file():
        raise FileNotFoundError(source)
    if not template.is_file():
        raise FileNotFoundError(template)

    warnings: list[dict] = []
    cfg = migrate_config_for_source(source)

    src = Presentation(str(source))
    if is_likely_reacademize_output(src, filename_hint=source.name):
        raise ValueError(
            "ppt-academizer 결과물(academy-output 등)은 다시 변환할 수 없습니다. "
            "원본 파트너·교육 .pptx를 업로드해 주세요."
        )
    if uses_academy_layouts(src):
        warnings.append(
            {
                "code": "ACADEMY_SOURCE_LAYOUT",
                "message": "아카데미 마스터 기반 원본 — §7로 레이아웃·도형을 다시 맞춥니다.",
            }
        )
    plan = build_slide_plan(src, cfg)
    if not plan:
        raise ValueError("No slides to migrate.")
    expected_slides = expected_output_slide_count(plan, cfg)

    output.parent.mkdir(parents=True, exist_ok=True)
    if output.exists():
        shutil.copy2(output, output.with_name(output.stem + "_backup.pptx"))

    shutil.copy2(template, output)
    prs = Presentation(str(output))
    adl._PLACEHOLDER_GEOM = (
        _capture_seed_placeholder_geometry(prs),
        _capture_layout_placeholder_geometry(prs),
    )
    try:
        return _migrate_cmp_deck_body(
            prs,
            src,
            output,
            template,
            cfg,
            plan,
            expected_slides,
            warnings,
            skip_ooxml_repair=skip_ooxml_repair,
            skip_powerpoint_repair=skip_powerpoint_repair,
        )
    finally:
        adl._PLACEHOLDER_GEOM = None


def _migrate_cmp_deck_body(
    prs,
    src,
    output: Path,
    template: Path,
    cfg,
    plan,
    expected_slides: int,
    warnings: list,
    *,
    skip_ooxml_repair: bool,
    skip_powerpoint_repair: bool,
) -> tuple[Path, list[dict], int]:
    global CLR_DK1, CLR_DK2, CLR_ACCENT1
    brand_dir = resolve_brand_dir(template)
    brand_pack = load_brand_pack(str(brand_dir)) if brand_dir else None
    if brand_pack:
        CLR_DK1 = brand_rgb(brand_pack, "dk1", "#000000")
        CLR_DK2 = brand_rgb(brand_pack, "dk2", "#44546A")
        CLR_ACCENT1 = brand_rgb(brand_pack, "accent1", "#006DFF")
        warnings.append(
            {
                "code": "BRAND_PACK",
                "level": "info",
                "message": f"Using academy brand pack at {brand_dir}",
            }
        )
    else:
        warnings.append(
            {
                "code": "BRAND_PACK_MISSING",
                "level": "info",
                "message": "brand/ pack not found next to template; using built-in theme colors.",
            }
        )

    sw, sh = int(prs.slide_width), int(prs.slide_height)
    seeds = layout_seed_slides(prs)
    for name in (LAYOUT_COVER, LAYOUT_TOC, LAYOUT_SECTION, LAYOUT_CONTENT):
        if name not in seeds:
            raise ValueError(f"Template missing starter slide for layout {name!r}")

    initial_part_ids = {id(s.part) for s in prs.slides}

    content_step = 0
    section_num = 0
    radial_slides: list = []
    skip_relayout_slides: list = []
    canvas_rescaled_slides: list = []
    preserve_body_style_slides: list = []
    for src_idx, kind in plan:
        src_slide = src.slides[src_idx]
        if kind == "cover":
            content_step = 0
            slide = duplicate_slide_from_seed(prs, seeds[LAYOUT_COVER])
            _prepare_slide_for_editing(slide)
            fill_cover_from_source(slide, src_slide, cfg)
            apply_brand_pack_to_slide(slide, brand_pack, kind="cover")
            copy_speaker_notes(src_slide, slide)
            add_toc_after_cover(prs, seeds, src, src_idx, cfg)
        elif kind == "toc":
            content_step = 0
            slide = duplicate_slide_from_seed(prs, seeds[LAYOUT_TOC])
            _prepare_slide_for_editing(slide)
            fill_toc_from_source(slide, src_slide)
            copy_speaker_notes(src_slide, slide)
        elif kind == "section":
            content_step = 0
            section_num += 1
            slide = duplicate_slide_from_seed(prs, seeds[LAYOUT_SECTION])
            _prepare_slide_for_editing(slide)
            fill_section_from_source(slide, src_slide, section_num)
            copy_speaker_notes(src_slide, slide)
        else:
            content_step += 1
            slide = duplicate_slide_from_seed(prs, seeds[LAYOUT_CONTENT])
            _prepare_slide_for_editing(slide)
            copy_speaker_notes(src_slide, slide)
            charts_in = _count_charts(src_slide)
            src_empty = (
                classify_slide(src_slide, src_idx, cfg.part_cover_indices) == "empty"
            )
            placeholder_n = sum(
                1
                for sh in src_slide.shapes
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
                and is_stock_placeholder_picture(sh)
            )
            (
                radial,
                skip_relayout,
                restored,
                canvas_rescaled,
                preserve_body_style,
                pattern_name,
            ) = migrate_content_body(
                slide,
                src_slide,
                str(content_step),
                sw,
                sh,
                src_classify_empty=src_empty,
                brand_pack=brand_pack,
            )
            if radial:
                radial_slides.append(slide)
            if skip_relayout:
                skip_relayout_slides.append(slide)
            if preserve_body_style:
                preserve_body_style_slides.append(slide)
            if pattern_name:
                warnings.append(
                    {
                        "code": "LAYOUT_PATTERN",
                        "level": "info",
                        "slide": content_step,
                        "message": (
                            f"Applied academy {pattern_name} card pattern "
                            f"on content step {content_step}."
                        ),
                    }
                )
            if canvas_rescaled:
                canvas_rescaled_slides.append(slide)
                if not any(w.get("code") == "CANVAS_RESCALED" for w in warnings):
                    src_sw, src_sh = _presentation_dims(src_slide)
                    warnings.append(
                        {
                            "code": "CANVAS_RESCALED",
                            "message": (
                                f"Source slide size {src_sw}×{src_sh} differs from "
                                f"academy {sw}×{sh}; body shapes were scaled to fit."
                            ),
                        }
                    )
            if placeholder_n:
                warnings.append(
                    {
                        "code": "SOURCE_PLACEHOLDER_IMAGE",
                        "message": (
                            f"Source slide {src_idx + 1}: {placeholder_n} placeholder "
                            "image(s) (e.g. 600×400 stub) — re-export from Google Slides "
                            "with embedded media to restore visuals."
                        ),
                        "src_index": src_idx,
                        "count": placeholder_n,
                    }
                )
            if restored:
                warnings.append(
                    {
                        "code": "RASTER_DIAGRAM_RESTORED",
                        "src_index": src_idx,
                        "count": restored,
                    }
                )
            charts_out = _count_charts(slide)
            if charts_in > charts_out:
                warnings.append(
                    {
                        "code": "CHART_NOT_COPIED",
                        "message": f"Source slide {src_idx}: charts {charts_in} → {charts_out}",
                        "src_index": src_idx,
                    }
                )
            if classify_slide(src_slide, src_idx, cfg.part_cover_indices) == "empty":
                warnings.append(
                    {
                        "code": "SLIDE_KEPT_EMPTY",
                        "message": f"Source slide {src_idx} migrated as blank content",
                        "src_index": src_idx,
                    }
                )

    to_remove = [i for i, s in enumerate(prs.slides) if id(s.part) in initial_part_ids]
    for i in reversed(to_remove):
        delete_slide(prs, i)

    for slide in prs.slides:
        radial = slide in radial_slides
        skip_relayout = slide in skip_relayout_slides
        canvas_rescaled = slide in canvas_rescaled_slides
        preserve_body_style = slide in preserve_body_style_slides
        bac.polish_slide(
            slide, skip_body_dedup=radial or canvas_rescaled or preserve_body_style
        )
        if slide.slide_layout.name == LAYOUT_CONTENT:
            title = ""
            for sh in slide.placeholders:
                if sh.placeholder_format.idx == 10:
                    title = normalize_slide_title(sh.text or "")
            hide_empty_governing_placeholder(slide)
            if not radial:
                apply_academy_fonts_and_colors(
                    slide, placeholders_only=preserve_body_style
                )
                if not preserve_body_style:
                    apply_academy_body_shape_typography(slide)
        if not radial and not canvas_rescaled and not preserve_body_style:
            bac.fit_body_shapes(slide, sw, sh)
        if slide.slide_layout.name == LAYOUT_CONTENT and not skip_relayout:
            relayout_content_columns(slide, sw)
            t = ""
            for sh in slide.placeholders:
                if sh.placeholder_format.idx == 10:
                    t = normalize_slide_title(sh.text or "")
            apply_slide_title_layout(slide, sw, t)
        elif slide.slide_layout.name == LAYOUT_CONTENT and skip_relayout:
            t = ""
            for sh in slide.placeholders:
                if sh.placeholder_format.idx == 10:
                    t = normalize_slide_title(sh.text or "")
            apply_slide_title_layout(slide, sw, t)
            if not radial:
                apply_academy_fonts_and_colors(
                    slide, placeholders_only=preserve_body_style
                )
                if not preserve_body_style:
                    apply_academy_body_shape_typography(slide)

    from scripts.academy_deck_build_lib import fix_open_in_slide_view, polish_academy_presentation

    polish_academy_presentation(prs)
    for slide in prs.slides:
        if (
            slide.slide_layout.name == LAYOUT_CONTENT
            and slide not in canvas_rescaled_slides
            and slide not in preserve_body_style_slides
        ):
            title = ""
            for sh in slide.placeholders:
                if sh.placeholder_format.idx == 10:
                    title = normalize_slide_title(sh.text or "")
            apply_slide_title_layout(slide, sw, title)
            hide_empty_governing_placeholder(slide)
            apply_academy_body_shape_typography(slide)
    icon_hits: list[str] = []
    if brand_pack:
        for slide in prs.slides:
            name = slide.slide_layout.name
            if name == LAYOUT_COVER:
                apply_brand_pack_to_slide(slide, brand_pack, kind="cover")
            elif name == LAYOUT_CONTENT:
                apply_brand_pack_to_slide(
                    slide, brand_pack, kind="content", icon_hits=icon_hits
                )
        if icon_hits:
            warnings.append(
                {
                    "code": "BRAND_ICONS_REPLACED",
                    "level": "info",
                    "message": (
                        f"Replaced {len(icon_hits)} icon(s) with academy brand pack: "
                        + ", ".join(sorted(set(icon_hits))[:12])
                    ),
                }
            )
    slide_count = len(prs.slides)
    if slide_count != expected_slides:
        warnings.append(
            {
                "code": "SLIDE_COUNT_MISMATCH",
                "message": f"Expected {expected_slides} slides, got {slide_count}",
                "expected": expected_slides,
                "actual": slide_count,
                "deck_kind": cfg.deck_kind,
            }
        )

    save_academy_deck(prs, output)

    validate_script = PPT_TEST / "scripts" / "validate_pptx.py"
    if validate_script.is_file():
        import subprocess

        proc = subprocess.run(
            [sys.executable, str(validate_script), str(output)],
            capture_output=True,
            text=True,
        )
        if proc.returncode != 0:
            raw = (proc.stderr or proc.stdout or "validate_pptx failed").strip()
            if "Traceback" in raw or "ImportError" in raw or "ModuleNotFoundError" in raw:
                msg = "validate_pptx_unavailable"
            else:
                msg = raw[:200]
            warnings.append(
                {
                    "code": "OOXML_VALIDATE_FAILED",
                    "message": msg,
                }
            )

    if not skip_ooxml_repair:
        from scripts.ooxml_repair_pptx import repair_pptx_in_place

        repair_pptx_in_place(output, template=template)
        fix_open_in_slide_view(output)
        from scripts.academy_deck_build_lib import finalize_academy_package

        finalize_academy_package(output)

    raw_output = output
    if not skip_powerpoint_repair:
        repaired_output = output.with_name(
            output.stem.replace("_아카데미적용_", "_아카데미적용_PP정리_") + ".pptx"
        )
        try:
            import platform

            if platform.system() == "Darwin":
                from scripts.powerpoint_repair_mac import powerpoint_repair_and_save

                powerpoint_repair_and_save(raw_output, repaired_output)
                output = repaired_output
        except Exception as exc:
            warnings.append(
                {"code": "PP_REPAIR_SKIPPED", "message": f"PowerPoint repair skipped: {exc}"}
            )

    warnings.insert(
        0,
        {
            "code": "MIGRATE_META",
            "deck_kind": cfg.deck_kind,
            "source_slides": len(src.slides),
            "plan_entries": len(plan),
            "expected_slides": expected_slides,
        },
    )
    return output, warnings, slide_count


def main() -> int:
    parser = argparse.ArgumentParser(description="Migrate CMP deck to academy template (§7)")
    parser.add_argument("--source", type=Path, default=DEFAULT_SOURCE, help="Source .pptx")
    parser.add_argument(
        "--out",
        type=Path,
        default=None,
        help="Output .pptx (default: source dir, stamped name)",
    )
    parser.add_argument("--template", type=Path, default=None, help="Academy template (else TEMPLATE_PPTX)")
    parser.add_argument(
        "--powerpoint-repair",
        action="store_true",
        help="Run macOS PowerPoint repair pass (slow)",
    )
    args = parser.parse_args()

    if not args.source.is_file():
        raise SystemExit(f"Missing source: {args.source}")

    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    out = args.out or args.source.with_name(
        f"{args.source.stem}_아카데미적용_{stamp}.pptx"
    )

    output, warnings, slide_count = migrate_cmp_deck(
        args.source,
        out,
        template_path=args.template,
        skip_powerpoint_repair=not args.powerpoint_repair,
    )
    print(f"Slides: {slide_count}")
    print(f"Saved: {output}")
    if warnings:
        for w in warnings:
            print(f"  warn: {w}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
